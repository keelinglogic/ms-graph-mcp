"""
Microsoft Graph Email MCP Server - Native FastMCP Implementation
Exposes Microsoft 365 email access as MCP tools for Claude Code via HTTP

Version: 2.1.0 (2026-01-16) - Deployment test from ms-graph-mcp repo
"""

import asyncio
import os
import sys
import json
import time
import base64
import mimetypes
from pathlib import Path
from datetime import datetime, timedelta
from zoneinfo import ZoneInfo
from typing import Any, Callable, TypeVar, Optional, List, Tuple, Dict
from functools import wraps
from dotenv import load_dotenv

import httpx
from mcp.server import FastMCP


from azure.identity import DeviceCodeCredential, TokenCachePersistenceOptions
from azure.core.credentials import AccessToken
from msgraph import GraphServiceClient
from msgraph.generated.users.item.messages.messages_request_builder import MessagesRequestBuilder
from msgraph.generated.users.item.messages.item.message_item_request_builder import MessageItemRequestBuilder
from msgraph.generated.models.message import Message
from msgraph.generated.models.message_collection_response import MessageCollectionResponse
from msgraph.generated.users.item.messages.item.forward.forward_post_request_body import ForwardPostRequestBody
from msgraph.generated.users.item.messages.item.move.move_post_request_body import MovePostRequestBody
from msgraph.generated.users.item.send_mail.send_mail_post_request_body import SendMailPostRequestBody
from msgraph.generated.models.recipient import Recipient
from msgraph.generated.models.email_address import EmailAddress
from msgraph.generated.models.item_body import ItemBody
from msgraph.generated.models.body_type import BodyType
from msgraph.generated.models.followup_flag import FollowupFlag
from msgraph.generated.models.followup_flag_status import FollowupFlagStatus
from msgraph.generated.users.item.messages.item.reply.reply_post_request_body import ReplyPostRequestBody
from msgraph.generated.users.item.messages.item.reply_all.reply_all_post_request_body import ReplyAllPostRequestBody
from msgraph.generated.users.item.messages.item.create_reply.create_reply_post_request_body import CreateReplyPostRequestBody
from msgraph.generated.users.item.messages.item.create_reply_all.create_reply_all_post_request_body import CreateReplyAllPostRequestBody
from msgraph.generated.models.file_attachment import FileAttachment
from msgraph.generated.models.mail_folder import MailFolder
from msgraph.generated.models.mail_folder_collection_response import MailFolderCollectionResponse
from msgraph.generated.models.outlook_category import OutlookCategory
from msgraph.generated.models.outlook_category_collection_response import OutlookCategoryCollectionResponse

# Calendar imports (Phase 9)
from msgraph.generated.models.event import Event
from msgraph.generated.models.event_collection_response import EventCollectionResponse
from msgraph.generated.models.calendar import Calendar
from msgraph.generated.models.calendar_collection_response import CalendarCollectionResponse
from msgraph.generated.models.date_time_time_zone import DateTimeTimeZone
from msgraph.generated.models.attendee import Attendee
from msgraph.generated.models.response_status import ResponseStatus
from msgraph.generated.models.free_busy_status import FreeBusyStatus
from msgraph.generated.users.item.calendar_view.calendar_view_request_builder import CalendarViewRequestBuilder
from msgraph.generated.users.item.calendars.item.calendar_view.calendar_view_request_builder import CalendarViewRequestBuilder as CalendarSpecificViewRequestBuilder
from kiota_abstractions.base_request_configuration import RequestConfiguration

# Contact imports (Phase 10)
from msgraph.generated.models.contact import Contact
from msgraph.generated.models.contact_collection_response import ContactCollectionResponse
from msgraph.generated.models.contact_folder import ContactFolder
from msgraph.generated.models.contact_folder_collection_response import ContactFolderCollectionResponse
from msgraph.generated.models.physical_address import PhysicalAddress
from msgraph.generated.users.item.contacts.contacts_request_builder import ContactsRequestBuilder

# To Do imports (Phase 11)
from msgraph.generated.models.todo_task_list import TodoTaskList
from msgraph.generated.models.todo_task_list_collection_response import TodoTaskListCollectionResponse
from msgraph.generated.models.todo_task import TodoTask
from msgraph.generated.models.todo_task_collection_response import TodoTaskCollectionResponse
from msgraph.generated.models.checklist_item import ChecklistItem
from msgraph.generated.models.checklist_item_collection_response import ChecklistItemCollectionResponse
from msgraph.generated.models.importance import Importance
from msgraph.generated.models.task_status import TaskStatus

# Planner imports (Phase 11)
from msgraph.generated.models.planner_plan import PlannerPlan
from msgraph.generated.models.planner_plan_collection_response import PlannerPlanCollectionResponse
from msgraph.generated.models.planner_bucket import PlannerBucket
from msgraph.generated.models.planner_bucket_collection_response import PlannerBucketCollectionResponse
from msgraph.generated.models.planner_task import PlannerTask
from msgraph.generated.models.planner_task_collection_response import PlannerTaskCollectionResponse
from msgraph.generated.models.planner_task_details import PlannerTaskDetails
from msgraph.generated.models.planner_assignments import PlannerAssignments
from msgraph.generated.models.planner_assigned_to_task_board_task_format import PlannerAssignedToTaskBoardTaskFormat

# Group imports (Phase 11 - Groups)
from msgraph.generated.models.group import Group
from msgraph.generated.models.group_collection_response import GroupCollectionResponse
from msgraph.generated.models.directory_object import DirectoryObject
from msgraph.generated.models.directory_object_collection_response import DirectoryObjectCollectionResponse


# Load environment variables from explicit path
env_path = Path(__file__).parent / '.env'
load_dotenv(dotenv_path=env_path)

CLIENT_ID = os.getenv('CLIENT_ID')
TENANT_ID = os.getenv('TENANT_ID')
SCOPES = os.getenv('SCOPES', 'User.Read Mail.Read Mail.Send offline_access').split()

# Display timezone for calendar events (IANA timezone name)
# Converts event times from their stored timezone to this timezone for display
MCP_MS_GRAPH_TIMEZONE = os.getenv('MCP_MS_GRAPH_TIMEZONE', 'UTC')

# Validate required environment variables
if not CLIENT_ID or not TENANT_ID:
    raise ValueError(
        "Missing required environment variables. "
        "Please ensure CLIENT_ID and TENANT_ID are set in .env file"
    )

# Data directory (configurable via environment for container deployments)
DATA_DIR = Path(os.environ.get("ATHENA_DATA_DIR", str(Path.home() / ".athena")))

# Token cache location
TOKEN_CACHE_DIR = DATA_DIR / 'credentials'
TOKEN_CACHE_NAME = "graph_mcp_cache"

# Email retrieval limits
MAX_EMAIL_COUNT = 50  # Maximum emails to return in a single request
DEFAULT_RECENT_COUNT = 5  # Default for read_recent_emails
DEFAULT_SEARCH_COUNT = 10  # Default for search_emails

# Email preview lengths
PREVIEW_LENGTH_FULL = 300  # Preview length for detailed email list
PREVIEW_LENGTH_SEARCH = 150  # Preview length for search results

# API retry configuration
MAX_RETRIES = 3  # Maximum number of retry attempts
RETRY_INITIAL_DELAY = 1.0  # Initial retry delay in seconds
RETRY_BACKOFF_FACTOR = 2.0  # Exponential backoff multiplier
RETRY_MAX_DELAY = 30.0  # Maximum retry delay in seconds

# API timeout configuration
API_TIMEOUT = 60.0  # Timeout for Graph API calls in seconds

# Global Graph client and initialization lock
_graph_client = None
_graph_client_lock = asyncio.Lock()

# Custom OAuth state tracking
_auth_state = {
    'device_code': None,      # Device code for polling
    'user_code': None,        # User-facing code to display
    'verification_uri': None, # URL user visits
    'expires_at': None,       # Unix timestamp when code expires
    'interval': 5             # Polling interval in seconds
}
_auth_state_lock = asyncio.Lock()

# Token cache file path
TOKEN_CACHE_FILE = TOKEN_CACHE_DIR / f"{TOKEN_CACHE_NAME}.json"
# Attachment download directory
ATTACHMENT_DOWNLOAD_DIR = DATA_DIR / "attachments"


def escape_odata_string(value: str) -> str:
    """
    Escape single quotes in OData string literals to prevent injection attacks.

    OData uses single quotes for string literals. User input with unescaped quotes
    could break out of the intended filter and access unauthorized data.

    Args:
        value: User input string to escape

    Returns:
        Escaped string safe for OData filter

    Example:
        Input:  "test') or (subject eq 'secret"
        Output: "test'') or (subject eq ''secret"
    """
    return value.replace("'", "''")


class CachedTokenCredential:
    """
    Custom credential that uses manually-obtained OAuth tokens from our cache.

    This prevents Azure SDK from triggering automatic device code flows.
    We control authentication entirely through test_connection's manual flow.
    """

    def __init__(self, token_data: Dict[str, Any]):
        """
        Initialize credential with cached token data.

        Args:
            token_data: Token dict from load_token_cache() with access_token and expires_in
        """
        self.token_data = token_data
        self.cached_at = token_data.get('cached_at', time.time())

    def get_token(self, *scopes, **kwargs) -> AccessToken:
        """
        Return cached access token.

        Args:
            *scopes: Required by Azure SDK (ignored, we use pre-authorized token)
            **kwargs: Additional args (ignored)

        Returns:
            AccessToken object with token and expiry
        """
        access_token = self.token_data.get('access_token')
        expires_in = self.token_data.get('expires_in', 3600)

        # Calculate expiry timestamp
        expires_on = int(self.cached_at + expires_in)

        return AccessToken(access_token, expires_on)


T = TypeVar('T')


def retry_with_backoff(func: Callable[..., T]) -> Callable[..., T]:
    """
    Decorator to retry async functions with exponential backoff.

    Retries on transient errors (network issues, rate limits, service unavailable).
    Uses exponential backoff with configurable parameters from constants.

    Args:
        func: Async function to wrap with retry logic

    Returns:
        Wrapped function with retry capability

    Raises:
        Last exception after MAX_RETRIES attempts
    """
    @wraps(func)
    async def wrapper(*args, **kwargs):
        delay = RETRY_INITIAL_DELAY
        last_exception = None

        for attempt in range(MAX_RETRIES):
            try:
                return await func(*args, **kwargs)

            except Exception as e:
                last_exception = e
                error_type = type(e).__name__

                # Check if error is retryable
                is_retryable = (
                    'timeout' in str(e).lower() or
                    'connection' in str(e).lower() or
                    'network' in str(e).lower() or
                    '429' in str(e) or  # Rate limit
                    '503' in str(e) or  # Service unavailable
                    '502' in str(e) or  # Bad gateway
                    '504' in str(e)     # Gateway timeout
                )

                if not is_retryable or attempt == MAX_RETRIES - 1:
                    # Non-retryable error or final attempt - raise
                    raise

                # Calculate backoff delay
                wait_time = min(delay, RETRY_MAX_DELAY)
                print(
                    f"⚠️  Retry {attempt + 1}/{MAX_RETRIES} after {error_type}: {str(e)[:100]}",
                    file=sys.stderr
                )
                print(f"   Waiting {wait_time:.1f}s before retry...", file=sys.stderr)

                await asyncio.sleep(wait_time)
                delay *= RETRY_BACKOFF_FACTOR

        # Should never reach here, but just in case
        raise last_exception

    return wrapper


def safe_email_address(recipient: Optional[Any]) -> Tuple[str, str]:
    """
    Safely extract name and address from a recipient object.

    Microsoft Graph API can return null email_address fields for certain
    message types (automated emails, system notifications). This function
    handles those cases gracefully.

    Args:
        recipient: Recipient object from Graph API (may be None)

    Returns:
        Tuple of (name, address) with safe defaults if data is missing

    Example:
        Normal: ("John Doe", "john@example.com")
        Null recipient: ("Unknown", "unknown")
        Null email_address: ("Unknown", "unknown")
    """
    if not recipient:
        return ("Unknown", "unknown")

    if not hasattr(recipient, 'email_address') or not recipient.email_address:
        return ("Unknown", "unknown")

    name = recipient.email_address.name or "Unknown"
    address = recipient.email_address.address or "unknown"

    return (name, address)


# ============================================================================
# Custom OAuth Device Code Flow
# ============================================================================

async def request_device_code() -> Dict[str, Any]:
    """
    Request device code from Microsoft OAuth endpoint.

    Returns device code information immediately without blocking for authorization.
    This replaces the blocking DeviceCodeCredential.prompt_callback approach.

    Returns:
        dict: Contains device_code, user_code, verification_uri, expires_in, interval

    Raises:
        httpx.HTTPError: If device code request fails
    """
    url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/devicecode"

    data = {
        'client_id': CLIENT_ID,
        'scope': ' '.join(SCOPES)
    }

    async with httpx.AsyncClient() as client:
        response = await client.post(url, data=data)
        response.raise_for_status()
        return response.json()


async def poll_for_token(device_code: str) -> Dict[str, Any]:
    """
    Poll Microsoft OAuth token endpoint to complete device code flow.

    Call this AFTER user has authorized the device code. This function
    will attempt one poll (does not loop/retry).

    Args:
        device_code: The device_code from request_device_code()

    Returns:
        dict: Token response containing access_token, refresh_token, expires_in

    Raises:
        httpx.HTTPStatusError: On HTTP errors (400 = still pending, 401 = expired, etc.)
    """
    url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/token"

    data = {
        'client_id': CLIENT_ID,
        'grant_type': 'urn:ietf:params:oauth:grant-type:device_code',
        'device_code': device_code
    }

    async with httpx.AsyncClient() as client:
        response = await client.post(url, data=data)
        response.raise_for_status()  # Raises on 400+ status codes
        return response.json()


async def refresh_access_token(refresh_token: str) -> Dict[str, Any]:
    """
    Use refresh token to obtain a new access token silently.

    This avoids the device code flow when the access token expires (~1 hour)
    but the refresh token is still valid (~90 days).

    Args:
        refresh_token: The refresh_token from a previous token response

    Returns:
        dict: New token response containing access_token, refresh_token, expires_in

    Raises:
        httpx.HTTPStatusError: On HTTP errors (400/401 = refresh token expired/invalid)
    """
    url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/token"

    data = {
        'client_id': CLIENT_ID,
        'grant_type': 'refresh_token',
        'refresh_token': refresh_token,
        'scope': ' '.join(SCOPES)
    }

    async with httpx.AsyncClient(timeout=30.0) as client:
        response = await client.post(url, data=data)
        response.raise_for_status()
        return response.json()


def save_token_cache(token_data: Dict[str, Any]) -> None:
    """
    Save OAuth token to persistent cache file.

    Args:
        token_data: Token response from poll_for_token()
    """
    global _graph_client

    TOKEN_CACHE_DIR.mkdir(parents=True, exist_ok=True)

    # Security: Validate cache directory isn't a symlink
    if TOKEN_CACHE_DIR.is_symlink():
        raise RuntimeError(
            f"Security Error: Token cache directory is a symlink: {TOKEN_CACHE_DIR}"
        )

    # Set restrictive permissions
    os.chmod(TOKEN_CACHE_DIR, 0o700)

    # Add timestamp for debugging
    token_data['cached_at'] = time.time()

    # Write token to file
    with open(TOKEN_CACHE_FILE, 'w') as f:
        json.dump(token_data, f, indent=2)

    # Set restrictive file permissions
    os.chmod(TOKEN_CACHE_FILE, 0o600)

    # Clear cached Graph client so it gets recreated with new token
    _graph_client = None

    print(f"✅ Token cached to {TOKEN_CACHE_FILE}", file=sys.stderr, flush=True)


def load_token_cache() -> Optional[Dict[str, Any]]:
    """
    Load OAuth token from persistent cache file.

    Returns:
        dict: Cached token data, or None if cache doesn't exist or is invalid
    """
    if not TOKEN_CACHE_FILE.exists():
        return None

    try:
        with open(TOKEN_CACHE_FILE, 'r') as f:
            token_data = json.load(f)

        # Validate token has required fields
        if 'access_token' not in token_data:
            print("⚠️  Token cache missing access_token, treating as invalid", file=sys.stderr)
            return None

        return token_data

    except (json.JSONDecodeError, IOError) as e:
        print(f"⚠️  Token cache corrupted or unreadable: {e}", file=sys.stderr)
        return None


def is_token_valid(token_data: Dict[str, Any]) -> bool:
    """
    Check if cached token is still valid (not expired).

    Args:
        token_data: Token from load_token_cache()

    Returns:
        bool: True if token is valid, False if expired or missing data
    """
    if not token_data:
        return False

    # Check if we have expiration info
    cached_at = token_data.get('cached_at')
    expires_in = token_data.get('expires_in')

    if not cached_at or not expires_in:
        # Can't determine expiry - assume invalid to be safe
        print("⚠️  Token cache missing expiry info, treating as expired", file=sys.stderr)
        return False

    # Check if token has expired (with 5 minute buffer for safety)
    age_seconds = time.time() - cached_at
    expires_seconds = expires_in - 300  # 5 min buffer

    if age_seconds >= expires_seconds:
        print(f"⚠️  Token expired ({age_seconds:.0f}s old, expires at {expires_seconds}s)", file=sys.stderr)
        return False

    return True


@retry_with_backoff
async def fetch_messages(
    mailbox,
    request_config: MessagesRequestBuilder.MessagesRequestBuilderGetRequestConfiguration
) -> Optional[MessageCollectionResponse]:
    """
    Fetch messages with retry logic and timeout.

    Applies timeout to prevent indefinite hangs on network issues.
    Retries on transient failures via decorator.

    Args:
        client: Authenticated Graph API client
        request_config: Request configuration with query parameters

    Returns:
        Message collection response or None

    Raises:
        TimeoutError: If request exceeds API_TIMEOUT
    """
    try:
        return await asyncio.wait_for(
            mailbox.messages.get(request_configuration=request_config),
            timeout=API_TIMEOUT
        )
    except asyncio.TimeoutError:
        raise TimeoutError(f"Graph API request timed out after {API_TIMEOUT}s")


@retry_with_backoff
async def fetch_message_by_id(
    mailbox,
    message_id: str,
    include_body: bool = False
) -> Optional[Message]:
    """
    Fetch specific message by ID with retry logic and timeout.

    Applies timeout to prevent indefinite hangs on network issues.
    Retries on transient failures via decorator.

    Args:
        client: Authenticated Graph API client
        message_id: Unique message identifier
        include_body: If True, include full body content (default: False for performance)

    Returns:
        Message object or None

    Raises:
        TimeoutError: If request exceeds API_TIMEOUT
    """
    try:
        # Build select list based on include_body parameter
        # Issue #9: Removed 'body' by default for archival performance
        # Issue #23: Allow opt-in body retrieval for get_email_by_id tool
        select_fields = ["id", "subject", "sender", "receivedDateTime", "bodyPreview",
                        "toRecipients", "ccRecipients", "bccRecipients", "importance",
                        "hasAttachments", "isRead", "isDraft", "internetMessageId",
                        "conversationId", "flag", "parentFolderId", "sentDateTime", "replyTo",
                        "categories"]

        if include_body:
            select_fields.append("body")

        # Create request configuration with all needed fields
        # This is required for proper SDK operation with delegated mailbox endpoints
        request_config = MessageItemRequestBuilder.MessageItemRequestBuilderGetRequestConfiguration(
            query_parameters=MessageItemRequestBuilder.MessageItemRequestBuilderGetQueryParameters(
                select=select_fields
            )
        )

        return await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).get(request_configuration=request_config),
            timeout=API_TIMEOUT
        )
    except asyncio.TimeoutError:
        raise TimeoutError(f"Graph API request timed out after {API_TIMEOUT}s")


async def get_graph_client() -> GraphServiceClient:
    """
    Get or create Microsoft Graph client using manually-obtained tokens.

    This function does NOT trigger automatic authentication flows.
    Authentication must be completed via test_connection first.

    Uses async lock to prevent race conditions when multiple concurrent
    requests try to initialize the client simultaneously.

    Returns:
        GraphServiceClient: Authenticated Graph API client

    Raises:
        RuntimeError: If no valid token in cache (call test_connection first)
    """
    global _graph_client

    # Fast path: client already initialized AND token still valid
    if _graph_client is not None:
        # Must check token validity - cached client may have expired token
        token_data = load_token_cache()
        if token_data and is_token_valid(token_data):
            return _graph_client
        # Token expired - clear cached client to force refresh below
        print("⏰ Cached client has expired token, will refresh...", file=sys.stderr, flush=True)
        _graph_client = None

    # Slow path: need to initialize or refresh (with lock to prevent races)
    async with _graph_client_lock:
        # Double-check after acquiring lock (another coroutine may have initialized)
        if _graph_client is not None:
            # Re-check token validity (another coroutine may have refreshed)
            token_data = load_token_cache()
            if token_data and is_token_valid(token_data):
                return _graph_client
            # Still expired, continue to refresh
            _graph_client = None

        # Load token from our custom cache
        token_data = load_token_cache()

        if not token_data:
            raise RuntimeError(
                "❌ No authentication token found.\n"
                "Please call test_connection first to authenticate."
            )

        # If token expired, attempt silent refresh before failing
        if not is_token_valid(token_data):
            refresh_token = token_data.get('refresh_token')

            if refresh_token:
                try:
                    print("🔄 Access token expired, attempting silent refresh...", file=sys.stderr, flush=True)
                    new_token_data = await refresh_access_token(refresh_token)
                    save_token_cache(new_token_data)
                    token_data = new_token_data
                    print("✅ Token refreshed successfully", file=sys.stderr, flush=True)
                except Exception as e:
                    print(f"⚠️ Token refresh failed: {e}", file=sys.stderr, flush=True)
                    raise RuntimeError(
                        "❌ Access token expired and refresh failed.\n"
                        "Please call test_connection to re-authenticate."
                    )
            else:
                raise RuntimeError(
                    "❌ Access token expired and no refresh token available.\n"
                    "Please call test_connection to re-authenticate."
                )

        print("🔐 Initializing Microsoft Graph client with cached token...", file=sys.stderr, flush=True)

        # Create credential using our cached token
        credential = CachedTokenCredential(token_data)

        # Create Graph client
        _graph_client = GraphServiceClient(credential, SCOPES)
        print("✅ Graph client initialized successfully", file=sys.stderr, flush=True)

        return _graph_client



def get_mailbox_endpoint(client, mailbox_id: str):
    """
    Get the appropriate mailbox endpoint based on mailbox_id.
    
    Args:
        client: Authenticated Graph API client
        mailbox_id: Email address or "me" for authenticated user's mailbox
    
    Returns:
        Mailbox endpoint object (either client.me or client.users.by_user_id(...))
    """
    if mailbox_id == "me":
        return client.me
    else:
        return client.users.by_user_id(mailbox_id)


# Create FastMCP server (simpler API with decorator pattern)
mcp = FastMCP("microsoft-graph-email")


@mcp.tool()
async def test_connection() -> str:
    """
    Test connection to Microsoft Graph API. Verifies authentication and API connectivity.

    This function implements smart OAuth state management:
    1. If token cached and valid → Test API connection
    2. If device code pending → Poll for token completion
    3. If no auth state → Return new device code for authorization

    Useful for debugging and health checks.
    """
    global _auth_state

    try:
        # Step 1: Check if we have a valid cached token
        token_cache = load_token_cache()

        if token_cache and is_token_valid(token_cache):
            # We have a valid token - try using it
            try:
                client = await get_graph_client()
                user = await asyncio.wait_for(
                    client.me.get(),
                    timeout=API_TIMEOUT
                )

                if not user:
                    return "❌ Connection test failed: No user profile returned"

                # Success!
                age_seconds = time.time() - token_cache.get('cached_at', time.time())
                expires_in = token_cache.get('expires_in', 0)
                time_remaining = expires_in - age_seconds

                return (
                    f"✅ Connection test successful!\n\n"
                    f"Authenticated as:\n"
                    f"  Name: {user.display_name or 'Unknown'}\n"
                    f"  Email: {user.mail or user.user_principal_name or 'Unknown'}\n"
                    f"  ID: {user.id or 'Unknown'}\n\n"
                    f"Graph API Status: Connected\n"
                    f"Token Cache: Active (expires in {time_remaining/3600:.1f} hours)"
                )

            except Exception as e:
                # Token might be invalid despite passing validation
                print(f"⚠️  Cached token failed API call: {e}", file=sys.stderr)
                # Clear invalid cache and continue to re-auth
                if TOKEN_CACHE_FILE.exists():
                    TOKEN_CACHE_FILE.unlink()
                _graph_client = None  # Clear cached client
                async with _auth_state_lock:
                    _auth_state['device_code'] = None
                    _auth_state['user_code'] = None

        # Step 2: Check if we have a pending device code authorization
        async with _auth_state_lock:
            device_code = _auth_state.get('device_code')
            user_code = _auth_state.get('user_code')
            verification_uri = _auth_state.get('verification_uri')
            expires_at = _auth_state.get('expires_at')

        if device_code and expires_at:
            # Check if code has expired
            time_remaining = expires_at - time.time()

            if time_remaining <= 0:
                # Code expired - clear state and request new one
                print("⚠️  Device code expired, requesting new one", file=sys.stderr)
                async with _auth_state_lock:
                    _auth_state['device_code'] = None
                    _auth_state['user_code'] = None
                device_code = None  # Force new code request below
            else:
                # Try polling for token (user might have authorized)
                try:
                    token_data = await poll_for_token(device_code)

                    # Success! Save token and test API
                    save_token_cache(token_data)

                    # Clear auth state
                    async with _auth_state_lock:
                        _auth_state['device_code'] = None
                        _auth_state['user_code'] = None
                        _auth_state['expires_at'] = None

                    # Test the new token
                    client = await get_graph_client()
                    user = await asyncio.wait_for(
                        client.me.get(),
                        timeout=API_TIMEOUT
                    )

                    return (
                        f"✅ Authorization successful!\n\n"
                        f"Authenticated as:\n"
                        f"  Name: {user.display_name or 'Unknown'}\n"
                        f"  Email: {user.mail or user.user_principal_name or 'Unknown'}\n"
                        f"  ID: {user.id or 'Unknown'}\n\n"
                        f"Graph API Status: Connected\n"
                        f"Token cached successfully"
                    )

                except httpx.HTTPStatusError as e:
                    if e.response.status_code == 400:
                        # Still waiting for user authorization
                        error_data = e.response.json()
                        error_code = error_data.get('error', 'unknown')

                        if error_code == 'authorization_pending':
                            minutes_remaining = int(time_remaining / 60)
                            seconds_remaining = int(time_remaining % 60)

                            return (
                                f"⏳ Waiting for authorization...\n\n"
                                f"Please complete these steps:\n"
                                f"1. Open: {verification_uri}\n"
                                f"2. Enter code: {user_code}\n"
                                f"3. Authorize the application\n"
                                f"4. Call test_connection again\n\n"
                                f"⏱️  Code expires in: {minutes_remaining}m {seconds_remaining}s\n"
                                f"💡 Tip: Authorization usually completes in 30 seconds"
                            )
                        elif error_code == 'expired_token':
                            # Code expired, clear and request new
                            async with _auth_state_lock:
                                _auth_state['device_code'] = None
                                _auth_state['user_code'] = None
                            device_code = None  # Force new code request below
                        else:
                            # Other error - clear state and retry
                            async with _auth_state_lock:
                                _auth_state['device_code'] = None
                                _auth_state['user_code'] = None
                            return (
                                f"❌ Authorization failed: {error_code}\n\n"
                                f"Call test_connection again for a new device code."
                            )
                    else:
                        raise  # Re-raise unexpected errors

        # Step 3: No valid token and no pending code - request new device code
        if not device_code:
            code_info = await request_device_code()

            # Store in auth state
            async with _auth_state_lock:
                _auth_state['device_code'] = code_info['device_code']
                _auth_state['user_code'] = code_info['user_code']
                _auth_state['verification_uri'] = code_info['verification_uri']
                _auth_state['expires_at'] = time.time() + code_info['expires_in']
                _auth_state['interval'] = code_info.get('interval', 5)

            expires_in_minutes = code_info['expires_in'] // 60

            return (
                f"🔐 Authentication Required\n\n"
                f"To access your Microsoft 365 emails, please authorize this application:\n\n"
                f"1. Open this URL in your browser:\n"
                f"   {code_info['verification_uri']}\n\n"
                f"2. Enter this code:\n"
                f"   {code_info['user_code']}\n\n"
                f"3. Sign in and authorize the requested permissions\n\n"
                f"4. Call test_connection again (within {expires_in_minutes} minutes)\n\n"
                f"⏱️  Code expires in: {expires_in_minutes} minutes\n"
                f"🔒 Permissions: User.Read, Mail.Read, Mail.Send (read + forward)"
            )

    except httpx.HTTPError as e:
        return (
            f"❌ Network error: {type(e).__name__}\n\n"
            f"Details: {str(e)}\n\n"
            f"Check network connectivity to Microsoft OAuth endpoints."
        )
    except asyncio.TimeoutError:
        return (
            f"❌ Connection test failed: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time. Check network connectivity."
        )
    except Exception as e:
        error_type = type(e).__name__
        return (
            f"❌ Connection test failed: {error_type}\n\n"
            f"Details: {str(e)}\n\n"
            f"Unable to connect to Microsoft Graph API."
        )


@mcp.tool()
async def logout() -> str:
    """
    Clear cached OAuth token to force re-authentication.

    Use this when you need to:
    - Switch to a different user account
    - Force fresh authentication
    - Clear corrupted token state

    After logout, call test_connection to get a new device code.
    """
    global _graph_client, _auth_state

    cleared = []

    # Clear token cache file
    if TOKEN_CACHE_FILE.exists():
        TOKEN_CACHE_FILE.unlink()
        cleared.append("token cache file")

    # Clear in-memory Graph client
    _graph_client = None
    cleared.append("Graph client")

    # Clear in-memory auth state
    async with _auth_state_lock:
        _auth_state = {
            'device_code': None,
            'user_code': None,
            'verification_uri': None,
            'expires_at': None,
            'interval': 5
        }
    cleared.append("auth state")

    return (
        f"✅ Logged out successfully.\n\n"
        f"Cleared: {', '.join(cleared)}\n\n"
        f"Call test_connection to authenticate with a new account."
    )


@mcp.tool()
async def read_recent_emails(
    mailbox_id: str = "thomas@sixpillar.co.uk",
    count: int = DEFAULT_RECENT_COUNT,
    folder_name: str = None
) -> str:
    """
    Read recent emails from Microsoft 365. Returns subject, sender, date, and preview for each email.
    Useful for checking latest communications.

    Args:
        mailbox_id: Email address of mailbox to access (default: thomas@sixpillar.co.uk, use "me" for authenticated user)
        count: Number of recent emails to retrieve (default 5, max 50)
        folder_name: Optional folder to filter by (e.g., "Inbox", "Archive", "Sent Items")
                    Default: None (retrieves from all folders)
                    Accepts: folder names, folder IDs, paths (Inbox/Projects), well-known names
    """
    try:
        # Validate and cap count
        count = min(max(1, count), MAX_EMAIL_COUNT)

        client = await get_graph_client()

        # If folder_name specified, get folder endpoint; else get mailbox endpoint
        if folder_name:
            # Resolve folder name/ID to folder_id using existing helper
            folder_id, resolved_name = await resolve_folder_id(folder_name, mailbox_id, client)
            # Get folder object (fetch_messages will add .messages.get())
            mailbox = get_mailbox_endpoint(client, mailbox_id)
            endpoint = mailbox.mail_folders.by_mail_folder_id(folder_id)
        else:
            # Default: Get all messages from mailbox (backward compatibility)
            endpoint = get_mailbox_endpoint(client, mailbox_id)

        # Configure request
        request_config = MessagesRequestBuilder.MessagesRequestBuilderGetRequestConfiguration(
            query_parameters=MessagesRequestBuilder.MessagesRequestBuilderGetQueryParameters(
                select=["id", "subject", "sender", "receivedDateTime", "bodyPreview", "isRead"],
                top=count,
                orderby=["receivedDateTime desc"]
            )
        )

        # Get messages with retry logic (fetch_messages adds .messages.get())
        messages = await fetch_messages(endpoint, request_config)

        if not messages or not messages.value:
            return "No messages found."

        # Format results
        result_lines = [f"📬 Retrieved {len(messages.value)} recent emails:\n"]

        for i, msg in enumerate(messages.value, 1):
            sender_name, sender_addr = safe_email_address(msg.sender)
            read_status = "📖" if msg.is_read else "📧"

            result_lines.append(f"\n{'=' * 80}")
            result_lines.append(f"{read_status} Email #{i}")
            result_lines.append(f"{'=' * 80}")
            result_lines.append(f"ID: {msg.id}")
            result_lines.append(f"Subject: {msg.subject or '(No subject)'}")
            result_lines.append(f"From: {sender_name} <{sender_addr}>")
            result_lines.append(f"Date: {msg.received_date_time}")

            if msg.body_preview:
                preview = msg.body_preview[:PREVIEW_LENGTH_FULL]
                result_lines.append(f"\nPreview: {preview}...")

        result_lines.append(f"\n{'=' * 80}")
        result_lines.append(f"\n💡 Use get_email_by_id with an ID above to read full email content")

        return "\n".join(result_lines)

    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error reading emails: {error_type}: {str(e)}"


@mcp.tool()
async def get_email_by_id(message_id: str, mailbox_id: str = "thomas@sixpillar.co.uk") -> str:
    """
    Get full details of a specific email by its ID, including complete body content.
    Use this after reading recent emails to get full email text.

    Args:
        mailbox_id: Email address of mailbox to access (default: thomas@sixpillar.co.uk, use "me" for authenticated user)
        message_id: The ID of the email message to retrieve (from read_recent_emails)
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)
        message = await fetch_message_by_id(mailbox, message_id, include_body=True)

        # Validate API response
        if not message:
            return (f"❌ Error: Email with ID '{message_id}' not found.\n\n"
                    f"The message may have been deleted or the ID is invalid.")

        # Format full email details
        sender_name, sender_addr = safe_email_address(message.sender)

        result_lines = [
            f"{'=' * 80}",
            f"📧 {message.subject or '(No subject)'} {'📖' if message.is_read else '[UNREAD]'}",
            f"{'=' * 80}",
            f"From: {sender_name} <{sender_addr}>"
        ]

        if message.to_recipients:
            to_list = ", ".join([
                f"{name} <{address}>"
                for r in message.to_recipients
                for name, address in [safe_email_address(r)]
            ])
            result_lines.append(f"To: {to_list}")

        if message.cc_recipients:
            cc_list = ", ".join([
                f"{name} <{address}>"
                for r in message.cc_recipients
                for name, address in [safe_email_address(r)]
            ])
            result_lines.append(f"Cc: {cc_list}")

        result_lines.extend([
            f"Date: {message.received_date_time}",
            f"Importance: {message.importance}",
            f"",
            f"{'=' * 80}",
            "Message Body:",
            f"{'=' * 80}",
            ""
        ])

        if message.body and message.body.content:
            # Include full body
            body_type = message.body.content_type
            result_lines.append(f"[Content Type: {body_type}]")
            result_lines.append("")
            result_lines.append(message.body.content)
        else:
            result_lines.append("(No body content)")

        if message.has_attachments:
            result_lines.append(f"\n📎 This email has attachments")

        return "\n".join(result_lines)

    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error getting email: {error_type}: {str(e)}"


@mcp.tool()
async def search_emails(mailbox_id: str = "thomas@sixpillar.co.uk", query: str = "", sender: str = "", count: int = DEFAULT_SEARCH_COUNT) -> str:
    """
    Search emails by subject keywords, sender, or both. Returns matching emails with previews.
    Useful for finding specific communications or threads.

    Args:
        mailbox_id: Email address of mailbox to access (default: thomas@sixpillar.co.uk, use "me" for authenticated user)
        query: Search keywords (searches subject/from/body using KQL)
        sender: Filter by sender email address (e.g., 'john@example.com')
        count: Max results to return (default 10, max 50)
    """
    try:
        # Validate count
        count = min(max(1, count), MAX_EMAIL_COUNT)

        if not query and not sender:
            return "❌ Error: Provide at least one of 'query' or 'sender' to search"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Microsoft Graph search uses KQL (Keyword Query Language)
        # Cannot combine $search and $filter - must use one approach

        if query and sender:
            # Use KQL syntax to combine: from:sender@example.com keyword
            # Escape quotes in user input for KQL
            escaped_sender = sender.replace('"', '\\"')
            escaped_query = query.replace('"', '\\"')
            search_str = f"\"from:{escaped_sender} {escaped_query}\""

            query_params = MessagesRequestBuilder.MessagesRequestBuilderGetQueryParameters(
                search=search_str,
                select=["id", "subject", "sender", "receivedDateTime", "bodyPreview", "isRead"],
                top=count
            )
        elif query:
            # Search by keyword only (searches subject, from, body)
            escaped_query = query.replace('"', '\\"')
            search_str = f"\"{escaped_query}\""

            query_params = MessagesRequestBuilder.MessagesRequestBuilderGetQueryParameters(
                search=search_str,
                select=["id", "subject", "sender", "receivedDateTime", "bodyPreview", "isRead"],
                top=count
            )
        else:
            # Sender only - use $filter (exact match, no search complexity)
            query_params = MessagesRequestBuilder.MessagesRequestBuilderGetQueryParameters(
                select=["id", "subject", "sender", "receivedDateTime", "bodyPreview", "isRead"],
                top=count,
                orderby=["receivedDateTime desc"],
                filter=f"from/emailAddress/address eq '{escape_odata_string(sender)}'"
            )

        request_config = MessagesRequestBuilder.MessagesRequestBuilderGetRequestConfiguration(
            query_parameters=query_params
        )

        # Get messages with retry logic
        messages = await fetch_messages(mailbox, request_config)

        if not messages or not messages.value:
            search_desc = f"query='{query}'" if query else ""
            search_desc += f" sender='{sender}'" if sender else ""
            return f"🔍 No emails found matching: {search_desc}"

        # Format results
        search_desc = f"query='{query}'" if query else ""
        search_desc += f" sender='{sender}'" if sender else ""

        result_lines = [f"🔍 Found {len(messages.value)} emails matching: {search_desc}\n"]

        for i, msg in enumerate(messages.value, 1):
            sender_name, sender_addr = safe_email_address(msg.sender)
            read_status = "📖" if msg.is_read else "📧"

            result_lines.append(f"\n{i}. {read_status} {msg.subject or '(No subject)'}")
            result_lines.append(f"   From: {sender_name} <{sender_addr}>")
            result_lines.append(f"   Date: {msg.received_date_time}")
            result_lines.append(f"   ID: {msg.id}")

            if msg.body_preview:
                preview = msg.body_preview[:PREVIEW_LENGTH_SEARCH]
                result_lines.append(f"   Preview: {preview}...")

        result_lines.append(f"\n💡 Use get_email_by_id with an ID above to read full email content")

        return "\n".join(result_lines)

    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error searching emails: {error_type}: {str(e)}"


@mcp.tool()
async def get_conversation_thread(conversation_id: str, mailbox_id: str) -> str:
    """
    Get all messages in an email conversation thread by conversationId.
    Returns messages in chronological order showing the complete conversation.

    Args:
        conversation_id: The conversation ID to retrieve (from message.conversationId property)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for authenticated user's mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        Formatted thread showing all messages in the conversation chronologically
    """
    try:
        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Filter by conversationId
        query_params = MessagesRequestBuilder.MessagesRequestBuilderGetQueryParameters(
            select=["id", "subject", "sender", "receivedDateTime", "bodyPreview", "isRead", "conversationId"],
            filter=f"conversationId eq '{escape_odata_string(conversation_id)}'",
            # Removed orderby - will sort client-side to avoid InefficientFilter error
            top=MAX_EMAIL_COUNT  # Get all messages in thread (up to max)
        )

        request_config = MessagesRequestBuilder.MessagesRequestBuilderGetRequestConfiguration(
            query_parameters=query_params
        )

        # Get messages with retry logic
        messages = await fetch_messages(mailbox, request_config)

        if not messages or not messages.value:
            return f"🔍 No messages found in conversation: {conversation_id}"

        # Sort messages chronologically (client-side to avoid API filter+orderby constraint)
        messages.value.sort(key=lambda m: m.received_date_time if m.received_date_time else datetime.min)

        # Format thread
        thread_count = len(messages.value)
        result_lines = [
            f"💬 Conversation Thread ({thread_count} messages)",
            f"Conversation ID: {conversation_id}",
            ""
        ]

        for i, msg in enumerate(messages.value, 1):
            sender_name, sender_addr = safe_email_address(msg.sender)
            read_status = "📖" if msg.is_read else "📧"

            # Thread visualization
            is_first = (i == 1)
            is_last = (i == thread_count)

            if is_first:
                prefix = "┌─"
            elif is_last:
                prefix = "└─"
            else:
                prefix = "├─"

            result_lines.append(f"{prefix} {read_status} Message #{i}/{thread_count}")
            result_lines.append(f"   Subject: {msg.subject or '(No subject)'}")
            result_lines.append(f"   From: {sender_name} <{sender_addr}>")
            result_lines.append(f"   Date: {msg.received_date_time}")
            result_lines.append(f"   ID: {msg.id}")

            if msg.body_preview:
                preview = msg.body_preview[:PREVIEW_LENGTH_SEARCH]
                # Indent preview with thread visualization
                preview_lines = preview.split('\n')
                for line in preview_lines[:3]:  # Max 3 lines
                    if not is_last:
                        result_lines.append(f"│  {line}")
                    else:
                        result_lines.append(f"   {line}")
                if len(preview) > PREVIEW_LENGTH_SEARCH:
                    if not is_last:
                        result_lines.append(f"│  ...")
                    else:
                        result_lines.append(f"   ...")

            result_lines.append("")  # Spacing between messages

        result_lines.append(f"💡 Use get_email_by_id with an ID above to read full email content")

        return "\n".join(result_lines)

    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error getting conversation thread: {error_type}: {str(e)}"


@mcp.tool()
async def get_message_conversation(message_id: str, mailbox_id: str) -> str:
    """
    Get the complete conversation thread for a specific message.
    Convenience tool that finds the conversationId and returns all messages in that thread.

    Args:
        message_id: The ID of any message in the conversation (required)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for authenticated user's mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        Formatted thread showing all messages in the conversation
    """
    try:
        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # First, get the message to find its conversationId
        query_params = MessageItemRequestBuilder.MessageItemRequestBuilderGetQueryParameters(
            select=["conversationId", "subject"]
        )
        request_config = MessageItemRequestBuilder.MessageItemRequestBuilderGetRequestConfiguration(
            query_parameters=query_params
        )

        message = await mailbox.messages.by_message_id(message_id).get(request_configuration=request_config)

        if not message:
            return f"❌ Message not found: {message_id}"

        if not message.conversation_id:
            return f"❌ Message has no conversation ID: {message_id}"

        # Now get the full conversation thread
        return await get_conversation_thread(message.conversation_id, mailbox_id)

    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error getting message conversation: {error_type}: {str(e)}"


@mcp.tool()
async def search_emails_advanced(
    mailbox_id: str,
    query: str = "",
    sender: str = "",
    from_domain: str = "",
    date_from: str = "",
    date_to: str = "",
    has_attachments: Optional[bool] = None,
    importance: str = "",
    size_min: int = 0,
    size_max: int = 0,
    count: int = DEFAULT_SEARCH_COUNT
) -> str:
    """
    Advanced email search with multiple filters using KQL (Keyword Query Language).
    Supports date ranges, attachments, domains, importance, and size filters.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for authenticated user's mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        query: Search keywords (searches subject/from/body)
        sender: Filter by specific sender email address
        from_domain: Filter by sender domain (e.g., "example.com" for all emails from @example.com)
        date_from: Start date for received filter (format: MM/DD/YYYY or YYYY-MM-DD)
        date_to: End date for received filter (format: MM/DD/YYYY or YYYY-MM-DD)
        has_attachments: Filter for emails with/without attachments (true/false)
        importance: Filter by importance level ("low", "normal", "high")
        size_min: Minimum email size in bytes (e.g., 1000 for emails > 1KB)
        size_max: Maximum email size in bytes (e.g., 5000000 for emails < 5MB)
        count: Max results to return (default 10, max 50)

    Examples:
        - Find large emails with attachments from last month:
          date_from="01/01/2025" date_to="01/31/2025" has_attachments=true size_min=1000000
        - Find all emails from a company domain:
          from_domain="sixpillar.co.uk"
        - Find high priority emails with specific keyword:
          query="urgent" importance="high"

    Returns:
        Formatted list of matching emails with previews
    """
    try:
        # Validate count
        count = min(max(1, count), MAX_EMAIL_COUNT)

        # Build KQL search query from parameters
        kql_parts = []

        # Text search
        if query:
            escaped_query = query.replace('"', '\\"')
            kql_parts.append(escaped_query)  # No quotes for text queries

        # Sender filters
        if sender:
            escaped_sender = sender.replace('"', '\\"')
            kql_parts.append(f'from:"{escaped_sender}"')

        if from_domain:
            # Domain search (e.g., from:@example.com)
            escaped_domain = from_domain.replace('"', '\\"')
            kql_parts.append(f'from:@{escaped_domain}')

        # Date range filters
        if date_from and date_to:
            # Range syntax: received:MM/DD/YYYY..MM/DD/YYYY
            kql_parts.append(f'received:{date_from}..{date_to}')
        elif date_from:
            # From date onwards: received>=MM/DD/YYYY
            kql_parts.append(f'received>={date_from}')
        elif date_to:
            # Up to date: received<=MM/DD/YYYY
            kql_parts.append(f'received<={date_to}')

        # Attachment filter
        if has_attachments is not None:
            kql_parts.append(f'hasattachments:{str(has_attachments).lower()}')

        # Importance filter
        if importance:
            importance_lower = importance.lower()
            if importance_lower in ['low', 'normal', 'high']:
                kql_parts.append(f'importance:{importance_lower}')

        # Size range filter
        if size_min > 0 and size_max > 0:
            kql_parts.append(f'size:{size_min}..{size_max}')
        elif size_min > 0:
            kql_parts.append(f'size>={size_min}')
        elif size_max > 0:
            kql_parts.append(f'size<={size_max}')

        # Check if any search criteria provided
        if not kql_parts:
            return "❌ Error: Provide at least one search parameter"

        # Combine all KQL parts with AND
        kql_query = ' AND '.join(kql_parts)

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Execute search using KQL
        query_params = MessagesRequestBuilder.MessagesRequestBuilderGetQueryParameters(
            search=f'"{kql_query}"',
            select=["id", "subject", "sender", "receivedDateTime", "bodyPreview", "isRead", "hasAttachments", "importance"],
            top=count
        )

        request_config = MessagesRequestBuilder.MessagesRequestBuilderGetRequestConfiguration(
            query_parameters=query_params
        )

        # Get messages with retry logic
        messages = await fetch_messages(mailbox, request_config)

        if not messages or not messages.value:
            return f"🔍 No emails found matching search criteria:\n{kql_query}"

        # Format results
        result_lines = [
            f"🔍 Found {len(messages.value)} emails matching advanced search:",
            f"Query: {kql_query}",
            ""
        ]

        for i, msg in enumerate(messages.value, 1):
            sender_name, sender_addr = safe_email_address(msg.sender)
            read_status = "📖" if msg.is_read else "📧"
            attachment_icon = "📎" if msg.has_attachments else ""
            importance_icon = "❗" if msg.importance and str(msg.importance).lower() == "high" else ""

            result_lines.append(f"\n{i}. {read_status}{attachment_icon}{importance_icon} {msg.subject or '(No subject)'}")
            result_lines.append(f"   From: {sender_name} <{sender_addr}>")
            result_lines.append(f"   Date: {msg.received_date_time}")
            result_lines.append(f"   ID: {msg.id}")

            if msg.body_preview:
                preview = msg.body_preview[:PREVIEW_LENGTH_SEARCH]
                result_lines.append(f"   Preview: {preview}...")

        result_lines.append(f"\n💡 Use get_email_by_id with an ID above to read full email content")

        return "\n".join(result_lines)

    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error in advanced search: {error_type}: {str(e)}"


@mcp.tool()
async def forward_email(message_id: str, to_address: str, comment: str = "", mailbox_id: str = "thomas@sixpillar.co.uk") -> str:
    """
    Forward an email to a specified recipient.

    Args:
        mailbox_id: Email address of mailbox to access (default: thomas@sixpillar.co.uk, use "me" for authenticated user)
        message_id: The ID of the email message to forward
        to_address: Email address to forward to (e.g., 'hubdoc.login+hubdoc.3a44y6zb@app.hubdoc.com' for SPD)
        comment: Optional comment to add to the forwarded message

    Returns:
        Success confirmation or error message
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        if not to_address:
            return "❌ Error: to_address is required"

        # Basic email validation
        if '@' not in to_address or '.' not in to_address.split('@')[1]:
            return f"❌ Error: '{to_address}' doesn't appear to be a valid email address"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Prepare forward request body using proper SDK types
        forward_body = ForwardPostRequestBody()

        # Create recipient with email address
        recipient = Recipient()
        email_addr = EmailAddress()
        email_addr.address = to_address
        recipient.email_address = email_addr

        forward_body.to_recipients = [recipient]

        if comment:
            forward_body.comment = comment

        # Forward the message using Graph API
        # POST /me/messages/{id}/forward
        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).forward.post(body=forward_body),
            timeout=API_TIMEOUT
        )

        # Get the original message details for confirmation
        original_message = await fetch_message_by_id(mailbox, message_id)
        subject = original_message.subject if original_message else "Unknown"

        return (
            f"✅ Email forwarded successfully!\n\n"
            f"Subject: {subject}\n"
            f"Forwarded to: {to_address}\n"
            f"Message ID: {message_id}"
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error forwarding email: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time. The email may or may not have been forwarded."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error forwarding email: {error_type}: {str(e)}"


@mcp.tool()
async def reply_email(
    message_id: str,
    body: str,
    body_type: str = "text",
    mailbox_id: str = "thomas@sixpillar.co.uk"
) -> str:
    """
    Reply to an email message, preserving conversation threading.

    Automatically preserves: conversation ID, subject prefix (Re:), original sender as recipient.
    The reply will appear in the same conversation thread in Outlook.

    Args:
        message_id: The ID of the message to reply to (required)
        body: Reply message body content (required)
        body_type: Body content type - "text" or "html" (default: "text")
        mailbox_id: Email address of mailbox to access (default: "thomas@sixpillar.co.uk")
                   Use "me" for athena@'s own mailbox

    Returns:
        Success confirmation with threading info
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        if not body:
            return "❌ Error: body (reply content) is required"

        # Validate body_type
        if body_type.lower() not in ["text", "html"]:
            return f"❌ Error: body_type must be 'text' or 'html', got '{body_type}'"

        client = await get_graph_client()

        # Resolve mailbox (use same pattern as other tools)
        target_mailbox = "me" if mailbox_id == "me" else mailbox_id

        # Create reply request body
        reply_body = ReplyPostRequestBody()
        reply_body.comment = body

        # If HTML, we need to set the message body type
        if body_type.lower() == "html":
            # Create a message with HTML body
            message = Message()
            message_body_obj = ItemBody()
            message_body_obj.content_type = BodyType.Html
            message_body_obj.content = body
            message.body = message_body_obj
            reply_body.message = message

        # Send reply using Graph API
        # POST /users/{mailbox}/messages/{id}/reply
        if target_mailbox == "me":
            await asyncio.wait_for(
                client.me.messages.by_message_id(message_id).reply.post(body=reply_body),
                timeout=API_TIMEOUT
            )
        else:
            await asyncio.wait_for(
                client.users.by_user_id(target_mailbox).messages.by_message_id(message_id).reply.post(body=reply_body),
                timeout=API_TIMEOUT
            )

        return (
            f"✅ Reply sent successfully!\n\n"
            f"Message ID: {message_id}\n"
            f"Body Type: {body_type}\n"
            f"Mailbox: {mailbox_id}\n\n"
            f"The reply has been sent and saved to Sent Items.\n"
            f"Threading preserved - reply appears in conversation."
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error sending reply: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time. The reply may or may not have been sent."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error sending reply: {error_type}: {str(e)}"


@mcp.tool()
async def reply_all_email(
    message_id: str,
    body: str,
    body_type: str = "text",
    mailbox_id: str = "thomas@sixpillar.co.uk"
) -> str:
    """
    Reply to all recipients of an email message, preserving conversation threading.

    Automatically includes: original sender + all To recipients + all CC recipients.
    The reply will appear in the same conversation thread in Outlook.

    Args:
        message_id: The ID of the message to reply to (required)
        body: Reply message body content (required)
        body_type: Body content type - "text" or "html" (default: "text")
        mailbox_id: Email address of mailbox to access (default: "thomas@sixpillar.co.uk")
                   Use "me" for athena@'s own mailbox

    Returns:
        Success confirmation with threading info
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        if not body:
            return "❌ Error: body (reply content) is required"

        # Validate body_type
        if body_type.lower() not in ["text", "html"]:
            return f"❌ Error: body_type must be 'text' or 'html', got '{body_type}'"

        client = await get_graph_client()

        # Resolve mailbox
        target_mailbox = "me" if mailbox_id == "me" else mailbox_id

        # Create reply all request body
        reply_all_body = ReplyAllPostRequestBody()
        reply_all_body.comment = body

        # If HTML, set the message body type
        if body_type.lower() == "html":
            message = Message()
            message_body_obj = ItemBody()
            message_body_obj.content_type = BodyType.Html
            message_body_obj.content = body
            message.body = message_body_obj
            reply_all_body.message = message

        # Send reply all using Graph API
        # POST /users/{mailbox}/messages/{id}/replyAll
        if target_mailbox == "me":
            await asyncio.wait_for(
                client.me.messages.by_message_id(message_id).reply_all.post(body=reply_all_body),
                timeout=API_TIMEOUT
            )
        else:
            await asyncio.wait_for(
                client.users.by_user_id(target_mailbox).messages.by_message_id(message_id).reply_all.post(body=reply_all_body),
                timeout=API_TIMEOUT
            )

        return (
            f"✅ Reply All sent successfully!\n\n"
            f"Message ID: {message_id}\n"
            f"Body Type: {body_type}\n"
            f"Mailbox: {mailbox_id}\n\n"
            f"The reply has been sent to ALL recipients and saved to Sent Items.\n"
            f"Threading preserved - reply appears in conversation."
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error sending reply all: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time. The reply may or may not have been sent."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error sending reply all: {error_type}: {str(e)}"


@mcp.tool()
async def create_reply_draft(
    message_id: str,
    body: str,
    body_type: str = "text",
    reply_all: bool = False,
    mailbox_id: str = "thomas@sixpillar.co.uk"
) -> str:
    """
    Create a draft reply to an email message for human review before sending.

    Creates a draft in the Drafts folder that preserves threading and can be:
    - Reviewed and edited in Outlook
    - Sent using send_draft_email tool
    - Further edited using edit_draft_email tool

    Args:
        message_id: The ID of the message to reply to (required)
        body: Reply message body content (required)
        body_type: Body content type - "text" or "html" (default: "text")
        reply_all: If True, includes all recipients; if False, replies only to sender (default: False)
        mailbox_id: Email address of mailbox to access (default: "thomas@sixpillar.co.uk")
                   Use "me" for athena@'s own mailbox

    Returns:
        Draft ID and confirmation
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        if not body:
            return "❌ Error: body (reply content) is required"

        # Validate body_type
        if body_type.lower() not in ["text", "html"]:
            return f"❌ Error: body_type must be 'text' or 'html', got '{body_type}'"

        client = await get_graph_client()

        # Resolve mailbox
        target_mailbox = "me" if mailbox_id == "me" else mailbox_id

        # Create message with body
        message = Message()
        message_body_obj = ItemBody()
        message_body_obj.content_type = BodyType.Html if body_type.lower() == "html" else BodyType.Text
        message_body_obj.content = body
        message.body = message_body_obj

        # Create draft reply using Graph API
        # POST /users/{mailbox}/messages/{id}/createReply or createReplyAll
        if reply_all:
            create_reply_body = CreateReplyAllPostRequestBody()
            create_reply_body.message = message

            if target_mailbox == "me":
                draft = await asyncio.wait_for(
                    client.me.messages.by_message_id(message_id).create_reply_all.post(body=create_reply_body),
                    timeout=API_TIMEOUT
                )
            else:
                draft = await asyncio.wait_for(
                    client.users.by_user_id(target_mailbox).messages.by_message_id(message_id).create_reply_all.post(body=create_reply_body),
                    timeout=API_TIMEOUT
                )
        else:
            create_reply_body = CreateReplyPostRequestBody()
            create_reply_body.message = message

            if target_mailbox == "me":
                draft = await asyncio.wait_for(
                    client.me.messages.by_message_id(message_id).create_reply.post(body=create_reply_body),
                    timeout=API_TIMEOUT
                )
            else:
                draft = await asyncio.wait_for(
                    client.users.by_user_id(target_mailbox).messages.by_message_id(message_id).create_reply.post(body=create_reply_body),
                    timeout=API_TIMEOUT
                )

        reply_type = "Reply All" if reply_all else "Reply"

        return (
            f"✅ {reply_type} draft created successfully!\n\n"
            f"Draft ID: {draft.id}\n"
            f"Original Message ID: {message_id}\n"
            f"Body Type: {body_type}\n"
            f"Reply Type: {reply_type}\n"
            f"Mailbox: {mailbox_id}\n\n"
            f"📋 The draft is now in your Outlook Drafts folder.\n"
            f"You can review, edit, and send it from Outlook,\n"
            f"or use edit_draft_email and send_draft_email tools."
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error creating reply draft: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error creating reply draft: {error_type}: {str(e)}"


@mcp.tool()
async def send_email(
    to_address: str,
    subject: str,
    body: str,
    body_type: str = "text",
    cc_addresses: str = "",
    importance: str = "normal",
    from_mailbox: str = "thomas@sixpillar.co.uk"
) -> str:
    """
    Send a new email message.

    Args:
        to_address: Primary recipient email address (required)
        subject: Email subject line (required)
        body: Email body content (required)
        body_type: Body content type - "text" or "html" (default: "text")
        cc_addresses: Optional CC recipients (comma-separated for multiple)
        importance: Email importance - "low", "normal", or "high" (default: "normal")
        from_mailbox: Email address to send from (default: thomas@sixpillar.co.uk, use "me" for authenticated user)

    Returns:
        Success confirmation or error message
    """
    try:
        if not to_address:
            return "❌ Error: to_address is required"

        if not subject:
            return "❌ Error: subject is required"

        if not body:
            return "❌ Error: body is required"

        # Basic email validation for primary recipient
        if '@' not in to_address or '.' not in to_address.split('@')[1]:
            return f"❌ Error: '{to_address}' doesn't appear to be a valid email address"

        # Validate body_type
        if body_type.lower() not in ["text", "html"]:
            return f"❌ Error: body_type must be 'text' or 'html', got '{body_type}'"

        # Validate importance
        importance_lower = importance.lower()
        if importance_lower not in ["low", "normal", "high"]:
            return f"❌ Error: importance must be 'low', 'normal', or 'high', got '{importance}'"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, from_mailbox)

        # Create the email message
        message = Message()
        message.subject = subject

        # Set body content
        message_body = ItemBody()
        message_body.content_type = BodyType.Html if body_type.lower() == "html" else BodyType.Text
        message_body.content = body
        message.body = message_body

        # Set importance
        message.importance = importance_lower

        # Add TO recipient
        to_recipient = Recipient()
        to_email = EmailAddress()
        to_email.address = to_address
        to_recipient.email_address = to_email
        message.to_recipients = [to_recipient]

        # Add CC recipients if provided
        if cc_addresses:
            cc_list = [addr.strip() for addr in cc_addresses.split(',')]
            cc_recipients = []

            for cc_addr in cc_list:
                if cc_addr:  # Skip empty strings
                    # Basic validation
                    if '@' not in cc_addr or '.' not in cc_addr.split('@')[1]:
                        return f"❌ Error: '{cc_addr}' doesn't appear to be a valid email address"

                    cc_recipient = Recipient()
                    cc_email = EmailAddress()
                    cc_email.address = cc_addr
                    cc_recipient.email_address = cc_email
                    cc_recipients.append(cc_recipient)

            if cc_recipients:
                message.cc_recipients = cc_recipients

        # Create send mail request body
        send_body = SendMailPostRequestBody()
        send_body.message = message
        send_body.save_to_sent_items = True  # Save copy to Sent Items

        # Send the message using Graph API
        # POST /me/sendMail
        await asyncio.wait_for(
            mailbox.send_mail.post(body=send_body),
            timeout=API_TIMEOUT
        )

        cc_info = f"\nCc: {cc_addresses}" if cc_addresses else ""

        return (
            f"✅ Email sent successfully!\n\n"
            f"To: {to_address}{cc_info}\n"
            f"Subject: {subject}\n"
            f"Body Type: {body_type}\n"
            f"Importance: {importance}"
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error sending email: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time. The email may or may not have been sent."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error sending email: {error_type}: {str(e)}"


@mcp.tool()
async def create_draft_email(
    to_address: str,
    subject: str,
    body: str,
    body_type: str = "text",
    cc_addresses: str = "",
    importance: str = "normal",
    mailbox_id: str = "thomas@sixpillar.co.uk"
) -> str:
    """
    Create a new email draft in Outlook Drafts folder. Draft can be reviewed and edited before sending.

    Args:
        mailbox_id: Email address of mailbox to access (default: thomas@sixpillar.co.uk, use "me" for authenticated user)
        to_address: Primary recipient email address (required)
        subject: Email subject line (required)
        body: Email body content (required)
        body_type: Body content type - "text" or "html" (default: "text")
        cc_addresses: Optional CC recipients (comma-separated for multiple)
        importance: Email importance - "low", "normal", or "high" (default: "normal")

    Returns:
        Success confirmation with draft ID
    """
    try:
        if not to_address:
            return "❌ Error: to_address is required"

        if not subject:
            return "❌ Error: subject is required"

        if not body:
            return "❌ Error: body is required"

        # Basic email validation for primary recipient
        if '@' not in to_address or '.' not in to_address.split('@')[1]:
            return f"❌ Error: '{to_address}' doesn't appear to be a valid email address"

        # Validate body_type
        if body_type.lower() not in ["text", "html"]:
            return f"❌ Error: body_type must be 'text' or 'html', got '{body_type}'"

        # Validate importance
        importance_lower = importance.lower()
        if importance_lower not in ["low", "normal", "high"]:
            return f"❌ Error: importance must be 'low', 'normal', or 'high', got '{importance}'"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Create the email message
        message = Message()
        message.subject = subject

        # Set body content
        message_body = ItemBody()
        message_body.content_type = BodyType.Html if body_type.lower() == "html" else BodyType.Text
        message_body.content = body
        message.body = message_body

        # Set importance
        message.importance = importance_lower

        # Add TO recipient
        to_recipient = Recipient()
        to_email = EmailAddress()
        to_email.address = to_address
        to_recipient.email_address = to_email
        message.to_recipients = [to_recipient]

        # Add CC recipients if provided
        if cc_addresses:
            cc_list = [addr.strip() for addr in cc_addresses.split(',')]
            cc_recipients = []

            for cc_addr in cc_list:
                if cc_addr:  # Skip empty strings
                    # Basic validation
                    if '@' not in cc_addr or '.' not in cc_addr.split('@')[1]:
                        return f"❌ Error: '{cc_addr}' doesn't appear to be a valid email address"

                    cc_recipient = Recipient()
                    cc_email = EmailAddress()
                    cc_email.address = cc_addr
                    cc_recipient.email_address = cc_email
                    cc_recipients.append(cc_recipient)

            if cc_recipients:
                message.cc_recipients = cc_recipients

        # Create draft using Graph API
        # POST /me/messages (creates draft, does not send)
        result = await asyncio.wait_for(
            mailbox.messages.post(body=message),
            timeout=API_TIMEOUT
        )

        cc_info = f"\nCc: {cc_addresses}" if cc_addresses else ""

        return (
            f"✅ Draft created successfully!\n\n"
            f"Draft ID: {result.id}\n"
            f"To: {to_address}{cc_info}\n"
            f"Subject: {subject}\n"
            f"Body Type: {body_type}\n"
            f"Importance: {importance}\n\n"
            f"📋 The draft is now in your Outlook Drafts folder.\n"
            f"You can review, edit, and send it from Outlook."
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error creating draft: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time. The draft may or may not have been created."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error creating draft: {error_type}: {str(e)}"




@mcp.tool()
async def send_draft_email(message_id: str, mailbox_id: str = "thomas@sixpillar.co.uk") -> str:
    """
    Send an existing draft email by its message ID.

    Args:
        message_id: The ID of the draft message to send (from create_draft_email or get_email_by_id)

    Returns:
        Success confirmation or error message
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Send the draft message
        # POST /me/messages/{id}/send
        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).send.post(),
            timeout=API_TIMEOUT
        )

        return (
            f"✅ Draft email sent successfully!\n\n"
            f"Message ID: {message_id}\n\n"
            f"The email has been sent and moved to your Sent Items folder."
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error sending draft: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time. The draft may or may not have been sent."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error sending draft: {error_type}: {str(e)}"


@mcp.tool()
async def edit_draft_email(
    message_id: str,
    subject: str = "",
    body: str = "",
    body_type: str = "",
    to_addresses: str = "",
    cc_addresses: str = "",
    importance: str = "",
    mailbox_id: str = "thomas@sixpillar.co.uk"
) -> str:
    """
    Edit an existing draft email. Only provided fields will be updated.

    Args:
        mailbox_id: Email address of mailbox to access (default: thomas@sixpillar.co.uk, use "me" for authenticated user)
        message_id: The ID of the draft message to edit (required)
        subject: New email subject (optional - leave empty to keep current)
        body: New email body content (optional - leave empty to keep current)
        body_type: Body content type - "text" or "html" (optional - leave empty to keep current)
        to_addresses: Primary recipients (comma-separated) (optional - leave empty to keep current)
        cc_addresses: CC recipients (comma-separated) (optional - leave empty to keep current)
        importance: Email importance - "low", "normal", or "high" (optional - leave empty to keep current)

    Returns:
        Success confirmation with updated fields or error message
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        # Validate body_type if provided
        if body_type and body_type.lower() not in ["text", "html"]:
            return f"❌ Error: body_type must be 'text' or 'html', got '{body_type}'"

        # Validate importance if provided
        if importance and importance.lower() not in ["low", "normal", "high"]:
            return f"❌ Error: importance must be 'low', 'normal', or 'high', got '{importance}'"

        # Validate email addresses if provided
        if to_addresses:
            to_list = [addr.strip() for addr in to_addresses.split(',')]
            for addr in to_list:
                if addr and ('@' not in addr or '.' not in addr.split('@')[1]):
                    return f"❌ Error: '{addr}' doesn't appear to be a valid email address"

        if cc_addresses:
            cc_list = [addr.strip() for addr in cc_addresses.split(',')]
            for addr in cc_list:
                if addr and ('@' not in addr or '.' not in addr.split('@')[1]):
                    return f"❌ Error: '{addr}' doesn't appear to be a valid email address"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Build update message with only provided fields
        message = Message()
        updated_fields = []

        # Update subject if provided
        if subject:
            message.subject = subject
            updated_fields.append(f"Subject: {subject}")

        # Update body if provided
        if body:
            message_body = ItemBody()
            if body_type:
                message_body.content_type = BodyType.Html if body_type.lower() == "html" else BodyType.Text
                updated_fields.append(f"Body Type: {body_type.lower()}")
            else:
                # If body provided but no body_type, default to text
                message_body.content_type = BodyType.Text
            message_body.content = body
            message.body = message_body
            preview = body[:50] + "..." if len(body) > 50 else body
            updated_fields.append(f"Body: {preview}")

        # Update importance if provided
        if importance:
            message.importance = importance.lower()
            updated_fields.append(f"Importance: {importance.lower()}")

        # Update TO recipients if provided
        if to_addresses:
            to_list = [addr.strip() for addr in to_addresses.split(',')]
            to_recipients = []
            for to_addr in to_list:
                if to_addr:
                    to_recipient = Recipient()
                    to_email = EmailAddress()
                    to_email.address = to_addr
                    to_recipient.email_address = to_email
                    to_recipients.append(to_recipient)
            if to_recipients:
                message.to_recipients = to_recipients
                updated_fields.append(f"To: {to_addresses}")

        # Update CC recipients if provided
        if cc_addresses:
            cc_list = [addr.strip() for addr in cc_addresses.split(',')]
            cc_recipients = []
            for cc_addr in cc_list:
                if cc_addr:
                    cc_recipient = Recipient()
                    cc_email = EmailAddress()
                    cc_email.address = cc_addr
                    cc_recipient.email_address = cc_email
                    cc_recipients.append(cc_recipient)
            if cc_recipients:
                message.cc_recipients = cc_recipients
                updated_fields.append(f"Cc: {cc_addresses}")

        # Check if any fields were actually provided
        if not updated_fields:
            return "❌ Error: No fields to update. Provide at least one field to edit."

        # Update the draft message
        # PATCH /me/messages/{id}
        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).patch(body=message),
            timeout=API_TIMEOUT
        )

        updated_list = "\n".join([f"  - {field}" for field in updated_fields])

        return (
            f"✅ Draft email updated successfully!\n\n"
            f"Message ID: {message_id}\n\n"
            f"Updated fields:\n{updated_list}\n\n"
            f"The draft has been updated in your Outlook Drafts folder."
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error updating draft: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time. The draft may or may not have been updated."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error updating draft: {error_type}: {str(e)}"

@mcp.tool()
async def delete_email(message_id: str, mailbox_id: str = "thomas@sixpillar.co.uk") -> str:
    """
    Delete an email by moving it to the Deleted Items folder.

    Args:
        message_id: The ID of the email message to delete

    Returns:
        Success confirmation or error message
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Delete the message using Graph API
        # DELETE /me/messages/{id}
        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).delete(),
            timeout=API_TIMEOUT
        )

        return (
            f"✅ Email deleted successfully!\n\n"
            f"Message ID: {message_id}\n"
            f"Status: Moved to Deleted Items"
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error deleting email: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error deleting email: {error_type}: {str(e)}"


@mcp.tool()
async def mark_email_read(message_id: str, is_read: bool = True, mailbox_id: str = "thomas@sixpillar.co.uk") -> str:
    """
    Mark an email as read or unread.

    Args:
        mailbox_id: Email address of mailbox to access (default: thomas@sixpillar.co.uk, use "me" for authenticated user)
        message_id: The ID of the email message to update
        is_read: True to mark as read, False to mark as unread (default: True)

    Returns:
        Success confirmation or error message
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Create a message object with just the isRead property
        message = Message()
        message.is_read = is_read

        # Update the message using Graph API
        # PATCH /me/messages/{id}
        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).patch(body=message),
            timeout=API_TIMEOUT
        )

        status = "read" if is_read else "unread"
        return (
            f"✅ Email marked as {status} successfully!\n\n"
            f"Message ID: {message_id}\n"
            f"Status: {status}"
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error updating email: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error updating email: {error_type}: {str(e)}"


@mcp.tool()
async def move_email(message_id: str, destination_folder: str, mailbox_id: str = "thomas@sixpillar.co.uk") -> str:
    """
    Move an email to a specified folder.

    Args:
        mailbox_id: Email address of mailbox to access (default: thomas@sixpillar.co.uk, use "me" for authenticated user)
        message_id: The ID of the email message to move
        destination_folder: Name of the destination folder (e.g., "Archive", "Inbox", "Deleted Items")
                          Can also be a folder ID for custom folders

    Returns:
        Success confirmation or error message
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        if not destination_folder:
            return "❌ Error: destination_folder is required"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # First, try to find the folder by name
        # Get all mail folders
        folders_response = await asyncio.wait_for(
            mailbox.mail_folders.get(),
            timeout=API_TIMEOUT
        )

        destination_folder_id = None
        if folders_response and folders_response.value:
            for folder in folders_response.value:
                if folder.display_name and folder.display_name.lower() == destination_folder.lower():
                    destination_folder_id = folder.id
                    break

        # If not found by name, assume it's a folder ID
        if not destination_folder_id:
            destination_folder_id = destination_folder

        # Move the message using Graph API
        # POST /me/messages/{id}/move
        # Create proper request body object for move operation
        move_body = MovePostRequestBody()
        move_body.destination_id = destination_folder_id
        
        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).move.post(
                body=move_body
            ),
            timeout=API_TIMEOUT
        )

        return (
            f"✅ Email moved successfully!\n\n"
            f"Message ID: {message_id}\n"
            f"Destination: {destination_folder}"
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error moving email: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error moving email: {error_type}: {str(e)}"


@mcp.tool()
async def archive_email(message_id: str, mailbox_id: str = "thomas@sixpillar.co.uk") -> str:
    """
    Archive an email by moving it to the Archive folder.

    Args:
        message_id: The ID of the email message to archive

    Returns:
        Success confirmation or error message
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Find the Archive folder
        folders_response = await asyncio.wait_for(
            mailbox.mail_folders.get(),
            timeout=API_TIMEOUT
        )

        archive_folder_id = None
        if folders_response and folders_response.value:
            for folder in folders_response.value:
                if folder.display_name and folder.display_name.lower() == 'archive':
                    archive_folder_id = folder.id
                    break

        if not archive_folder_id:
            return (
                f"❌ Error archiving email: Archive folder not found\n\n"
                f"Please ensure your mailbox has an Archive folder."
            )

        # Move the message to Archive folder
        # POST /me/messages/{id}/move
        archive_body = MovePostRequestBody()
        archive_body.destination_id = archive_folder_id

        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).move.post(
                body=archive_body
            ),
            timeout=API_TIMEOUT
        )

        return (
            f"✅ Email archived successfully!\n\n"
            f"Message ID: {message_id}\n"
            f"Status: Moved to Archive folder"
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error archiving email: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error archiving email: {error_type}: {str(e)}"


@mcp.tool()
async def flag_email(message_id: str, flag_status: str = "flagged", mailbox_id: str = "thomas@sixpillar.co.uk") -> str:
    """
    Flag an email for follow-up or mark as complete.

    Args:
        mailbox_id: Email address of mailbox to access (default: thomas@sixpillar.co.uk, use "me" for authenticated user)
        message_id: The ID of the email message to flag
        flag_status: Flag status - "flagged" (for follow-up), "complete", or "notFlagged" (default: "flagged")

    Returns:
        Success confirmation or error message
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        # Validate flag_status
        valid_statuses = ["flagged", "complete", "notflagged"]
        flag_status_lower = flag_status.lower()
        if flag_status_lower not in valid_statuses:
            return f"❌ Error: flag_status must be 'flagged', 'complete', or 'notFlagged', got '{flag_status}'"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Create a message object with flag property
        message = Message()
        flag = FollowupFlag()

        # Map string status to FollowupFlagStatus enum
        status_map = {
            'flagged': FollowupFlagStatus.Flagged,
            'complete': FollowupFlagStatus.Complete,
            'notflagged': FollowupFlagStatus.NotFlagged
        }
        flag.flag_status = status_map[flag_status_lower]
        message.flag = flag

        # Update the message using Graph API
        # PATCH /me/messages/{id}
        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).patch(body=message),
            timeout=API_TIMEOUT
        )

        status_display = {
            'flagged': 'flagged for follow-up',
            'complete': 'marked as complete',
            'notflagged': 'unflagged'
        }

        return (
            f"✅ Email {status_display.get(flag_status_lower, flag_status_lower)} successfully!\n\n"
            f"Message ID: {message_id}\n"
            f"Flag Status: {flag_status_lower}"
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error flagging email: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error flagging email: {error_type}: {str(e)}"

@mcp.tool()
async def list_attachments(

    message_id: str,

    mailbox_id: str = "thomas@sixpillar.co.uk"

) -> str:

    """

    List all attachments for a specific email message.



    Args:

        message_id: The ID of the email message

        mailbox_id: Email address of the mailbox to access (default: "thomas@sixpillar.co.uk")

                   Use "me" for athena@'s own mailbox



    Returns:

        Formatted list of attachments with details (name, ID, type, size, content type)

    """

    try:

        client = await get_graph_client()



        # Determine which mailbox to query

        if mailbox_id == "me":

            mailbox = client.me

        else:

            mailbox = client.users.by_user_id(mailbox_id)



        # Fetch attachments

        attachments = await asyncio.wait_for(

            mailbox.messages.by_message_id(message_id).attachments.get(),

            timeout=API_TIMEOUT

        )



        if not attachments or not attachments.value:

            return (

                f"📎 No attachments found for this email\n\n"

                f"Message ID: {message_id}\n"

                f"Mailbox: {mailbox_id}"

            )



        # Format attachment details

        result_lines = [

            f"📎 Attachments ({len(attachments.value)}) for message",

            f"Message ID: {message_id}",

            f"Mailbox: {mailbox_id}",

            ""

        ]



        for i, att in enumerate(attachments.value, 1):

            # Extract attachment details

            att_name = att.name or "(No name)"

            att_id = att.id or "(No ID)"

            att_type = att.odata_type.split('.')[-1] if att.odata_type else "unknown"

            content_type = att.content_type or "unknown"

            size_kb = att.size / 1024 if att.size else 0

            is_inline = getattr(att, 'is_inline', False)



            result_lines.append(f"{i}. {att_name}")

            result_lines.append(f"   ID: {att_id}")

            result_lines.append(f"   Type: {att_type}")

            result_lines.append(f"   Content-Type: {content_type}")

            result_lines.append(f"   Size: {size_kb:.2f} KB")

            result_lines.append(f"   Inline: {is_inline}")

            result_lines.append("")



        result_lines.append("💡 Use download_attachment to save attachments (coming in Phase 2)")



        return "\n".join(result_lines)



    except asyncio.TimeoutError:

        return (

            f"❌ Error listing attachments: Timeout after {API_TIMEOUT}s\n\n"

            f"The Graph API did not respond in time."

        )

    except Exception as e:

        error_type = type(e).__name__

        return f"❌ Error listing attachments: {error_type}: {str(e)}"





@mcp.tool()

async def download_attachment(

    message_id: str,

    attachment_id: str,

    filename: Optional[str] = None,

    mailbox_id: str = "thomas@sixpillar.co.uk"

) -> str:

    """

    Download an email attachment to secure temporary storage.



    Args:

        message_id: The ID of the email message

        attachment_id: The ID of the attachment to download

        filename: Optional custom filename (defaults to attachment's original name)

        mailbox_id: Email address of the mailbox to access (default: "thomas@sixpillar.co.uk")

                   Use "me" for athena@'s own mailbox



    Returns:

        Download confirmation with file path and details

    """

    try:

        client = await get_graph_client()



        # Determine which mailbox to query

        if mailbox_id == "me":

            mailbox = client.me

        else:

            mailbox = client.users.by_user_id(mailbox_id)



        # Fetch attachment metadata and content

        attachment = await asyncio.wait_for(

            mailbox.messages.by_message_id(message_id).attachments.by_attachment_id(attachment_id).get(),

            timeout=API_TIMEOUT

        )



        if not attachment:

            return (

                f"❌ Attachment not found\n\n"

                f"Message ID: {message_id}\n"

                f"Attachment ID: {attachment_id}"

            )



        # Extract attachment details

        att_name = filename or attachment.name or "unnamed_attachment"

        att_type = attachment.odata_type.split('.')[-1] if attachment.odata_type else "unknown"

        content_type = attachment.content_type or "application/octet-stream"

        size_bytes = attachment.size or 0



        # Only handle fileAttachment type (direct download)

        if att_type != "fileAttachment":

            return (

                f"❌ Unsupported attachment type: {att_type}\n\n"

                f"Only 'fileAttachment' type is supported for download.\n"

                f"itemAttachment and referenceAttachment require different handling."

            )



        # Validate filename for security (prevent directory traversal)

        safe_filename = Path(att_name).name  # Strip any directory components

        if not safe_filename or safe_filename.startswith('.') or '/' in safe_filename or '\\' in safe_filename:

            safe_filename = f"attachment_{attachment_id[:8]}"



        # Create download directory if it doesn't exist

        ATTACHMENT_DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)



        # Generate unique filename to avoid collisions

        timestamp = int(time.time())

        unique_filename = f"{timestamp}_{safe_filename}"

        file_path = ATTACHMENT_DOWNLOAD_DIR / unique_filename



        # Get attachment content (base64 encoded)

        content_bytes = attachment.content_bytes

        if not content_bytes:

            return (

                f"❌ No content available for attachment\n\n"

                f"Attachment: {att_name}\n"

                f"ID: {attachment_id}"

            )



        # Decode and save content

        try:

            # contentBytes from Graph API is base64 encoded - must decode before writing

            decoded_content = base64.b64decode(content_bytes)

            with open(file_path, 'wb') as f:

                f.write(decoded_content)

        except Exception as write_error:

            return (

                f"❌ Error writing file: {type(write_error).__name__}\n\n"

                f"{str(write_error)}"

            )



        # Verify file was written

        if not file_path.exists():

            return f"❌ File was not created: {file_path}"



        actual_size = file_path.stat().st_size

        size_kb = actual_size / 1024



        return (

            f"✅ Attachment downloaded successfully!\n\n"

            f"File: {unique_filename}\n"

            f"Path: {file_path}\n"

            f"Size: {size_kb:.2f} KB ({actual_size:,} bytes)\n"

            f"Content-Type: {content_type}\n\n"

            f"Original attachment: {att_name}\n"

            f"Message ID: {message_id}\n"

            f"Attachment ID: {attachment_id}\n\n"

            f"💡 Use extract_attachment_text to extract text from this file (coming in Phase 3)"

        )



    except asyncio.TimeoutError:

        return (

            f"❌ Error downloading attachment: Timeout after {API_TIMEOUT}s\n\n"

            f"The Graph API did not respond in time."

        )

    except Exception as e:

        error_type = type(e).__name__

        return f"❌ Error downloading attachment: {error_type}: {str(e)}"





@mcp.tool()
async def download_attachment_base64(
    message_id: str,
    attachment_id: str,
    mailbox_id: str = "thomas@sixpillar.co.uk"
) -> str:
    """
    Download email attachment and return base64-encoded content for writing to laptop/workspace.

    **USE THIS TOOL WHEN:**
    - User requests: "download attachment to Downloads", "save invoice locally"
    - File needs to be on laptop or workspace filesystem (not container)
    - User will open/view/edit the file
    - Interactive workflow where user wants the actual file

    **DO NOT USE THIS TOOL WHEN:**
    - Server-side automation only (use download_attachment instead)
    - Immediately extracting text afterwards (use download_attachment + extract_attachment_text)
    - Batch processing multiple attachments (use download_attachment)
    - User doesn't need to see/access the file

    **WORKFLOW:**
    1. Call this tool to get base64 content
    2. Use Write() to save to laptop/workspace: Write(f"~/Downloads/{filename}", base64_decode(content))

    Args:
        message_id: The ID of the email message (required)
        attachment_id: The ID of the attachment to download (required)
        mailbox_id: Email address of mailbox (default: "thomas@sixpillar.co.uk", use "me" for athena@)

    Returns:
        JSON string with: {filename, base64_content, size_bytes, mime_type}

    Limitations:
        - **CRITICAL: Context window limit ~600-800 KB** (base64 consumes Claude's context)
        - Files > 500 KB may fail or truncate (use download_attachment + SSH instead)
        - Max file size: 25MB (Exchange Online limit, but context limit hits first)
        - Base64 adds ~33% overhead in JSON response
        - Not efficient for batch automation

    Best Practice:
        - Check file size first via list_attachments
        - If > 500 KB: Use download_attachment (container filesystem) instead
        - If < 500 KB: This tool works great

    Examples:
        - User: "Download this invoice and save to my Downloads"
          → download_attachment_base64() + Write() to laptop
        - User: "Save this contract locally"
          → download_attachment_base64() + Write() to workspace
    """
    try:
        client = await get_graph_client()

        # Determine which mailbox to query
        if mailbox_id == "me":
            mailbox = client.me
        else:
            mailbox = client.users.by_user_id(mailbox_id)

        # Fetch attachment metadata and content
        attachment = await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).attachments.by_attachment_id(attachment_id).get(),
            timeout=API_TIMEOUT
        )

        if not attachment:
            return json.dumps({
                "error": "Attachment not found",
                "message_id": message_id,
                "attachment_id": attachment_id
            })

        # Extract attachment details
        filename = attachment.name or "unnamed_attachment"
        att_type = attachment.odata_type.split('.')[-1] if attachment.odata_type else "unknown"
        content_type = attachment.content_type or "application/octet-stream"
        size_bytes = attachment.size or 0

        # Only handle fileAttachment type
        if att_type != "fileAttachment":
            return json.dumps({
                "error": f"Unsupported attachment type: {att_type}",
                "message": "Only 'fileAttachment' type is supported. itemAttachment and referenceAttachment require different handling."
            })

        # Get attachment content (already base64 encoded by Graph API)
        content_bytes = attachment.content_bytes
        if not content_bytes:
            return json.dumps({
                "error": "Attachment has no content",
                "filename": filename
            })

        # content_bytes is already base64-encoded bytes from Graph API
        # Convert to string for JSON serialization
        base64_content = content_bytes.decode('utf-8')

        # Return structured JSON
        result = {
            "filename": filename,
            "base64_content": base64_content,
            "size_bytes": size_bytes,
            "mime_type": content_type
        }

        return json.dumps(result, indent=2)

    except asyncio.TimeoutError:
        return json.dumps({
            "error": "Timeout downloading attachment",
            "message": f"Request timed out after {API_TIMEOUT} seconds"
        })
    except Exception as e:
        return json.dumps({
            "error": "Failed to download attachment",
            "message": str(e)
        })




@mcp.tool()

async def extract_attachment_text(

    file_path: str,

    max_length: int = 10000

) -> str:

    """

    Extract text content from a downloaded attachment file.



    Supports: PDF (.pdf), Word (.docx), Excel (.xlsx), PowerPoint (.pptx), Text (.txt, .md, .log, etc)



    Args:

        file_path: Path to the downloaded attachment file (from download_attachment)

        max_length: Maximum characters to return (default: 10000)



    Returns:

        Extracted text content with metadata

    """

    try:

        # Validate file path

        path = Path(file_path)

        if not path.exists():

            return (

                f"❌ File not found: {file_path}\n\n"

                f"Make sure the file was downloaded using download_attachment first."

            )



        if not path.is_file():

            return f"❌ Not a file: {file_path}"



        # Get file extension

        extension = path.suffix.lower()

        file_size_kb = path.stat().st_size / 1024



        # Import text extraction libraries (lazy import to avoid startup failures)

        try:

            if extension == '.pdf':

                from pypdf import PdfReader

            elif extension == '.docx':

                from docx import Document

            elif extension == '.xlsx':

                from openpyxl import load_workbook

            elif extension == '.pptx':

                from pptx import Presentation

        except ImportError as import_error:

            return (

                f"❌ Missing dependency for {extension} files\n\n"

                f"Error: {import_error}\n\n"

                f"The required Python package may not be installed in the container."

            )



        # Extract text based on file type

        extracted_text = ""

        page_count = 0



        try:

            if extension == '.pdf':

                # Extract PDF text

                reader = PdfReader(str(path))

                page_count = len(reader.pages)

                pages_text = []

                for i, page in enumerate(reader.pages, 1):

                    page_text = page.extract_text()

                    if page_text:

                        pages_text.append(f"--- Page {i} ---\n{page_text}")

                extracted_text = "\n\n".join(pages_text)



            elif extension == '.docx':

                # Extract Word document text

                doc = Document(str(path))

                page_count = len(doc.paragraphs)

                paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

                extracted_text = "\n\n".join(paragraphs)



            elif extension == '.xlsx':

                # Extract Excel spreadsheet text

                wb = load_workbook(str(path), data_only=True)

                sheets_text = []

                for sheet_name in wb.sheetnames:

                    ws = wb[sheet_name]

                    sheet_text = [f"--- Sheet: {sheet_name} ---"]

                    for row in ws.iter_rows(values_only=True):

                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])

                        if row_text.strip(" |"):

                            sheet_text.append(row_text)

                    sheets_text.append("\n".join(sheet_text))

                extracted_text = "\n\n".join(sheets_text)

                page_count = len(wb.sheetnames)



            elif extension == '.pptx':

                # Extract PowerPoint text

                prs = Presentation(str(path))

                page_count = len(prs.slides)

                slides_text = []

                for i, slide in enumerate(prs.slides, 1):

                    slide_text = [f"--- Slide {i} ---"]

                    for shape in slide.shapes:

                        if hasattr(shape, "text") and shape.text:

                            slide_text.append(shape.text)

                    slides_text.append("\n".join(slide_text))

                extracted_text = "\n\n".join(slides_text)



            elif extension in ['.txt', '.md', '.log', '.csv', '.json', '.xml', '.html']:

                # Extract plain text files

                try:

                    extracted_text = path.read_text(encoding='utf-8')

                except UnicodeDecodeError:

                    # Try with latin-1 encoding as fallback

                    extracted_text = path.read_text(encoding='latin-1')



            else:

                return (

                    f"❌ Unsupported file type: {extension}\n\n"

                    f"Supported formats: .pdf, .docx, .xlsx, .pptx, .txt, .md, .log, .csv, .json, .xml, .html\n"

                    f"File: {path.name}"

                )



        except Exception as extract_error:

            return (

                f"❌ Error extracting text from {extension} file\n\n"

                f"Error: {type(extract_error).__name__}: {str(extract_error)}\n"

                f"File: {path.name}"

            )



        # Truncate if needed

        if not extracted_text:

            return (

                f"⚠️  No text extracted from file\n\n"

                f"File: {path.name}\n"

                f"Type: {extension}\n"

                f"Size: {file_size_kb:.2f} KB\n\n"

                f"The file may be empty, contain only images, or have other non-text content."

            )



        char_count = len(extracted_text)

        truncated = char_count > max_length

        display_text = extracted_text[:max_length] if truncated else extracted_text



        # Format result

        result_lines = [

            f"📄 Text extracted from: {path.name}",

            f"Type: {extension}",

            f"Size: {file_size_kb:.2f} KB",

        ]



        if page_count:

            count_label = "pages" if extension == '.pdf' else ("paragraphs" if extension == '.docx' else ("sheets" if extension == '.xlsx' else "slides"))

            result_lines.append(f"Content: {page_count} {count_label}")



        result_lines.append(f"Characters: {char_count:,}" + (f" (showing first {max_length:,})" if truncated else ""))

        result_lines.append("")

        result_lines.append("--- EXTRACTED TEXT ---")

        result_lines.append(display_text)



        if truncated:

            result_lines.append("")

            result_lines.append(f"... (truncated, {char_count - max_length:,} characters omitted)")



        return "\n".join(result_lines)



    except Exception as e:

        error_type = type(e).__name__

        return f"❌ Error extracting text: {error_type}: {str(e)}"






@mcp.tool()
async def send_email_with_attachments(
    to_address: str,
    subject: str,
    body: str,
    file_paths: str,
    body_type: str = "text",
    cc_addresses: str = "",
    importance: str = "normal",
    from_mailbox: str = "thomas@sixpillar.co.uk"
) -> str:
    """
    Send a new email message with file attachments.

    Args:
        to_address: Primary recipient email address (required)
        subject: Email subject line (required)
        body: Email body content (required)
        file_paths: Comma-separated list of file paths to attach (required)
        body_type: Body content type - "text" or "html" (default: "text")
        cc_addresses: Optional CC recipients (comma-separated for multiple)
        importance: Email importance - "low", "normal", or "high" (default: "normal")
        from_mailbox: Email address to send from (default: thomas@sixpillar.co.uk, use "me" for authenticated user)

    Size Limits:
        - 3MB per file (Microsoft Graph API limit)
        - 25MB total message size (Exchange Online limit)

    Returns:
        Success confirmation or error message
    """
    try:
        # Validate required fields
        if not to_address:
            return "❌ Error: to_address is required"

        if not subject:
            return "❌ Error: subject is required"

        if not body:
            return "❌ Error: body is required"

        if not file_paths:
            return "❌ Error: file_paths is required (comma-separated list of file paths)"

        # Basic email validation for primary recipient
        if '@' not in to_address or '.' not in to_address.split('@')[1]:
            return f"❌ Error: '{to_address}' doesn't appear to be a valid email address"

        # Validate body_type
        if body_type.lower() not in ["text", "html"]:
            return f"❌ Error: body_type must be 'text' or 'html', got '{body_type}'"

        # Validate importance
        importance_lower = importance.lower()
        if importance_lower not in ["low", "normal", "high"]:
            return f"❌ Error: importance must be 'low', 'normal', or 'high', got '{importance}'"

        # Parse file paths
        file_list = [fp.strip() for fp in file_paths.split(',')]
        if not file_list:
            return "❌ Error: No file paths provided"

        # Validate all files exist and check sizes
        attachments_list = []
        total_size = 0

        for file_path_str in file_list:
            if not file_path_str:
                continue

            file_path = Path(file_path_str)

            if not file_path.exists():
                return f"❌ Error: File not found: {file_path_str}"

            if not file_path.is_file():
                return f"❌ Error: Not a file: {file_path_str}"

            # Check file size (3MB limit per file)
            file_size = file_path.stat().st_size
            if file_size > 3 * 1024 * 1024:  # 3MB
                size_mb = file_size / (1024 * 1024)
                return f"❌ Error: File '{file_path.name}' is too large ({size_mb:.2f} MB). Maximum per file: 3 MB"

            total_size += file_size

            # Check total message size (25MB limit)
            if total_size > 25 * 1024 * 1024:  # 25MB
                size_mb = total_size / (1024 * 1024)
                return f"❌ Error: Total attachment size ({size_mb:.2f} MB) exceeds 25 MB limit"

            # Read file content (SDK expects bytes directly)
            with open(file_path, 'rb') as f:
                file_content = f.read()

            # Detect MIME type
            mime_type, _ = mimetypes.guess_type(file_path.name)
            if not mime_type:
                mime_type = 'application/octet-stream'

            # Create FileAttachment object
            attachment = FileAttachment()
            attachment.odata_type = "#microsoft.graph.fileAttachment"
            attachment.name = file_path.name
            attachment.content_type = mime_type
            attachment.content_bytes = file_content  # SDK expects bytes directly

            attachments_list.append(attachment)

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, from_mailbox)

        # Create the email message
        message = Message()
        message.subject = subject

        # Set body content
        message_body = ItemBody()
        message_body.content_type = BodyType.Html if body_type.lower() == "html" else BodyType.Text
        message_body.content = body
        message.body = message_body

        # Set importance
        message.importance = importance_lower

        # Add TO recipient
        to_recipient = Recipient()
        to_email = EmailAddress()
        to_email.address = to_address
        to_recipient.email_address = to_email
        message.to_recipients = [to_recipient]

        # Add CC recipients if provided
        if cc_addresses:
            cc_list = [addr.strip() for addr in cc_addresses.split(',')]
            cc_recipients = []

            for cc_addr in cc_list:
                if cc_addr:  # Skip empty strings
                    # Basic validation
                    if '@' not in cc_addr or '.' not in cc_addr.split('@')[1]:
                        return f"❌ Error: '{cc_addr}' doesn't appear to be a valid email address"

                    cc_recipient = Recipient()
                    cc_email = EmailAddress()
                    cc_email.address = cc_addr
                    cc_recipient.email_address = cc_email
                    cc_recipients.append(cc_recipient)

            if cc_recipients:
                message.cc_recipients = cc_recipients

        # Add attachments
        message.attachments = attachments_list

        # Create send mail request body
        send_body = SendMailPostRequestBody()
        send_body.message = message
        send_body.save_to_sent_items = True  # Save copy to Sent Items

        # Send the message using Graph API
        # POST /me/sendMail
        await asyncio.wait_for(
            mailbox.send_mail.post(body=send_body),
            timeout=API_TIMEOUT
        )

        cc_info = f"\nCc: {cc_addresses}" if cc_addresses else ""
        attachment_info = "\n".join([f"  - {att.name} ({att.content_type}, {len(att.content_bytes)/1024:.2f} KB)" for att in attachments_list])

        return (
            f"✅ Email sent successfully with {len(attachments_list)} attachment(s)!\n\n"
            f"To: {to_address}{cc_info}\n"
            f"Subject: {subject}\n"
            f"Body Type: {body_type}\n"
            f"Importance: {importance}\n\n"
            f"Attachments:\n{attachment_info}"
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error sending email: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time. The email may or may not have been sent."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error sending email with attachments: {error_type}: {str(e)}"


@mcp.tool()
async def add_attachment_to_draft(
    message_id: str,
    file_path: str,
    mailbox_id: str = "thomas@sixpillar.co.uk"
) -> str:
    """
    Add a file attachment to an existing draft email message.

    Args:
        message_id: The ID of the draft message to add attachment to (required)
        file_path: Path to the file to attach (required)
        mailbox_id: Email address of mailbox to access (default: "thomas@sixpillar.co.uk")
                   Use "me" for athena@'s own mailbox

    Size Limits:
        - 3MB per file (Microsoft Graph API limit for simple upload)
        - Files larger than 3MB require chunked upload (not yet supported)

    Returns:
        Success confirmation with attachment details or error message
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        if not file_path:
            return "❌ Error: file_path is required"

        # Validate file exists
        file = Path(file_path)

        if not file.exists():
            return f"❌ Error: File not found: {file_path}"

        if not file.is_file():
            return f"❌ Error: Not a file: {file_path}"

        # Check file size (3MB limit)
        file_size = file.stat().st_size
        if file_size > 3 * 1024 * 1024:  # 3MB
            size_mb = file_size / (1024 * 1024)
            return (
                f"❌ Error: File is too large ({size_mb:.2f} MB)\n\n"
                f"Maximum file size: 3 MB\n"
                f"For larger files, use chunked upload (not yet supported)"
            )

        # Read file content (SDK expects bytes directly)
        with open(file, 'rb') as f:
            file_content = f.read()

        # Detect MIME type
        mime_type, _ = mimetypes.guess_type(file.name)
        if not mime_type:
            mime_type = 'application/octet-stream'

        client = await get_graph_client()

        # Determine which mailbox to query
        if mailbox_id == "me":
            mailbox = client.me
        else:
            mailbox = client.users.by_user_id(mailbox_id)

        # Create FileAttachment object
        attachment = FileAttachment()
        attachment.odata_type = "#microsoft.graph.fileAttachment"
        attachment.name = file.name
        attachment.content_type = mime_type
        attachment.content_bytes = file_content  # SDK expects bytes directly

        # Add attachment to draft
        # POST /messages/{id}/attachments
        result = await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).attachments.post(body=attachment),
            timeout=API_TIMEOUT
        )

        size_kb = file_size / 1024

        return (
            f"✅ Attachment added to draft successfully!\n\n"
            f"Draft Message ID: {message_id}\n"
            f"Mailbox: {mailbox_id}\n\n"
            f"Attachment Details:\n"
            f"  Name: {file.name}\n"
            f"  Type: {mime_type}\n"
            f"  Size: {size_kb:.2f} KB\n"
            f"  Attachment ID: {result.id if result else 'N/A'}"
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error adding attachment: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error adding attachment to draft: {error_type}: {str(e)}"


@mcp.tool()
async def remove_attachment_from_draft(
    message_id: str,
    attachment_id: str,
    mailbox_id: str = "thomas@sixpillar.co.uk"
) -> str:
    """
    Remove a file attachment from an existing draft email message.

    Args:
        message_id: The ID of the draft message (required)
        attachment_id: The ID of the attachment to remove (required, get from list_attachments)
        mailbox_id: Email address of mailbox to access (default: "thomas@sixpillar.co.uk")
                   Use "me" for athena@'s own mailbox

    Returns:
        Success confirmation or error message
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        if not attachment_id:
            return "❌ Error: attachment_id is required (use list_attachments to get attachment ID)"

        client = await get_graph_client()

        # Determine which mailbox to query
        if mailbox_id == "me":
            mailbox = client.me
        else:
            mailbox = client.users.by_user_id(mailbox_id)

        # Delete attachment from draft
        # DELETE /messages/{message_id}/attachments/{attachment_id}
        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).attachments.by_attachment_id(attachment_id).delete(),
            timeout=API_TIMEOUT
        )

        return (
            f"✅ Attachment removed from draft successfully!\n\n"
            f"Draft Message ID: {message_id}\n"
            f"Attachment ID: {attachment_id}\n"
            f"Mailbox: {mailbox_id}"
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Error removing attachment: Timeout after {API_TIMEOUT}s\n\n"
            f"The Graph API did not respond in time."
        )
    except Exception as e:
        error_type = type(e).__name__
        return f"❌ Error removing attachment from draft: {error_type}: {str(e)}"


@mcp.tool()
async def add_attachment_base64(
    message_id: str,
    filename: str,
    base64_content: str,
    mailbox_id: str = "thomas@sixpillar.co.uk"
) -> str:
    """
    Attach file from laptop/workspace to draft email by passing base64-encoded content.

    **USE THIS TOOL WHEN:**
    - User requests: "attach this file from Documents", "attach the report I created"
    - File exists on laptop or workspace filesystem (not container)
    - Interactive workflow where user has the file locally

    **DO NOT USE THIS TOOL WHEN:**
    - File already in container (use add_attachment_to_draft instead)
    - File was just downloaded via download_attachment (use add_attachment_to_draft)
    - Server-side automation (use add_attachment_to_draft)

    **WORKFLOW:**
    1. Use Read() to get file content: content = Read("~/Documents/report.pdf")
    2. Call this tool: add_attachment_base64(draft_id, "report.pdf", content, mailbox_id)

    Args:
        message_id: The ID of the draft message (required)
        filename: Name for the attachment (required)
        base64_content: Base64-encoded file content (required)
        mailbox_id: Email address of mailbox (default: "thomas@sixpillar.co.uk", use "me" for athena@)

    Returns:
        Success confirmation or error message

    Limitations:
        - Max file size: 3MB (Graph API limit for simple upload)
        - Larger files require chunked upload (not yet supported)
        - Base64 encoding adds ~33% overhead

    Examples:
        - User: "Attach the quarterly-report.pdf from my Documents to the draft"
          → Read() file + add_attachment_base64()
        - User: "Attach this screenshot to the email"
          → Read() screenshot + add_attachment_base64()
    """
    try:
        client = await get_graph_client()

        # Determine which mailbox to use
        if mailbox_id == "me":
            mailbox = client.me
        else:
            mailbox = client.users.by_user_id(mailbox_id)

        # Validate filename (security)
        safe_filename = Path(filename).name  # Strip directory components
        if not safe_filename or safe_filename.startswith('.'):
            return (
                f"❌ Invalid filename: {filename}\n\n"
                f"Filename must be a valid file name without directory paths."
            )

        # Create FileAttachment object
        attachment = FileAttachment()
        attachment.name = safe_filename
        attachment.content_type = mimetypes.guess_type(safe_filename)[0] or "application/octet-stream"

        # Decode base64 content for size check
        try:
            content_bytes = base64.b64decode(base64_content)
            size_mb = len(content_bytes) / (1024 * 1024)

            if size_mb > 3:
                return (
                    f"❌ File too large: {size_mb:.2f}MB\n\n"
                    f"Maximum file size for simple upload: 3MB\n"
                    f"File: {safe_filename}\n\n"
                    f"Chunked upload for larger files is not yet implemented."
                )
        except Exception as e:
            return (
                f"❌ Invalid base64 content\n\n"
                f"Error decoding base64: {str(e)}"
            )

        # Re-encode as bytes for Graph API (Graph expects base64-encoded bytes)
        attachment.content_bytes = base64_content.encode('utf-8')

        # Add attachment to draft
        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).attachments.post(attachment),
            timeout=API_TIMEOUT
        )

        return (
            f"✅ Attachment added successfully!\n\n"
            f"Draft ID: {message_id}\n"
            f"File: {safe_filename}\n"
            f"Size: {size_mb:.2f}MB\n"
            f"Type: {attachment.content_type}\n\n"
            f"The attachment has been added to the draft email."
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Timeout adding attachment\n\n"
            f"Request timed out after {API_TIMEOUT} seconds.\n"
            f"The file may be too large or network connection slow."
        )
    except Exception as e:
        error_msg = str(e)
        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Draft not found\n\n"
                f"Message ID: {message_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure the message ID is correct and is a draft (not sent)."
            )
        return (
            f"❌ Failed to add attachment\n\n"
            f"File: {filename}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_mailbox_statistics(
    mailbox_id: str
) -> str:
    """
    Get comprehensive mailbox-wide folder statistics for monitoring and reporting.

    Returns statistics for all mail folders in a mailbox, including:
    - Folder name and ID
    - Total item count
    - Unread item count
    - Child folder count
    - Hidden status

    Useful for:
    - Monitoring inbox zero progress
    - Detecting unusual email volume (alerts, spam waves)
    - Reporting on folder usage and organization health
    - Workflow automation decisions based on folder state

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        Formatted statistics for all folders in the mailbox

    Example:
        get_mailbox_statistics(mailbox_id="me")
        get_mailbox_statistics(mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        client = await get_graph_client()

        # Determine which mailbox to use
        if mailbox_id == "me":
            mailbox = client.me
        else:
            mailbox = client.users.by_user_id(mailbox_id)

        # Get all mail folders (including hidden folders)
        folders_response = await asyncio.wait_for(
            mailbox.mail_folders.get(),
            timeout=API_TIMEOUT
        )

        if not folders_response or not folders_response.value:
            return (
                f"📭 No folders found in mailbox\n\n"
                f"Mailbox: {mailbox_id}"
            )

        folders = folders_response.value

        # Build statistics report
        result_lines = []
        result_lines.append(f"📊 MAILBOX STATISTICS")
        result_lines.append(f"Mailbox: {mailbox_id}")
        result_lines.append(f"Total Folders: {len(folders)}")
        result_lines.append("")
        result_lines.append("=" * 80)

        # Calculate totals
        total_items = sum(f.total_item_count or 0 for f in folders)
        total_unread = sum(f.unread_item_count or 0 for f in folders)

        result_lines.append(f"Overall Totals:")
        result_lines.append(f"  Total Messages: {total_items:,}")
        result_lines.append(f"  Unread Messages: {total_unread:,}")
        result_lines.append("")
        result_lines.append("=" * 80)
        result_lines.append("")

        # Sort folders by total item count (descending)
        sorted_folders = sorted(folders, key=lambda f: f.total_item_count or 0, reverse=True)

        # Display each folder
        for folder in sorted_folders:
            display_name = folder.display_name or "Unknown"
            total = folder.total_item_count or 0
            unread = folder.unread_item_count or 0
            children = folder.child_folder_count or 0
            hidden = " [HIDDEN]" if folder.is_hidden else ""

            result_lines.append(f"📁 {display_name}{hidden}")
            result_lines.append(f"   Total: {total:,} | Unread: {unread:,} | Child Folders: {children}")
            result_lines.append(f"   ID: {folder.id}")
            result_lines.append("")

        return "\n".join(result_lines)

    except asyncio.TimeoutError:
        return (
            f"❌ Timeout retrieving mailbox statistics\n\n"
            f"Request timed out after {API_TIMEOUT} seconds.\n"
            f"Mailbox: {mailbox_id}"
        )
    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Mailbox not found\n\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure the mailbox ID is correct and you have access rights."
            )

        return (
            f"❌ Error retrieving mailbox statistics: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_folder_statistics(
    mailbox_id: str,
    folder_name: str
) -> str:
    """
    Get detailed statistics for a specific mail folder.

    Returns comprehensive stats for a single folder including:
    - Folder name, ID, and parent folder
    - Total item count
    - Unread item count
    - Child folder count
    - Hidden status

    Useful for:
    - Monitoring specific folder health (Inbox, Sent Items, etc.)
    - Checking folder state before automation decisions
    - Detailed folder analysis

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        folder_name: Name of folder to get statistics for (required)
                    Common folders: "Inbox", "Sent Items", "Drafts", "Deleted Items", "Archive"
                    Can also use folder ID instead of name

    Returns:
        Detailed statistics for the specified folder

    Example:
        get_folder_statistics(mailbox_id="me", folder_name="Inbox")
        get_folder_statistics(mailbox_id="thomas@sixpillar.co.uk", folder_name="Archive")
    """
    try:
        client = await get_graph_client()

        # Determine which mailbox to use
        if mailbox_id == "me":
            mailbox = client.me
        else:
            mailbox = client.users.by_user_id(mailbox_id)

        # Try to get folder by name first (more common use case)
        # Graph API supports well-known folder names like "inbox", "drafts", etc.
        try:
            folder = await asyncio.wait_for(
                mailbox.mail_folders.by_mail_folder_id(folder_name).get(),
                timeout=API_TIMEOUT
            )
        except Exception:
            # If folder_name failed, try searching all folders
            folders_response = await asyncio.wait_for(
                mailbox.mail_folders.get(),
                timeout=API_TIMEOUT
            )

            if not folders_response or not folders_response.value:
                return (
                    f"❌ Folder not found: {folder_name}\n\n"
                    f"Mailbox: {mailbox_id}\n\n"
                    f"No folders found in mailbox."
                )

            # Search for folder by display name (case-insensitive)
            folder = None
            for f in folders_response.value:
                if f.display_name and f.display_name.lower() == folder_name.lower():
                    folder = f
                    break

            if not folder:
                available_folders = [f.display_name for f in folders_response.value if f.display_name]
                return (
                    f"❌ Folder not found: {folder_name}\n\n"
                    f"Mailbox: {mailbox_id}\n\n"
                    f"Available folders:\n" +
                    "\n".join(f"  - {name}" for name in sorted(available_folders))
                )

        # Build detailed statistics report
        result_lines = []
        result_lines.append(f"📊 FOLDER STATISTICS")
        result_lines.append("=" * 80)
        result_lines.append("")
        result_lines.append(f"Folder: {folder.display_name or 'Unknown'}")
        result_lines.append(f"Mailbox: {mailbox_id}")
        result_lines.append(f"Folder ID: {folder.id}")

        if folder.parent_folder_id:
            result_lines.append(f"Parent Folder ID: {folder.parent_folder_id}")

        if folder.is_hidden:
            result_lines.append(f"Status: HIDDEN")

        result_lines.append("")
        result_lines.append("=" * 80)
        result_lines.append("")

        # Statistics
        total = folder.total_item_count or 0
        unread = folder.unread_item_count or 0
        read = total - unread
        children = folder.child_folder_count or 0

        result_lines.append("📈 Item Counts:")
        result_lines.append(f"  Total Items:   {total:,}")
        result_lines.append(f"  Unread Items:  {unread:,}")
        result_lines.append(f"  Read Items:    {read:,}")

        if total > 0:
            unread_pct = (unread / total) * 100
            result_lines.append(f"  Unread %:      {unread_pct:.1f}%")

        result_lines.append("")
        result_lines.append(f"📁 Child Folders: {children}")

        return "\n".join(result_lines)

    except asyncio.TimeoutError:
        return (
            f"❌ Timeout retrieving folder statistics\n\n"
            f"Request timed out after {API_TIMEOUT} seconds.\n"
            f"Mailbox: {mailbox_id}\n"
            f"Folder: {folder_name}"
        )
    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Folder or mailbox not found\n\n"
                f"Mailbox: {mailbox_id}\n"
                f"Folder: {folder_name}\n\n"
                f"Ensure both the mailbox and folder exist and you have access rights."
            )

        return (
            f"❌ Error retrieving folder statistics: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Folder: {folder_name}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def list_categories(mailbox_id: str) -> str:
    """
    List all available email categories (Outlook labels/tags) for a mailbox.

    Categories are user-defined or system-defined labels that can be applied to emails
    for organization and filtering. Each category has a name and color.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for authenticated user's mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        Formatted list of available categories with names and colors

    Examples:
        list_categories(mailbox_id="me")
        list_categories(mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        client = await get_graph_client()

        # Determine which mailbox to use
        if mailbox_id == "me":
            mailbox = client.me
        else:
            mailbox = client.users.by_user_id(mailbox_id)

        # Get master categories list
        # GET /users/{id}/outlook/masterCategories
        categories_response = await asyncio.wait_for(
            mailbox.outlook.master_categories.get(),
            timeout=API_TIMEOUT
        )

        if not categories_response or not categories_response.value:
            return (
                f"📋 No categories found\n\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"You can create categories in Outlook to organize your emails."
            )

        result_lines = []
        result_lines.append(f"📋 EMAIL CATEGORIES ({len(categories_response.value)})")
        result_lines.append(f"Mailbox: {mailbox_id}")
        result_lines.append("=" * 80)
        result_lines.append("")

        for idx, category in enumerate(categories_response.value, 1):
            name = category.display_name or "Unknown"
            color = category.color or "None"
            result_lines.append(f"{idx}. {name}")
            result_lines.append(f"   Color: {color}")
            if category.id:
                result_lines.append(f"   ID: {category.id}")
            result_lines.append("")

        result_lines.append("💡 Use apply_categories to tag emails with these categories")

        return "\n".join(result_lines)

    except asyncio.TimeoutError:
        return (
            f"❌ Timeout listing categories\n\n"
            f"Request timed out after {API_TIMEOUT} seconds.\n"
            f"Mailbox: {mailbox_id}"
        )
    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Mailbox not found\n\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure the mailbox exists and you have access rights."
            )

        return (
            f"❌ Error listing categories: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def create_category(mailbox_id: str, display_name: str, color: str = "preset1") -> str:
    """
    Create a new email category (Outlook label/tag) in a mailbox.

    Creates a new master category that can be applied to emails for organization.
    Each category has a display name and color.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for authenticated user's mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        display_name: Name for the new category (required)
                     Example: "Error", "Urgent", "Project Alpha"
        color: Color preset for the category (default: "preset1" - red)
              Available: preset0-preset24
              Common: preset1 (red), preset2 (orange), preset3 (brown),
                     preset6 (yellow), preset7 (green), preset8 (teal),
                     preset9 (blue), preset10 (purple)

    Returns:
        Success confirmation with category details

    Examples:
        create_category(mailbox_id="me", display_name="Error")
        create_category(mailbox_id="thomas@sixpillar.co.uk", display_name="Urgent", color="preset1")
    """
    try:
        if not display_name:
            return "❌ Error: display_name is required"

        client = await get_graph_client()

        # Determine which mailbox to use
        if mailbox_id == "me":
            mailbox = client.me
        else:
            mailbox = client.users.by_user_id(mailbox_id)

        # Create new category object
        new_category = OutlookCategory()
        new_category.display_name = display_name
        new_category.color = color

        # POST /users/{id}/outlook/masterCategories
        created_category = await asyncio.wait_for(
            mailbox.outlook.master_categories.post(body=new_category),
            timeout=API_TIMEOUT
        )

        if not created_category:
            return (
                f"❌ Failed to create category\n\n"
                f"Mailbox: {mailbox_id}\n"
                f"Category: {display_name}"
            )

        return (
            f"✅ Category created successfully!\n\n"
            f"Mailbox: {mailbox_id}\n\n"
            f"Category Details:\n"
            f"  Name: {created_category.display_name}\n"
            f"  Color: {created_category.color}\n"
            f"  ID: {created_category.id}\n\n"
            f"💡 Use apply_categories to tag emails with this category"
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Timeout creating category\n\n"
            f"Request timed out after {API_TIMEOUT} seconds.\n"
            f"Mailbox: {mailbox_id}\n"
            f"Category: {display_name}"
        )
    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        # Check for duplicate category error
        if "already exists" in error_msg.lower() or "conflict" in error_msg.lower():
            return (
                f"⚠️ Category already exists\n\n"
                f"Mailbox: {mailbox_id}\n"
                f"Category: {display_name}\n\n"
                f"💡 Use list_categories to see all existing categories"
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Mailbox not found\n\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure the mailbox exists and you have access rights."
            )

        return (
            f"❌ Error creating category: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Category: {display_name}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_message_categories(message_id: str, mailbox_id: str) -> str:
    """
    Get the categories currently applied to a specific email message.

    Args:
        message_id: The ID of the email message (required)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for authenticated user's mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        List of categories applied to the message, or empty if none

    Examples:
        get_message_categories(message_id="AAMk...", mailbox_id="me")
        get_message_categories(message_id="AAMk...", mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # Use helper function with retry logic (same pattern as get_email_by_id)
        message = await fetch_message_by_id(mailbox, message_id)

        if not message:
            return (
                f"❌ Message not found\n\n"
                f"Message ID: {message_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        subject = message.subject or "No subject"
        categories = message.categories or []

        result_lines = []
        result_lines.append(f"📧 EMAIL CATEGORIES")
        result_lines.append("=" * 80)
        result_lines.append(f"Subject: {subject}")
        result_lines.append(f"Message ID: {message_id}")
        result_lines.append(f"Mailbox: {mailbox_id}")
        result_lines.append("")

        if not categories:
            result_lines.append("🏷️  No categories applied")
            result_lines.append("")
            result_lines.append("💡 Use apply_categories to tag this email")
        else:
            result_lines.append(f"🏷️  Categories ({len(categories)}):")
            for idx, category in enumerate(categories, 1):
                result_lines.append(f"  {idx}. {category}")

        return "\n".join(result_lines)

    except asyncio.TimeoutError:
        return (
            f"❌ Timeout getting message categories\n\n"
            f"Request timed out after {API_TIMEOUT} seconds.\n"
            f"Message ID: {message_id}\n"
            f"Mailbox: {mailbox_id}"
        )
    except Exception as e:
        import traceback
        error_type = type(e).__name__
        error_msg = str(e)

        # Log full traceback to container logs for debugging
        print(f"\n{'='*80}", file=sys.stderr, flush=True)
        print(f"ERROR in get_message_categories:", file=sys.stderr, flush=True)
        print(f"{'='*80}", file=sys.stderr, flush=True)
        traceback.print_exc(file=sys.stderr)
        print(f"{'='*80}\n", file=sys.stderr, flush=True)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Message or mailbox not found\n\n"
                f"Message ID: {message_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure both exist and you have access rights."
            )

        return (
            f"❌ Error getting message categories: {error_type}\n\n"
            f"Message ID: {message_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def apply_categories(
    message_id: str,
    categories: str,
    mailbox_id: str,
    replace: bool = True
) -> str:
    """
    Apply categories (labels/tags) to an email message.

    Categories are user-defined or system-defined labels for organizing emails.
    You can apply one or multiple categories to help organize and filter messages.

    Args:
        message_id: The ID of the email message (required)
        categories: Category names to apply, comma-separated (required)
                   Example: "Red category" or "Red category,Blue category,Important"
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for authenticated user's mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        replace: If True, replaces existing categories; if False, adds to existing (default: True)

    Returns:
        Success confirmation with applied categories

    Examples:
        apply_categories(message_id="AAMk...", categories="Red category", mailbox_id="me")
        apply_categories(message_id="AAMk...", categories="Important,Follow up", mailbox_id="me")
        apply_categories(message_id="AAMk...", categories="Work", mailbox_id="me", replace=False)
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"
        if not categories:
            return "❌ Error: categories is required"

        # Parse comma-separated categories
        category_list = [cat.strip() for cat in categories.split(',') if cat.strip()]
        if not category_list:
            return "❌ Error: At least one category name is required"

        client = await get_graph_client()
        mailbox = get_mailbox_endpoint(client, mailbox_id)

        # If not replacing, get existing categories first
        if not replace:
            # Use helper function with retry logic (same pattern as get_email_by_id)
            existing_message = await fetch_message_by_id(mailbox, message_id)
            existing_categories = existing_message.categories or []
            # Merge existing and new categories (remove duplicates)
            category_list = list(set(existing_categories + category_list))

        # Update message with new categories
        # PATCH /users/{id}/messages/{id}
        message = Message()
        message.categories = category_list

        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).patch(body=message),
            timeout=API_TIMEOUT
        )

        action = "replaced with" if replace else "added to message, total"
        return (
            f"✅ Categories {action} successfully!\n\n"
            f"Message ID: {message_id}\n"
            f"Mailbox: {mailbox_id}\n\n"
            f"Applied Categories ({len(category_list)}):\n" +
            "\n".join(f"  - {cat}" for cat in category_list)
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Timeout applying categories\n\n"
            f"Request timed out after {API_TIMEOUT} seconds.\n"
            f"Message ID: {message_id}\n"
            f"Mailbox: {mailbox_id}"
        )
    except Exception as e:
        import traceback
        error_type = type(e).__name__
        error_msg = str(e)

        # Log full traceback to container logs for debugging
        print(f"\n{'='*80}", file=sys.stderr, flush=True)
        print(f"ERROR in apply_categories:", file=sys.stderr, flush=True)
        print(f"{'='*80}", file=sys.stderr, flush=True)
        traceback.print_exc(file=sys.stderr)
        print(f"{'='*80}\n", file=sys.stderr, flush=True)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Message or mailbox not found\n\n"
                f"Message ID: {message_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure both exist and you have access rights."
            )

        if "InvalidPermission" in error_msg or "Forbidden" in error_msg:
            return (
                f"❌ Permission denied\n\n"
                f"Message ID: {message_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure MailboxSettings.ReadWrite permission is granted and consented.\n"
                f"Error: {error_msg}"
            )

        return (
            f"❌ Error applying categories: {error_type}\n\n"
            f"Message ID: {message_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Categories: {categories}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def remove_all_categories(message_id: str, mailbox_id: str) -> str:
    """
    Remove all categories from an email message.

    Clears all category labels/tags from the specified message.

    Args:
        message_id: The ID of the email message (required)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for authenticated user's mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        Success confirmation

    Examples:
        remove_all_categories(message_id="AAMk...", mailbox_id="me")
        remove_all_categories(message_id="AAMk...", mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        if not message_id:
            return "❌ Error: message_id is required"

        client = await get_graph_client()

        # Determine which mailbox to use
        if mailbox_id == "me":
            mailbox = client.me
        else:
            mailbox = client.users.by_user_id(mailbox_id)

        # Update message with empty categories array
        # PATCH /users/{id}/messages/{id}
        message = Message()
        message.categories = []

        await asyncio.wait_for(
            mailbox.messages.by_message_id(message_id).patch(body=message),
            timeout=API_TIMEOUT
        )

        return (
            f"✅ All categories removed successfully!\n\n"
            f"Message ID: {message_id}\n"
            f"Mailbox: {mailbox_id}\n\n"
            f"The message no longer has any categories applied."
        )

    except asyncio.TimeoutError:
        return (
            f"❌ Timeout removing categories\n\n"
            f"Request timed out after {API_TIMEOUT} seconds.\n"
            f"Message ID: {message_id}\n"
            f"Mailbox: {mailbox_id}"
        )
    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Message or mailbox not found\n\n"
                f"Message ID: {message_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure both exist and you have access rights."
            )

        if "InvalidPermission" in error_msg or "Forbidden" in error_msg:
            return (
                f"❌ Permission denied\n\n"
                f"Message ID: {message_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure MailboxSettings.ReadWrite permission is granted and consented.\n"
                f"Error: {error_msg}"
            )

        return (
            f"❌ Error removing categories: {error_type}\n\n"
            f"Message ID: {message_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )



#
# ============================================================================
# FOLDER MANAGEMENT HELPER FUNCTIONS
# ============================================================================
#

async def resolve_folder_id(
    folder_identifier: str,
    mailbox_id: str,
    graph_client: GraphServiceClient
) -> tuple[str, str]:
    """
    Resolve any folder identifier (ID, name, path, well-known) to folder ID and name.

    Args:
        folder_identifier: Folder ID, name, path (Inbox/Projects), or well-known name
        mailbox_id: Mailbox to search in
        graph_client: Graph API client

    Returns:
        Tuple of (folder_id, display_name)

    Raises:
        ValueError: If folder not found or ambiguous
    """
    # Get proper mailbox endpoint (handles "me" vs delegated mailbox)
    mailbox = get_mailbox_endpoint(graph_client, mailbox_id)

    # Case 1: Already a folder ID (starts with AAMk or AQMk)
    if folder_identifier.startswith("AAMk") or folder_identifier.startswith("AQMk"):
        try:
            folder = await mailbox.mail_folders.by_mail_folder_id(folder_identifier).get()
            if folder and folder.display_name:
                return (folder_identifier, folder.display_name)
        except Exception:
            raise ValueError(f"Folder ID not found: {folder_identifier}")

    # Case 2: Well-known folder names (inbox, sentitems, drafts, etc.)
    # Note: Well-known names work reliably for /me but may fail for delegated
    # mailbox access (/users/{id}). If this fails, we fall through to Case 4
    # which searches by display name.
    well_known_folders = {
        "inbox": "inbox",
        "sent": "sentitems",
        "sentitems": "sentitems",
        "drafts": "drafts",
        "deleteditems": "deleteditems",
        "deleted": "deleteditems",
        "trash": "deleteditems",
        "junkemail": "junkemail",
        "junk": "junkemail",
        "spam": "junkemail",
        "archive": "archive",
        "outbox": "outbox",
        "conversationhistory": "conversationhistory",
    }

    if folder_identifier.lower() in well_known_folders:
        well_known_name = well_known_folders[folder_identifier.lower()]
        try:
            folder = await mailbox.mail_folders.by_mail_folder_id(well_known_name).get()
            if folder and folder.display_name:
                return (folder.id, folder.display_name)
        except Exception:
            # Well-known folder API failed (common for delegated mailbox access)
            # Fall through to Case 4 (display name search) instead of failing
            pass

    # Case 3: Path notation (Inbox/Projects/ClientA)
    if "/" in folder_identifier:
        return await resolve_folder_path(folder_identifier, mailbox_id, graph_client)

    # Case 4: Display name (search all folders recursively)
    matches = await search_folders_by_name(folder_identifier, mailbox_id, graph_client)

    if len(matches) == 0:
        raise ValueError(f"Folder not found: '{folder_identifier}'")
    if len(matches) > 1:
        paths = [m['path'] for m in matches]
        raise ValueError(
            f"Ambiguous folder name '{folder_identifier}' - {len(matches)} matches found:\n" +
            "\n".join(f"  - {p}" for p in paths) +
            "\n\nUse full path notation or folder ID to be specific."
        )

    return (matches[0]['id'], matches[0]['name'])


async def resolve_folder_path(
    path: str,
    mailbox_id: str,
    graph_client: GraphServiceClient
) -> tuple[str, str]:
    """
    Resolve folder path like 'Inbox/Projects/ClientA' to folder ID and name.

    Args:
        path: Folder path with / separators
        mailbox_id: Mailbox to search in
        graph_client: Graph API client

    Returns:
        Tuple of (folder_id, display_name)

    Raises:
        ValueError: If path invalid or folder not found
    """
    parts = path.split("/")
    if not parts or parts[0] == "":
        raise ValueError(f"Invalid folder path: '{path}'")

    # Resolve first part (root level)
    current_id, current_name = await resolve_folder_id(parts[0], mailbox_id, graph_client)

    # Get proper mailbox endpoint
    mailbox = get_mailbox_endpoint(graph_client, mailbox_id)

    # Traverse each subsequent level
    for i, part in enumerate(parts[1:], start=1):
        if not part:
            raise ValueError(f"Invalid folder path (empty segment): '{path}'")

        # Get child folders of current folder
        children = await mailbox.mail_folders.by_mail_folder_id(current_id).child_folders.get()

        if not children or not children.value:
            raise ValueError(f"No child folders found in path '{path}' at '{'/'.join(parts[:i])}'")

        # Find matching child
        matches = [f for f in children.value if f.display_name == part]

        if len(matches) == 0:
            raise ValueError(f"Folder '{part}' not found in path '{path}' under '{'/'.join(parts[:i])}'")
        if len(matches) > 1:
            raise ValueError(f"Multiple '{part}' folders found in path '{path}' under '{'/'.join(parts[:i])}'")

        current_id = matches[0].id
        current_name = matches[0].display_name

    return (current_id, current_name)


async def search_folders_by_name(
    name: str,
    mailbox_id: str,
    graph_client: GraphServiceClient,
    parent_id: Optional[str] = None,
    current_path: str = ""
) -> List[Dict[str, str]]:
    """
    Recursively search for folders by display name.

    Args:
        name: Folder display name to search for
        mailbox_id: Mailbox to search in
        graph_client: Graph API client
        parent_id: Parent folder ID (None for root)
        current_path: Current path for building full paths

    Returns:
        List of dicts with 'id', 'name', 'path' keys
    """
    matches = []

    # Get proper mailbox endpoint
    mailbox = get_mailbox_endpoint(graph_client, mailbox_id)

    # Get folders at current level
    if parent_id:
        folders_response = await mailbox.mail_folders.by_mail_folder_id(parent_id).child_folders.get()
    else:
        folders_response = await mailbox.mail_folders.get()

    if not folders_response or not folders_response.value:
        return matches

    for folder in folders_response.value:
        folder_path = f"{current_path}/{folder.display_name}" if current_path else folder.display_name

        # Check if this folder matches
        if folder.display_name == name:
            matches.append({
                'id': folder.id,
                'name': folder.display_name,
                'path': folder_path
            })

        # Recursively search child folders
        child_matches = await search_folders_by_name(
            name, mailbox_id, graph_client,
            parent_id=folder.id,
            current_path=folder_path
        )
        matches.extend(child_matches)

    return matches


async def is_system_folder(
    folder_id: str,
    mailbox_id: str,
    graph_client: GraphServiceClient
) -> bool:
    """
    Check if folder is a protected system folder.

    System folders cannot be renamed or deleted.

    Args:
        folder_id: Folder ID to check
        mailbox_id: Mailbox containing folder
        graph_client: Graph API client

    Returns:
        True if system folder, False if custom folder
    """
    try:
        mailbox = get_mailbox_endpoint(graph_client, mailbox_id)
        folder = await mailbox.mail_folders.by_mail_folder_id(folder_id).get()

        # Check if folder has a well-known name property
        # System folders have this property set
        if hasattr(folder, 'well_known_name') and folder.well_known_name:
            return True

        # Also check against known system folder names as fallback
        system_folder_names = {
            "Inbox", "Sent Items", "Drafts", "Deleted Items",
            "Junk Email", "Outbox", "Archive", "Conversation History",
            "Clutter", "Notes", "Journal", "Tasks", "Contacts", "Calendar"
        }

        if folder.display_name in system_folder_names:
            return True

        return False
    except Exception:
        # If we can't determine, assume it's a system folder (safe default)
        return True


async def get_folder_stats(
    folder_id: str,
    mailbox_id: str,
    graph_client: GraphServiceClient
) -> Dict[str, int]:
    """
    Get folder statistics (item counts, child folder counts).

    Args:
        folder_id: Folder ID
        mailbox_id: Mailbox containing folder
        graph_client: Graph API client

    Returns:
        Dict with 'total_items', 'unread_items', 'child_folders' keys
    """
    try:
        mailbox = get_mailbox_endpoint(graph_client, mailbox_id)
        folder = await mailbox.mail_folders.by_mail_folder_id(folder_id).get()

        # Get child folder count
        children = await mailbox.mail_folders.by_mail_folder_id(folder_id).child_folders.get()
        child_count = len(children.value) if children and children.value else 0

        return {
            'total_items': folder.total_item_count or 0,
            'unread_items': folder.unread_item_count or 0,
            'child_folders': child_count
        }
    except Exception:
        return {
            'total_items': 0,
            'unread_items': 0,
            'child_folders': 0
        }


#
# ============================================================================
# FOLDER MANAGEMENT MCP TOOLS
# ============================================================================
#

@mcp.tool()
async def list_folders(
    mailbox_id: str,
    parent_folder: str = None,
    show_hidden: bool = False,
    format: str = "tree"
) -> str:
    """
    List mail folders in mailbox hierarchy.

    Args:
        mailbox_id: Mailbox to access ("me" or email address) - REQUIRED
        parent_folder: Optional parent folder (ID/name/path) to list children only
        show_hidden: Include hidden folders (default: False)
        format: Output format - "tree" for hierarchical or "flat" for simple list

    Returns:
        Formatted folder list with names, IDs, item counts, and hierarchy

    Examples:
        list_folders(mailbox_id="me")  # All folders in tree format
        list_folders(mailbox_id="me", parent_folder="Inbox", format="flat")
        list_folders(mailbox_id="thomas@sixpillar.co.uk", format="tree")
    """
    try:
        graph_client = await get_graph_client()
        mailbox = get_mailbox_endpoint(graph_client, mailbox_id)

        # If parent_folder specified, resolve it first
        if parent_folder:
            parent_id, parent_name = await resolve_folder_id(parent_folder, mailbox_id, graph_client)
            folders_response = await mailbox.mail_folders.by_mail_folder_id(parent_id).child_folders.get()
            root_display = f"Folders under '{parent_name}'"
        else:
            folders_response = await mailbox.mail_folders.get()
            root_display = "All Mail Folders"

        if not folders_response or not folders_response.value:
            return f"📁 {root_display}\nMailbox: {mailbox_id}\n\nNo folders found."

        folders = folders_response.value

        # Filter hidden folders if requested
        if not show_hidden:
            folders = [f for f in folders if not (hasattr(f, 'is_hidden') and f.is_hidden)]

        if format == "tree":
            output = f"📁 {root_display} (TREE VIEW)\n"
            output += f"Mailbox: {mailbox_id}\n"
            output += "=" * 80 + "\n\n"

            async def format_folder_tree(folder, indent=0):
                stats = await get_folder_stats(folder.id, mailbox_id, graph_client)
                is_system = await is_system_folder(folder.id, mailbox_id, graph_client)

                prefix = "  " * indent + ("📌 " if is_system else "📁 ")
                folder_type = "[SYSTEM]" if is_system else "[CUSTOM]"

                line = f"{prefix}{folder.display_name} {folder_type}\n"
                line += "  " * indent + f"   Items: {stats['total_items']} ({stats['unread_items']} unread)"
                if stats['child_folders'] > 0:
                    line += f" | Subfolders: {stats['child_folders']}"
                line += f"\n   ID: {folder.id}\n"

                # Get and format child folders
                children = await mailbox.mail_folders.by_mail_folder_id(folder.id).child_folders.get()
                if children and children.value:
                    for child in children.value:
                        if show_hidden or not (hasattr(child, 'is_hidden') and child.is_hidden):
                            line += await format_folder_tree(child, indent + 1)

                return line

            for folder in folders:
                output += await format_folder_tree(folder)

            return output.rstrip()

        else:  # flat format
            output = f"📁 {root_display} (FLAT LIST)\n"
            output += f"Mailbox: {mailbox_id}\n"
            output += "=" * 80 + "\n\n"

            for i, folder in enumerate(folders, 1):
                stats = await get_folder_stats(folder.id, mailbox_id, graph_client)
                is_system = await is_system_folder(folder.id, mailbox_id, graph_client)
                folder_type = "[SYSTEM]" if is_system else "[CUSTOM]"

                output += f"{i}. {folder.display_name} {folder_type}\n"
                output += f"   Items: {stats['total_items']} ({stats['unread_items']} unread)"
                if stats['child_folders'] > 0:
                    output += f" | Subfolders: {stats['child_folders']}"
                output += f"\n   ID: {folder.id}\n\n"

            return output.rstrip()

    except ValueError as e:
        return f"❌ Folder resolution error: {str(e)}"
    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error listing folders: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def create_folder(
    folder_name: str,
    mailbox_id: str,
    parent_folder: str = None
) -> str:
    """
    Create new mail folder.

    Args:
        folder_name: Name for new folder - REQUIRED
        mailbox_id: Mailbox to create in ("me" or email address) - REQUIRED
        parent_folder: Parent folder (ID/name/path) - None for root level

    Returns:
        Folder ID and confirmation message

    Examples:
        create_folder(folder_name="Projects", mailbox_id="me")  # Root level
        create_folder(folder_name="ClientA", mailbox_id="me", parent_folder="Inbox/Projects")
        create_folder(folder_name="Archive2024", mailbox_id="thomas@sixpillar.co.uk", parent_folder="Archive")

    Restrictions:
        - Folder name must be unique within parent folder
        - Folder name cannot contain: / \\ : * ? " < > |
    """
    try:
        graph_client = await get_graph_client()
        mailbox = get_mailbox_endpoint(graph_client, mailbox_id)

        # Create folder object
        new_folder = MailFolder()
        new_folder.display_name = folder_name

        # Create folder at appropriate location
        if parent_folder:
            # Resolve parent folder ID
            parent_id, parent_name = await resolve_folder_id(parent_folder, mailbox_id, graph_client)

            # Create as child of parent
            created_folder = await mailbox.mail_folders.by_mail_folder_id(parent_id).child_folders.post(new_folder)

            location = f"under '{parent_name}'"
        else:
            # Create at root level
            created_folder = await mailbox.mail_folders.post(new_folder)

            location = "at root level"

        return (
            f"✅ Folder created successfully\n\n"
            f"Name: {created_folder.display_name}\n"
            f"Location: {location}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Folder ID: {created_folder.id}\n\n"
            f"💡 Use this folder ID or name to move emails into it."
        )

    except ValueError as e:
        return f"❌ Folder resolution error: {str(e)}"
    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "already exists" in error_msg.lower() or "duplicate" in error_msg.lower():
            return (
                f"❌ Folder already exists\n\n"
                f"Folder name: {folder_name}\n"
                f"Parent: {parent_folder or 'root level'}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Folder names must be unique within the same parent folder."
            )

        return (
            f"❌ Error creating folder: {error_type}\n\n"
            f"Folder name: {folder_name}\n"
            f"Parent: {parent_folder or 'root level'}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def rename_folder(
    folder_id: str,
    new_name: str,
    mailbox_id: str
) -> str:
    """
    Rename existing mail folder.

    Args:
        folder_id: Folder to rename (ID, name, or path) - REQUIRED
        new_name: New display name for folder - REQUIRED
        mailbox_id: Mailbox containing folder ("me" or email address) - REQUIRED

    Returns:
        Confirmation with old and new names

    Examples:
        rename_folder(folder_id="Projects", new_name="ActiveProjects", mailbox_id="me")
        rename_folder(folder_id="Inbox/Old", new_name="Inbox/Archive2023", mailbox_id="me")
        rename_folder(folder_id="AAMkADU3...", new_name="NewName", mailbox_id="thomas@sixpillar.co.uk")

    Restrictions:
        - Cannot rename system folders (Inbox, Sent Items, Drafts, etc.)
        - New name must be unique within parent folder
        - Folder ID cannot be a well-known name (use display name or path for custom folders)
    """
    try:
        graph_client = await get_graph_client()
        mailbox = get_mailbox_endpoint(graph_client, mailbox_id)

        # Resolve folder ID
        resolved_id, old_name = await resolve_folder_id(folder_id, mailbox_id, graph_client)

        # Check if system folder
        if await is_system_folder(resolved_id, mailbox_id, graph_client):
            return (
                f"❌ Cannot rename system folder\n\n"
                f"Folder: {old_name}\n"
                f"Folder ID: {resolved_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"System folders (Inbox, Sent Items, Drafts, etc.) cannot be renamed.\n"
                f"Only custom folders can be renamed."
            )

        # Update folder
        folder_update = MailFolder()
        folder_update.display_name = new_name

        updated_folder = await mailbox.mail_folders.by_mail_folder_id(resolved_id).patch(folder_update)

        return (
            f"✅ Folder renamed successfully\n\n"
            f"Old name: {old_name}\n"
            f"New name: {updated_folder.display_name}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Folder ID: {resolved_id}"
        )

    except ValueError as e:
        return f"❌ Folder resolution error: {str(e)}"
    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "already exists" in error_msg.lower() or "duplicate" in error_msg.lower():
            return (
                f"❌ Folder name already exists\n\n"
                f"New name: {new_name}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Folder names must be unique within the same parent folder."
            )

        return (
            f"❌ Error renaming folder: {error_type}\n\n"
            f"Folder: {folder_id}\n"
            f"New name: {new_name}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def delete_folder(
    folder_id: str,
    mailbox_id: str
) -> str:
    """
    Delete folder and ALL contents permanently.

    ⚠️ WARNING: This permanently deletes the folder AND all emails/subfolders inside!
    Deleted items are NOT moved to Deleted Items folder - they are permanently removed.

    Args:
        folder_id: Folder to delete (ID, name, or path) - REQUIRED
        mailbox_id: Mailbox containing folder ("me" or email address) - REQUIRED

    Returns:
        Confirmation with count of deleted items and subfolders

    Examples:
        delete_folder(folder_id="OldProjects", mailbox_id="me")
        delete_folder(folder_id="Inbox/Archive2020", mailbox_id="me")
        delete_folder(folder_id="AAMkADU3...", mailbox_id="thomas@sixpillar.co.uk")

    Restrictions:
        - Cannot delete system folders (Inbox, Sent Items, Drafts, etc.)
        - Deletion is permanent and cannot be undone
        - ALL emails and subfolders are deleted recursively

    Safety:
        - Tool shows item count before deletion
        - System folders are protected
        - Consider archiving emails before deleting folders
    """
    try:
        graph_client = await get_graph_client()
        mailbox = get_mailbox_endpoint(graph_client, mailbox_id)

        # Resolve folder ID
        resolved_id, folder_name = await resolve_folder_id(folder_id, mailbox_id, graph_client)

        # Check if system folder
        if await is_system_folder(resolved_id, mailbox_id, graph_client):
            return (
                f"❌ Cannot delete system folder\n\n"
                f"Folder: {folder_name}\n"
                f"Folder ID: {resolved_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"System folders (Inbox, Sent Items, Drafts, etc.) cannot be deleted.\n"
                f"Only custom folders can be deleted."
            )

        # Get folder stats before deletion (for confirmation message)
        stats = await get_folder_stats(resolved_id, mailbox_id, graph_client)

        # Delete folder
        await mailbox.mail_folders.by_mail_folder_id(resolved_id).delete()

        # Build confirmation message
        items_msg = f"{stats['total_items']} email{'s' if stats['total_items'] != 1 else ''}"
        subfolders_msg = f"{stats['child_folders']} subfolder{'s' if stats['child_folders'] != 1 else ''}"

        return (
            f"✅ Folder deleted permanently\n\n"
            f"Folder: {folder_name}\n"
            f"Contents deleted:\n"
            f"  - {items_msg}\n"
            f"  - {subfolders_msg}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Folder ID: {resolved_id}\n\n"
            f"⚠️ Deletion is permanent and cannot be undone."
        )

    except ValueError as e:
        return f"❌ Folder resolution error: {str(e)}"
    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error deleting folder: {error_type}\n\n"
            f"Folder: {folder_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )

# =============================================================================
# CALENDAR TOOLS (Phase 9)
# =============================================================================


def format_event_time(dt_tz: DateTimeTimeZone, is_all_day: bool = False) -> str:
    """
    Format a DateTimeTimeZone object for display in MCP_MS_GRAPH_TIMEZONE.

    Args:
        dt_tz: DateTimeTimeZone object from Graph API
        is_all_day: Whether this is an all-day event

    Returns:
        Formatted time string in MCP_MS_GRAPH_TIMEZONE
    """
    if not dt_tz or not dt_tz.date_time:
        return "Unknown"

    try:
        # Parse the datetime string
        dt_str = dt_tz.date_time
        source_tz_str = dt_tz.time_zone or "UTC"

        # Parse datetime (Graph API returns ISO format without timezone offset)
        if "T" in dt_str:
            dt = datetime.fromisoformat(dt_str.replace("Z", ""))
        else:
            # Date only (all-day events)
            dt = datetime.fromisoformat(dt_str)

        if is_all_day:
            # All-day events: just show the date (no timezone conversion needed)
            return dt.strftime("%a %d %b %Y")
        else:
            # Convert from source timezone to display timezone
            source_tz = ZoneInfo(source_tz_str)
            display_tz = ZoneInfo(MCP_MS_GRAPH_TIMEZONE)
            dt_with_tz = dt.replace(tzinfo=source_tz)
            dt_local = dt_with_tz.astimezone(display_tz)
            return dt_local.strftime("%a %d %b %Y, %H:%M")

    except Exception:
        return dt_tz.date_time or "Unknown"


def format_event_time_range(start: DateTimeTimeZone, end: DateTimeTimeZone, is_all_day: bool = False) -> str:
    """
    Format start and end times as a range in MCP_MS_GRAPH_TIMEZONE.

    Args:
        start: Start DateTimeTimeZone
        end: End DateTimeTimeZone
        is_all_day: Whether this is an all-day event

    Returns:
        Formatted time range string in MCP_MS_GRAPH_TIMEZONE
    """
    if is_all_day:
        return f"{format_event_time(start, True)} (all day)"

    try:
        start_dt = datetime.fromisoformat(start.date_time.replace("Z", ""))
        end_dt = datetime.fromisoformat(end.date_time.replace("Z", ""))
        source_tz_str = start.time_zone or "UTC"

        # Convert from source timezone to display timezone
        source_tz = ZoneInfo(source_tz_str)
        display_tz = ZoneInfo(MCP_MS_GRAPH_TIMEZONE)
        start_local = start_dt.replace(tzinfo=source_tz).astimezone(display_tz)
        end_local = end_dt.replace(tzinfo=source_tz).astimezone(display_tz)

        # Same day: show date once with time range
        if start_local.date() == end_local.date():
            return f"{start_local.strftime('%a %d %b %Y')}, {start_local.strftime('%H:%M')}-{end_local.strftime('%H:%M')} {MCP_MS_GRAPH_TIMEZONE}"
        else:
            # Different days: show both dates
            return f"{start_local.strftime('%a %d %b %Y %H:%M')} - {end_local.strftime('%a %d %b %Y %H:%M')} {MCP_MS_GRAPH_TIMEZONE}"

    except Exception:
        return f"{format_event_time(start)} - {format_event_time(end)}"


def format_attendees_summary(attendees: list) -> str:
    """
    Format attendees list for display.

    Args:
        attendees: List of Attendee objects

    Returns:
        Formatted attendee summary
    """
    if not attendees:
        return "None"

    count = len(attendees)
    if count <= 3:
        names = []
        for a in attendees:
            if a.email_address:
                name = a.email_address.name or a.email_address.address or "Unknown"
                names.append(name)
        return ", ".join(names)
    else:
        # Show first 2 and count
        names = []
        for a in attendees[:2]:
            if a.email_address:
                name = a.email_address.name or a.email_address.address or "Unknown"
                names.append(name)
        return f"{', '.join(names)} +{count - 2} more"


@mcp.tool()
async def list_calendars(mailbox_id: str) -> str:
    """
    List all calendars available to the user.

    Shows default calendar, shared calendars, and group calendars with their
    permissions and colors.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        Formatted list of calendars with name, ID, color, and permissions

    Examples:
        list_calendars(mailbox_id="me")
        list_calendars(mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        graph_client = await get_graph_client()

        # Get all calendars
        calendars_response = await graph_client.users.by_user_id(mailbox_id).calendars.get()

        if not calendars_response or not calendars_response.value:
            return f"No calendars found for {mailbox_id}"

        calendars = calendars_response.value

        # Build output
        lines = [
            f"Calendars for {mailbox_id}:",
            ""
        ]

        for i, cal in enumerate(calendars, 1):
            name = cal.name or "Unnamed"
            cal_id = cal.id or "Unknown"
            color = cal.color.value if cal.color else "None"
            can_edit = cal.can_edit if cal.can_edit is not None else True
            is_default = cal.is_default_calendar if cal.is_default_calendar is not None else False

            # Format calendar entry
            default_marker = " (default)" if is_default else ""
            edit_status = "Can edit: Yes" if can_edit else "Can edit: No (read-only)"

            lines.append(f"{i}. {name}{default_marker}")
            lines.append(f"   ID: {cal_id[:50]}{'...' if len(cal_id) > 50 else ''}")
            lines.append(f"   Color: {color}")
            lines.append(f"   {edit_status}")
            lines.append("")

        lines.append(f"---")
        lines.append(f"Total: {len(calendars)} calendar{'s' if len(calendars) != 1 else ''}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error listing calendars: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def list_calendar_events(
    mailbox_id: str,
    start_date: str = "",
    end_date: str = "",
    calendar_id: str = "",
    count: int = 25
) -> str:
    """
    List calendar events in a date range.

    Uses calendarView endpoint which automatically expands recurring events
    into individual instances.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
        start_date: Start of range in YYYY-MM-DD format (default: today)
        end_date: End of range in YYYY-MM-DD format (default: 7 days from start)
        calendar_id: Specific calendar ID (default: primary calendar)
        count: Maximum events to return (default: 25, max: 100)

    Returns:
        Formatted list of events with time, subject, location, attendees

    Examples:
        list_calendar_events(mailbox_id="thomas@sixpillar.co.uk")
        list_calendar_events(mailbox_id="me", start_date="2025-12-01", end_date="2025-12-07")
    """
    try:
        graph_client = await get_graph_client()

        # Parse dates - default to today and +7 days
        if start_date:
            try:
                start_dt = datetime.fromisoformat(start_date)
            except ValueError:
                return f"❌ Invalid start_date format: {start_date}. Use YYYY-MM-DD."
        else:
            start_dt = datetime.now().replace(hour=0, minute=0, second=0, microsecond=0)

        if end_date:
            try:
                end_dt = datetime.fromisoformat(end_date)
            except ValueError:
                return f"❌ Invalid end_date format: {end_date}. Use YYYY-MM-DD."
        else:
            end_dt = start_dt + timedelta(days=7)

        # Ensure end is at end of day
        end_dt = end_dt.replace(hour=23, minute=59, second=59)

        # Format for Graph API (ISO 8601)
        start_iso = start_dt.strftime("%Y-%m-%dT%H:%M:%S")
        end_iso = end_dt.strftime("%Y-%m-%dT%H:%M:%S")

        # Limit count
        count = min(count, 100)

        # Build query params for calendarView (uses snake_case per SDK)
        query_params = CalendarViewRequestBuilder.CalendarViewRequestBuilderGetQueryParameters(
            start_date_time=start_iso,
            end_date_time=end_iso,
            top=count,
            orderby=["start/dateTime"],
            select=["id", "subject", "start", "end", "location", "attendees", "isAllDay", "organizer", "isCancelled"]
        )
        request_config = RequestConfiguration(query_parameters=query_params)

        # Build query - calendarView expands recurring events
        if calendar_id:
            # Specific calendar - use same query params structure
            specific_query_params = CalendarSpecificViewRequestBuilder.CalendarViewRequestBuilderGetQueryParameters(
                start_date_time=start_iso,
                end_date_time=end_iso,
                top=count,
                orderby=["start/dateTime"],
                select=["id", "subject", "start", "end", "location", "attendees", "isAllDay", "organizer", "isCancelled"]
            )
            specific_request_config = RequestConfiguration(query_parameters=specific_query_params)
            events_response = await graph_client.users.by_user_id(mailbox_id).calendars.by_calendar_id(calendar_id).calendar_view.get(
                request_configuration=specific_request_config
            )
        else:
            # Default calendar
            events_response = await graph_client.users.by_user_id(mailbox_id).calendar.calendar_view.get(
                request_configuration=request_config
            )

        if not events_response or not events_response.value:
            return (
                f"No events found for {mailbox_id}\n\n"
                f"Date range: {start_dt.strftime('%d %b %Y')} - {end_dt.strftime('%d %b %Y')}"
            )

        events = events_response.value

        # Group events by date
        events_by_date = {}
        for event in events:
            if event.is_cancelled:
                continue  # Skip cancelled events

            # Get event date
            if event.is_all_day:
                event_date = datetime.fromisoformat(event.start.date_time.split("T")[0]).date()
            else:
                event_date = datetime.fromisoformat(event.start.date_time.replace("Z", "")).date()

            date_key = event_date.strftime("%Y-%m-%d")
            if date_key not in events_by_date:
                events_by_date[date_key] = []
            events_by_date[date_key].append(event)

        # Build output
        lines = [
            f"Events for {mailbox_id} ({start_dt.strftime('%d %b')} - {end_dt.strftime('%d %b %Y')}):",
            ""
        ]

        total_events = 0
        for date_key in sorted(events_by_date.keys()):
            date_events = events_by_date[date_key]
            date_obj = datetime.fromisoformat(date_key)

            lines.append(date_obj.strftime("%a %d %b %Y"))
            lines.append("-" * len(date_obj.strftime("%a %d %b %Y")))

            for event in date_events:
                total_events += 1
                subject = event.subject or "(No subject)"

                # Time display (with timezone conversion)
                if event.is_all_day:
                    time_str = "(All day)"
                else:
                    # Convert times to display timezone
                    source_tz_str = event.start.time_zone or "UTC"
                    source_tz = ZoneInfo(source_tz_str)
                    display_tz = ZoneInfo(MCP_MS_GRAPH_TIMEZONE)
                    start_dt = datetime.fromisoformat(event.start.date_time.replace("Z", "")).replace(tzinfo=source_tz).astimezone(display_tz)
                    end_dt = datetime.fromisoformat(event.end.date_time.replace("Z", "")).replace(tzinfo=source_tz).astimezone(display_tz)
                    time_str = f"{start_dt.strftime('%H:%M')}-{end_dt.strftime('%H:%M')}"

                lines.append(f"{time_str:14} {subject}")

                # Location
                if event.location and event.location.display_name:
                    lines.append(f"{'':14} Location: {event.location.display_name}")

                # Attendees
                if event.attendees:
                    attendee_summary = format_attendees_summary(event.attendees)
                    lines.append(f"{'':14} Attendees: {attendee_summary}")

                # Event ID (full - needed for get_event_by_id)
                event_id = event.id or "Unknown"
                lines.append(f"{'':14} [ID: {event_id}]")
                lines.append("")

            lines.append("")

        lines.append("---")
        lines.append(f"Showing {total_events} event{'s' if total_events != 1 else ''}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error listing calendar events: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Date range: {start_date or 'today'} to {end_date or '+7 days'}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_event_by_id(
    event_id: str,
    mailbox_id: str
) -> str:
    """
    Get complete details of a calendar event.

    Returns full event information including body content, all attendees with
    their response status, recurrence pattern, and online meeting details.

    Args:
        event_id: The event ID (from list_calendar_events)
        mailbox_id: Email address of mailbox to access (required)

    Returns:
        Full event details including body, attendees with responses, recurrence

    Examples:
        get_event_by_id(event_id="AAMk...", mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        graph_client = await get_graph_client()

        # Get full event details
        event = await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).get()

        if not event:
            return f"❌ Event not found: {event_id}"

        # Build output
        lines = [
            "Event Details",
            "=============",
            ""
        ]

        # Subject and status
        subject = event.subject or "(No subject)"
        lines.append(f"Subject: {subject}")

        # Time
        is_all_day = event.is_all_day or False
        time_range = format_event_time_range(event.start, event.end, is_all_day)
        lines.append(f"When: {time_range}")

        # Location
        if event.location and event.location.display_name:
            lines.append(f"Location: {event.location.display_name}")

        # Online meeting
        if event.is_online_meeting and event.online_meeting:
            if hasattr(event.online_meeting, 'join_url') and event.online_meeting.join_url:
                lines.append(f"Online Meeting: {event.online_meeting.join_url}")

        # Organizer
        if event.organizer and event.organizer.email_address:
            org_name = event.organizer.email_address.name or ""
            org_email = event.organizer.email_address.address or ""
            lines.append(f"Organizer: {org_name} <{org_email}>")

        lines.append("")

        # Status
        if event.is_cancelled:
            lines.append("Status: CANCELLED")
        else:
            lines.append("Status: Confirmed")

        # Your response (if you're an attendee)
        if event.response_status:
            response = event.response_status.response.value if event.response_status.response else "none"
            lines.append(f"Your response: {response.capitalize()}")

        lines.append("")

        # Attendees
        if event.attendees:
            lines.append("Attendees:")

            # Group by response
            accepted = []
            tentative = []
            declined = []
            no_response = []

            for attendee in event.attendees:
                if not attendee.email_address:
                    continue

                name = attendee.email_address.name or ""
                email = attendee.email_address.address or ""
                attendee_type = attendee.type.value if attendee.type else "required"

                if attendee.status and attendee.status.response:
                    response = attendee.status.response.value
                else:
                    response = "none"

                entry = f"{name} <{email}>" if name else email
                if attendee_type == "optional":
                    entry += " (optional)"

                if response == "accepted":
                    accepted.append(entry)
                elif response in ["tentativelyAccepted", "tentative"]:
                    tentative.append(entry)
                elif response == "declined":
                    declined.append(entry)
                else:
                    no_response.append(entry)

            if accepted:
                for a in accepted:
                    lines.append(f"  ✓ {a} - Accepted")
            if tentative:
                for a in tentative:
                    lines.append(f"  ? {a} - Tentative")
            if declined:
                for a in declined:
                    lines.append(f"  ✗ {a} - Declined")
            if no_response:
                for a in no_response:
                    lines.append(f"  - {a} - No response")

            lines.append("")

        # Recurrence
        if event.recurrence:
            pattern = event.recurrence.pattern
            if pattern:
                pattern_type = pattern.type.value if pattern.type else "unknown"
                interval = pattern.interval or 1
                lines.append(f"Recurrence: {pattern_type} (every {interval})")
            else:
                lines.append("Recurrence: Yes (pattern details unavailable)")
        else:
            lines.append("Recurrence: None (single event)")

        lines.append("")

        # Body
        if event.body and event.body.content:
            content = event.body.content

            # Strip HTML if present
            if event.body.content_type and event.body.content_type.value == "html":
                # Simple HTML stripping
                import re
                content = re.sub(r'<[^>]+>', '', content)
                content = content.replace('&nbsp;', ' ').replace('&amp;', '&')
                content = re.sub(r'\n\s*\n', '\n\n', content).strip()

            if content:
                lines.append("Body:")
                # Limit body length
                if len(content) > 1000:
                    content = content[:1000] + "..."
                lines.append(content)
                lines.append("")

        # Metadata
        lines.append("---")
        lines.append(f"Event ID: {event.id}")
        if event.created_date_time:
            # SDK may return datetime object or string
            if isinstance(event.created_date_time, datetime):
                created = event.created_date_time.strftime("%Y-%m-%d %H:%M")
            else:
                created = datetime.fromisoformat(event.created_date_time.replace("Z", "")).strftime("%Y-%m-%d %H:%M")
            lines.append(f"Created: {created}")
        if event.last_modified_date_time:
            # SDK may return datetime object or string
            if isinstance(event.last_modified_date_time, datetime):
                modified = event.last_modified_date_time.strftime("%Y-%m-%d %H:%M")
            else:
                modified = datetime.fromisoformat(event.last_modified_date_time.replace("Z", "")).strftime("%Y-%m-%d %H:%M")
            lines.append(f"Last Modified: {modified}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error getting event: {error_type}\n\n"
            f"Event ID: {event_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_calendar_view(
    mailbox_id: str,
    start_datetime: str,
    end_datetime: str,
    calendar_id: str = ""
) -> str:
    """
    Get calendar view showing free/busy status for a time range.

    More detailed than list_calendar_events - shows exact busy periods
    and free time blocks. Useful for checking availability before scheduling.

    Args:
        mailbox_id: Email address of mailbox to access (required)
        start_datetime: Start of range (YYYY-MM-DDTHH:MM or YYYY-MM-DD)
        end_datetime: End of range (YYYY-MM-DDTHH:MM or YYYY-MM-DD)
        calendar_id: Specific calendar ID (default: primary calendar)

    Returns:
        Time-blocked view showing free and busy periods

    Examples:
        get_calendar_view(mailbox_id="thomas@sixpillar.co.uk",
                         start_datetime="2025-11-28T09:00",
                         end_datetime="2025-11-28T18:00")
    """
    try:
        graph_client = await get_graph_client()

        # Parse datetimes
        try:
            if "T" in start_datetime:
                start_dt = datetime.fromisoformat(start_datetime)
            else:
                start_dt = datetime.fromisoformat(f"{start_datetime}T00:00:00")
        except ValueError:
            return f"❌ Invalid start_datetime format: {start_datetime}. Use YYYY-MM-DDTHH:MM or YYYY-MM-DD."

        try:
            if "T" in end_datetime:
                end_dt = datetime.fromisoformat(end_datetime)
            else:
                end_dt = datetime.fromisoformat(f"{end_datetime}T23:59:59")
        except ValueError:
            return f"❌ Invalid end_datetime format: {end_datetime}. Use YYYY-MM-DDTHH:MM or YYYY-MM-DD."

        # Format for Graph API
        start_iso = start_dt.strftime("%Y-%m-%dT%H:%M:%S")
        end_iso = end_dt.strftime("%Y-%m-%dT%H:%M:%S")

        # Build query params for calendarView (uses snake_case per SDK)
        query_params = CalendarViewRequestBuilder.CalendarViewRequestBuilderGetQueryParameters(
            start_date_time=start_iso,
            end_date_time=end_iso,
            orderby=["start/dateTime"],
            select=["id", "subject", "start", "end", "isAllDay", "showAs", "isCancelled"]
        )
        request_config = RequestConfiguration(query_parameters=query_params)

        # Get events in range
        if calendar_id:
            specific_query_params = CalendarSpecificViewRequestBuilder.CalendarViewRequestBuilderGetQueryParameters(
                start_date_time=start_iso,
                end_date_time=end_iso,
                orderby=["start/dateTime"],
                select=["id", "subject", "start", "end", "isAllDay", "showAs", "isCancelled"]
            )
            specific_request_config = RequestConfiguration(query_parameters=specific_query_params)
            events_response = await graph_client.users.by_user_id(mailbox_id).calendars.by_calendar_id(calendar_id).calendar_view.get(
                request_configuration=specific_request_config
            )
        else:
            events_response = await graph_client.users.by_user_id(mailbox_id).calendar.calendar_view.get(
                request_configuration=request_config
            )

        events = events_response.value if events_response and events_response.value else []

        # Filter out cancelled events
        events = [e for e in events if not e.is_cancelled]

        # Build time blocks
        lines = [
            f"Calendar View: {start_dt.strftime('%a %d %b %Y')}, {start_dt.strftime('%H:%M')}-{end_dt.strftime('%H:%M')}",
            ""
        ]

        if not events:
            # All free
            duration = (end_dt - start_dt).total_seconds() / 3600
            lines.append(f"{start_dt.strftime('%H:%M')}-{end_dt.strftime('%H:%M')}  [FREE]   {duration:.1f} hours available")
            lines.append("")
            lines.append(f"Summary: {duration:.1f} hours free, 0 hours busy")
            return "\n".join(lines)

        # Calculate free/busy blocks
        total_busy_minutes = 0
        current_time = start_dt

        for event in events:
            # Parse event times
            if event.is_all_day:
                event_start = start_dt
                event_end = end_dt
            else:
                event_start = datetime.fromisoformat(event.start.date_time.replace("Z", ""))
                event_end = datetime.fromisoformat(event.end.date_time.replace("Z", ""))

            # Clamp to view range
            event_start = max(event_start, start_dt)
            event_end = min(event_end, end_dt)

            # Free time before this event
            if event_start > current_time:
                free_minutes = (event_start - current_time).total_seconds() / 60
                if free_minutes >= 15:  # Only show if >= 15 min
                    free_hours = free_minutes / 60
                    lines.append(f"{current_time.strftime('%H:%M')}-{event_start.strftime('%H:%M')}  [FREE]   {free_hours:.1f} hours available")

            # Busy time for this event
            subject = event.subject or "(No subject)"
            show_as = event.show_as.value if event.show_as else "busy"
            busy_minutes = (event_end - event_start).total_seconds() / 60
            total_busy_minutes += busy_minutes

            status = "[BUSY]" if show_as in ["busy", "oof", "workingElsewhere"] else f"[{show_as.upper()}]"
            lines.append(f"{event_start.strftime('%H:%M')}-{event_end.strftime('%H:%M')}  {status:10} {subject}")

            current_time = max(current_time, event_end)

        # Free time after last event
        if current_time < end_dt:
            free_minutes = (end_dt - current_time).total_seconds() / 60
            if free_minutes >= 15:
                free_hours = free_minutes / 60
                lines.append(f"{current_time.strftime('%H:%M')}-{end_dt.strftime('%H:%M')}  [FREE]   {free_hours:.1f} hours available")

        # Summary
        total_minutes = (end_dt - start_dt).total_seconds() / 60
        free_minutes = total_minutes - total_busy_minutes

        lines.append("")
        lines.append(f"Summary: {free_minutes/60:.1f} hours free, {total_busy_minutes/60:.1f} hours busy")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error getting calendar view: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Range: {start_datetime} to {end_datetime}\n"
            f"Error: {error_msg}"
        )


# =============================================================================
# CONTACT TOOLS (Phase 10)
# =============================================================================


def format_phone_numbers(contact: Contact) -> list:
    """
    Extract and format phone numbers from a contact.

    Args:
        contact: Contact object from Graph API

    Returns:
        List of tuples (number, type)
    """
    phones = []

    if contact.mobile_phone:
        phones.append((contact.mobile_phone, "mobile"))
    if contact.business_phones:
        for phone in contact.business_phones:
            phones.append((phone, "business"))
    if contact.home_phones:
        for phone in contact.home_phones:
            phones.append((phone, "home"))

    return phones


def format_email_addresses(contact: Contact) -> list:
    """
    Extract email addresses from a contact.

    Args:
        contact: Contact object from Graph API

    Returns:
        List of tuples (email, type)
    """
    emails = []

    if contact.email_addresses:
        for i, email in enumerate(contact.email_addresses):
            if email.address:
                label = email.name if email.name else f"email{i+1}"
                emails.append((email.address, label))

    return emails


def format_address(address: PhysicalAddress, label: str) -> str:
    """
    Format a physical address for display.

    Args:
        address: PhysicalAddress object
        label: Address type label

    Returns:
        Formatted address string
    """
    if not address:
        return ""

    parts = []
    if address.street:
        parts.append(address.street)
    if address.city:
        city_line = address.city
        if address.state:
            city_line += f", {address.state}"
        if address.postal_code:
            city_line += f" {address.postal_code}"
        parts.append(city_line)
    if address.country_or_region:
        parts.append(address.country_or_region)

    if parts:
        return f"{label}:\n  " + "\n  ".join(parts)
    return ""


@mcp.tool()
async def list_contact_folders(mailbox_id: str) -> str:
    """
    List all contact folders in the mailbox.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        Formatted list of contact folders with name, ID, and contact count

    Examples:
        list_contact_folders(mailbox_id="me")
        list_contact_folders(mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        graph_client = await get_graph_client()

        # Get all contact folders
        folders_response = await graph_client.users.by_user_id(mailbox_id).contact_folders.get()

        if not folders_response or not folders_response.value:
            return f"No contact folders found for {mailbox_id}"

        folders = folders_response.value
        total_contacts = 0

        # Build output
        lines = [
            f"Contact Folders for {mailbox_id}:",
            ""
        ]

        for i, folder in enumerate(folders, 1):
            name = folder.display_name or "Unnamed"
            folder_id = folder.id or "Unknown"

            # Get contact count for this folder
            try:
                contacts_response = await graph_client.users.by_user_id(mailbox_id).contact_folders.by_contact_folder_id(folder.id).contacts.get()
                contact_count = len(contacts_response.value) if contacts_response and contacts_response.value else 0
            except Exception:
                contact_count = 0

            total_contacts += contact_count

            lines.append(f"{i}. {name}")
            lines.append(f"   ID: {folder_id[:50]}{'...' if len(folder_id) > 50 else ''}")
            lines.append(f"   Contact count: {contact_count}")
            lines.append("")

        lines.append(f"---")
        lines.append(f"Total: {len(folders)} folder{'s' if len(folders) != 1 else ''}, {total_contacts} contacts")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error listing contact folders: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def list_contacts(
    mailbox_id: str,
    folder_id: str = "",
    count: int = 25,
    skip: int = 0
) -> str:
    """
    List contacts in the mailbox or specific folder.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        folder_id: Specific contact folder ID (default: main Contacts folder)
        count: Maximum contacts to return (default: 25, max: 100)
        skip: Number of contacts to skip for pagination (default: 0)

    Returns:
        Formatted list of contacts with name, email, company, phone

    Examples:
        list_contacts(mailbox_id="thomas@sixpillar.co.uk")
        list_contacts(mailbox_id="me", count=10)
        list_contacts(mailbox_id="me", skip=25, count=25)
    """
    try:
        graph_client = await get_graph_client()

        # Limit count
        count = min(count, 100)

        # Build query parameters
        query_params = ContactsRequestBuilder.ContactsRequestBuilderGetQueryParameters(
            top=count,
            skip=skip,
            orderby=["displayName"],
            select=["id", "displayName", "emailAddresses", "companyName", "jobTitle", "mobilePhone", "businessPhones"]
        )
        request_config = RequestConfiguration(query_parameters=query_params)

        # Get contacts from specific folder or main contacts
        if folder_id:
            contacts_response = await graph_client.users.by_user_id(mailbox_id).contact_folders.by_contact_folder_id(folder_id).contacts.get(
                request_configuration=request_config
            )
        else:
            contacts_response = await graph_client.users.by_user_id(mailbox_id).contacts.get(
                request_configuration=request_config
            )

        if not contacts_response or not contacts_response.value:
            return (
                f"No contacts found for {mailbox_id}\n\n"
                f"Folder: {folder_id if folder_id else 'All Contacts'}\n"
                f"Skip: {skip}"
            )

        contacts = contacts_response.value

        # Build output
        lines = [
            f"Contacts for {mailbox_id}:",
            ""
        ]

        for i, contact in enumerate(contacts, 1):
            display_name = contact.display_name or "(No name)"
            contact_id = contact.id or "Unknown"

            lines.append(f"{skip + i}. {display_name}")

            # Primary email
            if contact.email_addresses and len(contact.email_addresses) > 0:
                primary_email = contact.email_addresses[0].address if contact.email_addresses[0].address else "No email"
                lines.append(f"   Email: {primary_email}")

            # Company
            if contact.company_name:
                lines.append(f"   Company: {contact.company_name}")

            # Job title
            if contact.job_title:
                lines.append(f"   Title: {contact.job_title}")

            # Phone (prefer mobile, then business)
            phone = None
            if contact.mobile_phone:
                phone = contact.mobile_phone
            elif contact.business_phones and len(contact.business_phones) > 0:
                phone = contact.business_phones[0]

            if phone:
                lines.append(f"   Phone: {phone}")

            lines.append(f"   [ID: {contact_id}]")
            lines.append("")

        lines.append("---")
        lines.append(f"Showing {len(contacts)} contact{'s' if len(contacts) != 1 else ''} (skip: {skip})")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error listing contacts: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Folder: {folder_id if folder_id else 'All Contacts'}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_contact_by_id(
    contact_id: str,
    mailbox_id: str
) -> str:
    """
    Get complete details of a contact.

    Args:
        contact_id: The contact ID (from list_contacts or search)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        Full contact details including all addresses, phones, and notes

    Examples:
        get_contact_by_id(contact_id="AAMk...", mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        graph_client = await get_graph_client()

        # Get full contact details
        contact = await graph_client.users.by_user_id(mailbox_id).contacts.by_contact_id(contact_id).get()

        if not contact:
            return f"❌ Contact not found: {contact_id}"

        # Build output
        lines = [
            "Contact Details",
            "===============",
            ""
        ]

        # Name
        display_name = contact.display_name or "(No name)"
        lines.append(f"Name: {display_name}")

        # Name parts (if available)
        name_parts = []
        if contact.given_name:
            name_parts.append(f"First: {contact.given_name}")
        if contact.surname:
            name_parts.append(f"Last: {contact.surname}")
        if contact.middle_name:
            name_parts.append(f"Middle: {contact.middle_name}")
        if name_parts and (contact.given_name or contact.surname):
            lines.append(f"  ({', '.join(name_parts)})")

        # Company info
        if contact.company_name:
            lines.append(f"Company: {contact.company_name}")
        if contact.job_title:
            lines.append(f"Job Title: {contact.job_title}")
        if contact.department:
            lines.append(f"Department: {contact.department}")
        if contact.manager:
            lines.append(f"Manager: {contact.manager}")

        lines.append("")

        # Email addresses
        emails = format_email_addresses(contact)
        if emails:
            lines.append("Email Addresses:")
            for email, label in emails:
                lines.append(f"  - {email} ({label})")
            lines.append("")

        # Phone numbers
        phones = format_phone_numbers(contact)
        if phones:
            lines.append("Phone Numbers:")
            for phone, phone_type in phones:
                lines.append(f"  - {phone} ({phone_type})")
            lines.append("")

        # Addresses
        addresses = []
        if contact.business_address:
            addr = format_address(contact.business_address, "Business Address")
            if addr:
                addresses.append(addr)
        if contact.home_address:
            addr = format_address(contact.home_address, "Home Address")
            if addr:
                addresses.append(addr)
        if contact.other_address:
            addr = format_address(contact.other_address, "Other Address")
            if addr:
                addresses.append(addr)

        if addresses:
            for addr in addresses:
                lines.append(addr)
            lines.append("")

        # Birthday
        if contact.birthday:
            try:
                bday = contact.birthday.strftime("%d %B %Y")
                lines.append(f"Birthday: {bday}")
            except Exception:
                lines.append(f"Birthday: {contact.birthday}")

        # Personal note
        if contact.personal_notes:
            lines.append("")
            lines.append("Notes:")
            lines.append(contact.personal_notes)

        # Categories
        if contact.categories:
            lines.append("")
            lines.append(f"Categories: {', '.join(contact.categories)}")

        # Metadata
        lines.append("")
        lines.append("---")
        lines.append(f"Contact ID: {contact.id}")

        if contact.created_date_time:
            try:
                created = contact.created_date_time.strftime("%Y-%m-%d %H:%M")
                lines.append(f"Created: {created}")
            except Exception:
                lines.append(f"Created: {contact.created_date_time}")

        if contact.last_modified_date_time:
            try:
                modified = contact.last_modified_date_time.strftime("%Y-%m-%d %H:%M")
                lines.append(f"Last Modified: {modified}")
            except Exception:
                lines.append(f"Last Modified: {contact.last_modified_date_time}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Contact not found\n\n"
                f"Contact ID: {contact_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure the contact ID is correct."
            )

        return (
            f"❌ Error getting contact: {error_type}\n\n"
            f"Contact ID: {contact_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_contact_photo(
    contact_id: str,
    mailbox_id: str
) -> str:
    """
    Check if a contact has a photo and get metadata.

    Args:
        contact_id: The contact ID
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        Photo availability status and metadata

    Examples:
        get_contact_photo(contact_id="AAMk...", mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        graph_client = await get_graph_client()

        # First verify contact exists
        contact = await graph_client.users.by_user_id(mailbox_id).contacts.by_contact_id(contact_id).get()

        if not contact:
            return f"❌ Contact not found: {contact_id}"

        display_name = contact.display_name or "(No name)"

        # Try to get photo metadata
        try:
            photo = await graph_client.users.by_user_id(mailbox_id).contacts.by_contact_id(contact_id).photo.get()

            if photo:
                lines = [
                    f"Contact Photo: {display_name}",
                    "",
                    "✅ Photo available",
                    ""
                ]

                if photo.height and photo.width:
                    lines.append(f"Dimensions: {photo.width}x{photo.height}")
                if photo.id:
                    lines.append(f"Photo ID: {photo.id}")

                lines.append("")
                lines.append(f"Contact ID: {contact_id}")

                return "\n".join(lines)

        except Exception as photo_error:
            # No photo or error getting photo
            error_str = str(photo_error)
            if "ImageNotFound" in error_str or "ResourceNotFound" in error_str or "404" in error_str:
                return (
                    f"Contact Photo: {display_name}\n\n"
                    f"❌ No photo available\n\n"
                    f"Contact ID: {contact_id}"
                )
            else:
                # Some other error
                raise photo_error

        return (
            f"Contact Photo: {display_name}\n\n"
            f"❌ No photo available\n\n"
            f"Contact ID: {contact_id}"
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Contact not found\n\n"
                f"Contact ID: {contact_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error getting contact photo: {error_type}\n\n"
            f"Contact ID: {contact_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


# =============================================================================
# CONTACT SEARCH TOOLS (Phase 10 M2)
# =============================================================================


@mcp.tool()
async def search_contacts(
    mailbox_id: str,
    query: str,
    count: int = 20
) -> str:
    """
    Search contacts by name, email, or company.

    Searches across display name, email addresses, company name,
    and job title fields.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        query: Search keywords (required)
        count: Maximum results (default: 20, max: 50)

    Returns:
        Matching contacts with relevance-based ordering

    Examples:
        search_contacts(mailbox_id="thomas@sixpillar.co.uk", query="acme")
        search_contacts(mailbox_id="me", query="john smith")
    """
    try:
        if not query:
            return "❌ Error: query is required"

        graph_client = await get_graph_client()

        # Limit count
        count = min(count, 50)

        # Escape query for OData filter
        escaped_query = escape_odata_string(query)

        # Build filter to search across multiple fields
        # Using contains() for partial matching
        filter_query = (
            f"contains(displayName, '{escaped_query}') or "
            f"contains(companyName, '{escaped_query}') or "
            f"contains(jobTitle, '{escaped_query}')"
        )

        # Build query parameters
        query_params = ContactsRequestBuilder.ContactsRequestBuilderGetQueryParameters(
            filter=filter_query,
            top=count,
            orderby=["displayName"],
            select=["id", "displayName", "emailAddresses", "companyName", "jobTitle", "mobilePhone", "businessPhones"]
        )
        request_config = RequestConfiguration(query_parameters=query_params)

        # Execute search
        contacts_response = await graph_client.users.by_user_id(mailbox_id).contacts.get(
            request_configuration=request_config
        )

        if not contacts_response or not contacts_response.value:
            return (
                f"No contacts found matching '{query}'\n\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Try a different search term or check spelling."
            )

        contacts = contacts_response.value

        # Build output
        lines = [
            f"Search results for \"{query}\" in {mailbox_id}:",
            ""
        ]

        for i, contact in enumerate(contacts, 1):
            display_name = contact.display_name or "(No name)"
            contact_id = contact.id or "Unknown"

            lines.append(f"{i}. {display_name}")

            # Primary email
            if contact.email_addresses and len(contact.email_addresses) > 0:
                primary_email = contact.email_addresses[0].address if contact.email_addresses[0].address else "No email"
                lines.append(f"   Email: {primary_email}")

            # Company
            if contact.company_name:
                lines.append(f"   Company: {contact.company_name}")

            # Job title
            if contact.job_title:
                lines.append(f"   Title: {contact.job_title}")

            lines.append(f"   [ID: {contact_id}]")
            lines.append("")

        lines.append("---")
        lines.append(f"Found {len(contacts)} contact{'s' if len(contacts) != 1 else ''} matching \"{query}\"")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error searching contacts: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Query: {query}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def lookup_contact_by_email(
    mailbox_id: str,
    email: str
) -> str:
    """
    Look up a contact by exact email address.

    Useful for checking if a sender/recipient exists in contacts.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        email: Email address to look up (exact match required)

    Returns:
        Contact details if found, or "not found" message

    Examples:
        lookup_contact_by_email(mailbox_id="thomas@sixpillar.co.uk", email="john@example.com")
    """
    try:
        if not email:
            return "❌ Error: email is required"

        graph_client = await get_graph_client()

        # Escape email for OData filter
        escaped_email = escape_odata_string(email.lower())

        # Build filter to search email addresses
        # Graph API emailAddresses is a collection, need to search within it
        filter_query = f"emailAddresses/any(e:e/address eq '{escaped_email}')"

        # Build query parameters
        query_params = ContactsRequestBuilder.ContactsRequestBuilderGetQueryParameters(
            filter=filter_query,
            top=5,  # Should only be 1, but allow a few in case
            select=["id", "displayName", "emailAddresses", "companyName", "jobTitle", "mobilePhone", "businessPhones"]
        )
        request_config = RequestConfiguration(query_parameters=query_params)

        # Execute search
        contacts_response = await graph_client.users.by_user_id(mailbox_id).contacts.get(
            request_configuration=request_config
        )

        if not contacts_response or not contacts_response.value or len(contacts_response.value) == 0:
            return (
                f"❌ No contact found for: {email}\n\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"This email address is not in your contacts."
            )

        # Found contact(s)
        contacts = contacts_response.value

        if len(contacts) == 1:
            contact = contacts[0]
            display_name = contact.display_name or "(No name)"
            contact_id = contact.id or "Unknown"

            lines = [
                f"✅ Contact found for: {email}",
                "",
                f"Name: {display_name}"
            ]

            # All emails
            if contact.email_addresses:
                emails = [e.address for e in contact.email_addresses if e.address]
                if len(emails) > 1:
                    lines.append(f"Emails: {', '.join(emails)}")

            # Company
            if contact.company_name:
                lines.append(f"Company: {contact.company_name}")

            # Job title
            if contact.job_title:
                lines.append(f"Title: {contact.job_title}")

            # Phone
            phone = None
            if contact.mobile_phone:
                phone = contact.mobile_phone
            elif contact.business_phones and len(contact.business_phones) > 0:
                phone = contact.business_phones[0]
            if phone:
                lines.append(f"Phone: {phone}")

            lines.append("")
            lines.append(f"Contact ID: {contact_id}")

            return "\n".join(lines)
        else:
            # Multiple contacts with same email (unusual but possible)
            lines = [
                f"⚠️ Multiple contacts found for: {email}",
                ""
            ]

            for i, contact in enumerate(contacts, 1):
                display_name = contact.display_name or "(No name)"
                lines.append(f"{i}. {display_name}")
                if contact.company_name:
                    lines.append(f"   Company: {contact.company_name}")
                lines.append(f"   [ID: {contact.id}]")
                lines.append("")

            lines.append(f"---")
            lines.append(f"Found {len(contacts)} contacts with this email address")

            return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error looking up contact: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Email: {email}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_contacts_by_company(
    mailbox_id: str,
    company_name: str,
    count: int = 50
) -> str:
    """
    Get all contacts from a specific company.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        company_name: Company name to filter by (required)
        count: Maximum results (default: 50, max: 100)

    Returns:
        All contacts matching the company name

    Examples:
        get_contacts_by_company(mailbox_id="thomas@sixpillar.co.uk", company_name="Acme Corporation")
        get_contacts_by_company(mailbox_id="me", company_name="Six Pillar")
    """
    try:
        if not company_name:
            return "❌ Error: company_name is required"

        graph_client = await get_graph_client()

        # Limit count
        count = min(count, 100)

        # Fetch all contacts - Graph API doesn't support contains() for contacts
        # We need to filter client-side
        query_params = ContactsRequestBuilder.ContactsRequestBuilderGetQueryParameters(
            top=999,  # Fetch as many as possible
            orderby=["displayName"],
            select=["id", "displayName", "emailAddresses", "companyName", "jobTitle", "mobilePhone", "businessPhones", "department"]
        )
        request_config = RequestConfiguration(query_parameters=query_params)

        # Execute query
        contacts_response = await graph_client.users.by_user_id(mailbox_id).contacts.get(
            request_configuration=request_config
        )

        if not contacts_response or not contacts_response.value:
            return (
                f"No contacts found in mailbox\n\n"
                f"Mailbox: {mailbox_id}"
            )

        # Filter client-side by company name (case-insensitive partial match)
        company_lower = company_name.lower()
        contacts = [
            c for c in contacts_response.value
            if c.company_name and company_lower in c.company_name.lower()
        ][:count]  # Limit to requested count

        if not contacts:
            return (
                f"No contacts found for company: {company_name}\n\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Try a different company name or check spelling."
            )

        # Build output
        lines = [
            f"Contacts at \"{company_name}\":",
            f"Mailbox: {mailbox_id}",
            ""
        ]

        # Group by department if available
        by_dept = {}
        no_dept = []

        for contact in contacts:
            dept = contact.department if contact.department else None
            if dept:
                if dept not in by_dept:
                    by_dept[dept] = []
                by_dept[dept].append(contact)
            else:
                no_dept.append(contact)

        # Output by department
        def format_contact(contact, num):
            display_name = contact.display_name or "(No name)"
            contact_id = contact.id or "Unknown"

            contact_lines = [f"{num}. {display_name}"]

            # Primary email
            if contact.email_addresses and len(contact.email_addresses) > 0:
                primary_email = contact.email_addresses[0].address if contact.email_addresses[0].address else "No email"
                contact_lines.append(f"   Email: {primary_email}")

            # Job title
            if contact.job_title:
                contact_lines.append(f"   Title: {contact.job_title}")

            # Phone
            phone = None
            if contact.mobile_phone:
                phone = contact.mobile_phone
            elif contact.business_phones and len(contact.business_phones) > 0:
                phone = contact.business_phones[0]
            if phone:
                contact_lines.append(f"   Phone: {phone}")

            contact_lines.append(f"   [ID: {contact_id}]")
            contact_lines.append("")

            return contact_lines

        num = 1

        # Output contacts by department
        for dept in sorted(by_dept.keys()):
            lines.append(f"--- {dept} ---")
            for contact in by_dept[dept]:
                lines.extend(format_contact(contact, num))
                num += 1

        # Output contacts without department
        if no_dept:
            if by_dept:
                lines.append("--- Other ---")
            for contact in no_dept:
                lines.extend(format_contact(contact, num))
                num += 1

        lines.append("---")
        lines.append(f"Found {len(contacts)} contact{'s' if len(contacts) != 1 else ''} at \"{company_name}\"")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error getting company contacts: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Company: {company_name}\n"
            f"Error: {error_msg}"
        )


# =============================================================================
# CONTACT MANAGEMENT TOOLS (Phase 10 M3)
# =============================================================================


@mcp.tool()
async def create_contact(
    mailbox_id: str,
    display_name: str,
    email: str = "",
    company: str = "",
    job_title: str = "",
    mobile_phone: str = "",
    business_phone: str = "",
    notes: str = "",
    folder_id: str = ""
) -> str:
    """
    Create a new contact.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        display_name: Contact's full name (required)
        email: Primary email address (optional)
        company: Company/organization name (optional)
        job_title: Job title (optional)
        mobile_phone: Mobile phone number (optional)
        business_phone: Business phone number (optional)
        notes: Free-text notes about the contact (optional)
        folder_id: Contact folder to create in (default: main Contacts folder)

    Returns:
        Created contact details with ID

    Examples:
        create_contact(mailbox_id="thomas@sixpillar.co.uk", display_name="John Smith", email="john@example.com")
        create_contact(mailbox_id="me", display_name="Jane Doe", company="Acme Corp", job_title="CEO")
    """
    try:
        if not display_name:
            return "❌ Error: display_name is required"

        graph_client = await get_graph_client()

        # Create new contact object
        new_contact = Contact()
        new_contact.display_name = display_name

        # Parse name into given/surname if possible
        name_parts = display_name.strip().split()
        if len(name_parts) >= 2:
            new_contact.given_name = name_parts[0]
            new_contact.surname = " ".join(name_parts[1:])
        elif len(name_parts) == 1:
            new_contact.given_name = name_parts[0]

        # Set email if provided
        if email:
            # Use EmailAddress (already imported at top of file)
            email_obj = EmailAddress()
            email_obj.address = email
            email_obj.name = display_name
            new_contact.email_addresses = [email_obj]

        # Set company info
        if company:
            new_contact.company_name = company
        if job_title:
            new_contact.job_title = job_title

        # Set phone numbers
        if mobile_phone:
            new_contact.mobile_phone = mobile_phone
        if business_phone:
            new_contact.business_phones = [business_phone]

        # Set notes
        if notes:
            new_contact.personal_notes = notes

        # Create contact
        if folder_id:
            created = await graph_client.users.by_user_id(mailbox_id).contact_folders.by_contact_folder_id(folder_id).contacts.post(body=new_contact)
        else:
            created = await graph_client.users.by_user_id(mailbox_id).contacts.post(body=new_contact)

        if not created:
            return "❌ Error: Contact creation failed - no response"

        # Build success response
        lines = [
            "✅ Contact created successfully!",
            "",
            f"Name: {created.display_name or display_name}"
        ]

        if email:
            lines.append(f"Email: {email}")
        if company:
            lines.append(f"Company: {company}")
        if job_title:
            lines.append(f"Title: {job_title}")
        if mobile_phone:
            lines.append(f"Mobile: {mobile_phone}")
        if business_phone:
            lines.append(f"Business: {business_phone}")

        lines.append("")
        lines.append(f"Contact ID: {created.id}")
        lines.append(f"Mailbox: {mailbox_id}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error creating contact: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Name: {display_name}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def update_contact(
    contact_id: str,
    mailbox_id: str,
    display_name: str = "",
    email: str = "",
    company: str = "",
    job_title: str = "",
    mobile_phone: str = "",
    business_phone: str = "",
    notes: str = ""
) -> str:
    """
    Update an existing contact.

    Only provided parameters are updated - others remain unchanged.

    Args:
        contact_id: Contact to update (required)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        display_name: New display name (optional)
        email: New primary email (optional - replaces existing)
        company: New company name (optional)
        job_title: New job title (optional)
        mobile_phone: New mobile phone (optional)
        business_phone: New business phone (optional)
        notes: New notes (optional - replaces existing)

    Returns:
        Updated contact details

    Examples:
        update_contact(contact_id="AAMk...", mailbox_id="thomas@sixpillar.co.uk", mobile_phone="+44 7700 123456")
        update_contact(contact_id="AAMk...", mailbox_id="me", job_title="Senior Director")
    """
    try:
        if not contact_id:
            return "❌ Error: contact_id is required"

        graph_client = await get_graph_client()

        # Build update object with only provided fields
        update_contact = Contact()
        updated_fields = []

        if display_name:
            update_contact.display_name = display_name
            # Parse name into given/surname
            name_parts = display_name.strip().split()
            if len(name_parts) >= 2:
                update_contact.given_name = name_parts[0]
                update_contact.surname = " ".join(name_parts[1:])
            elif len(name_parts) == 1:
                update_contact.given_name = name_parts[0]
            updated_fields.append(f"Name: {display_name}")

        if email:
            # Use EmailAddress (already imported at top of file)
            email_obj = EmailAddress()
            email_obj.address = email
            update_contact.email_addresses = [email_obj]
            updated_fields.append(f"Email: {email}")

        if company:
            update_contact.company_name = company
            updated_fields.append(f"Company: {company}")

        if job_title:
            update_contact.job_title = job_title
            updated_fields.append(f"Title: {job_title}")

        if mobile_phone:
            update_contact.mobile_phone = mobile_phone
            updated_fields.append(f"Mobile: {mobile_phone}")

        if business_phone:
            update_contact.business_phones = [business_phone]
            updated_fields.append(f"Business: {business_phone}")

        if notes:
            update_contact.personal_notes = notes
            updated_fields.append("Notes: (updated)")

        if not updated_fields:
            return (
                "❌ No fields to update\n\n"
                "Provide at least one field to update (display_name, email, company, etc.)"
            )

        # Update contact
        updated = await graph_client.users.by_user_id(mailbox_id).contacts.by_contact_id(contact_id).patch(body=update_contact)

        if not updated:
            return "❌ Error: Contact update failed - no response"

        # Build success response
        lines = [
            "✅ Contact updated successfully!",
            "",
            "Updated fields:"
        ]

        for field in updated_fields:
            lines.append(f"  - {field}")

        lines.append("")
        lines.append(f"Contact ID: {contact_id}")
        lines.append(f"Mailbox: {mailbox_id}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Contact not found\n\n"
                f"Contact ID: {contact_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure the contact ID is correct."
            )

        return (
            f"❌ Error updating contact: {error_type}\n\n"
            f"Contact ID: {contact_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def delete_contact(
    contact_id: str,
    mailbox_id: str
) -> str:
    """
    Delete a contact.

    This action is permanent and cannot be undone.

    Args:
        contact_id: Contact to delete (required)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        Confirmation of deletion

    Examples:
        delete_contact(contact_id="AAMk...", mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        if not contact_id:
            return "❌ Error: contact_id is required"

        graph_client = await get_graph_client()

        # Get contact name before deleting (for confirmation message)
        try:
            contact = await graph_client.users.by_user_id(mailbox_id).contacts.by_contact_id(contact_id).get()
            contact_name = contact.display_name if contact else "Unknown"
        except Exception:
            contact_name = "Unknown"

        # Delete contact
        await graph_client.users.by_user_id(mailbox_id).contacts.by_contact_id(contact_id).delete()

        return (
            f"✅ Contact deleted successfully!\n\n"
            f"Name: {contact_name}\n"
            f"Contact ID: {contact_id}\n"
            f"Mailbox: {mailbox_id}\n\n"
            f"This action cannot be undone."
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Contact not found\n\n"
                f"Contact ID: {contact_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"The contact may have already been deleted."
            )

        return (
            f"❌ Error deleting contact: {error_type}\n\n"
            f"Contact ID: {contact_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def move_contact(
    contact_id: str,
    mailbox_id: str,
    destination_folder_id: str
) -> str:
    """
    Move a contact to a different folder.

    Args:
        contact_id: Contact to move (required)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        destination_folder_id: Target folder ID (required)

    Returns:
        Confirmation with new location

    Examples:
        move_contact(contact_id="AAMk...", mailbox_id="thomas@sixpillar.co.uk", destination_folder_id="AAMkFolder...")
    """
    try:
        if not contact_id:
            return "❌ Error: contact_id is required"
        if not destination_folder_id:
            return "❌ Error: destination_folder_id is required"

        graph_client = await get_graph_client()

        # Get contact name before moving
        try:
            contact = await graph_client.users.by_user_id(mailbox_id).contacts.by_contact_id(contact_id).get()
            contact_name = contact.display_name if contact else "Unknown"
        except Exception:
            contact_name = "Unknown"

        # Get destination folder name
        try:
            folder = await graph_client.users.by_user_id(mailbox_id).contact_folders.by_contact_folder_id(destination_folder_id).get()
            folder_name = folder.display_name if folder else "Unknown"
        except Exception:
            folder_name = "Unknown"

        # Move contact by updating parentFolderId via PATCH
        # Note: Contacts API doesn't have a move action, use PATCH instead
        update_contact = Contact()
        update_contact.parent_folder_id = destination_folder_id

        moved = await graph_client.users.by_user_id(mailbox_id).contacts.by_contact_id(contact_id).patch(body=update_contact)

        return (
            f"✅ Contact moved successfully!\n\n"
            f"Name: {contact_name}\n"
            f"Destination: {folder_name}\n"
            f"Contact ID: {moved.id if moved else contact_id}\n"
            f"Mailbox: {mailbox_id}"
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Contact or folder not found\n\n"
                f"Contact ID: {contact_id}\n"
                f"Folder ID: {destination_folder_id}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Ensure both the contact and destination folder exist."
            )

        return (
            f"❌ Error moving contact: {error_type}\n\n"
            f"Contact ID: {contact_id}\n"
            f"Destination: {destination_folder_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def create_contact_folder(
    mailbox_id: str,
    folder_name: str
) -> str:
    """
    Create a new contact folder for organizing contacts.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        folder_name: Name for the new folder (required)

    Returns:
        Created folder details with ID

    Examples:
        create_contact_folder(mailbox_id="thomas@sixpillar.co.uk", folder_name="Suppliers")
        create_contact_folder(mailbox_id="me", folder_name="Clients")
    """
    try:
        if not folder_name:
            return "❌ Error: folder_name is required"

        graph_client = await get_graph_client()

        # Create new folder object
        new_folder = ContactFolder()
        new_folder.display_name = folder_name

        # Create folder
        created = await graph_client.users.by_user_id(mailbox_id).contact_folders.post(body=new_folder)

        if not created:
            return "❌ Error: Folder creation failed - no response"

        return (
            f"✅ Contact folder created successfully!\n\n"
            f"Name: {created.display_name or folder_name}\n"
            f"Folder ID: {created.id}\n"
            f"Mailbox: {mailbox_id}\n\n"
            f"You can now use this folder_id to create or move contacts into it."
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "NameAlreadyExists" in error_msg or "already exists" in error_msg.lower():
            return (
                f"❌ Folder already exists\n\n"
                f"Folder name: {folder_name}\n"
                f"Mailbox: {mailbox_id}\n\n"
                f"Use a different name or list_contact_folders to find the existing folder ID."
            )

        return (
            f"❌ Error creating contact folder: {error_type}\n\n"
            f"Folder name: {folder_name}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


# =============================================================================
# Phase 9 M2: Calendar Search & Availability Tools
# =============================================================================

@mcp.tool()
async def search_calendar_events(
    mailbox_id: str,
    query: str,
    start_date: str = "",
    end_date: str = "",
    count: int = 20
) -> str:
    """
    Search calendar events by keyword.

    Searches subject, body, location, and attendee names.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        query: Search keywords (required)
        start_date: Start of search range in YYYY-MM-DD format (default: 30 days ago)
        end_date: End of search range in YYYY-MM-DD format (default: 90 days ahead)
        count: Maximum results to return (default: 20, max: 50)

    Returns:
        Matching events with relevance-based ordering

    Examples:
        search_calendar_events(mailbox_id="thomas@sixpillar.co.uk", query="project review")
        search_calendar_events(mailbox_id="me", query="standup", start_date="2025-12-01")
    """
    try:
        if not query:
            return "❌ Error: query is required"

        graph_client = await get_graph_client()

        # Parse dates - default to -30 days to +90 days
        if start_date:
            try:
                start_dt = datetime.fromisoformat(start_date)
            except ValueError:
                return f"❌ Invalid start_date format: {start_date}. Use YYYY-MM-DD."
        else:
            start_dt = datetime.now() - timedelta(days=30)
        start_dt = start_dt.replace(hour=0, minute=0, second=0, microsecond=0)

        if end_date:
            try:
                end_dt = datetime.fromisoformat(end_date)
            except ValueError:
                return f"❌ Invalid end_date format: {end_date}. Use YYYY-MM-DD."
        else:
            end_dt = datetime.now() + timedelta(days=90)
        end_dt = end_dt.replace(hour=23, minute=59, second=59)

        # Limit count
        count = min(count, 50)

        # Format for Graph API
        start_iso = start_dt.strftime("%Y-%m-%dT%H:%M:%S")
        end_iso = end_dt.strftime("%Y-%m-%dT%H:%M:%S")

        # Build filter - search in subject, location, and body
        # OData $filter with contains() for basic search
        escaped_query = escape_odata_string(query)
        filter_str = f"contains(subject, '{escaped_query}') or contains(location/displayName, '{escaped_query}')"

        # Build query params - use calendarView to get events in range, then filter
        query_params = CalendarViewRequestBuilder.CalendarViewRequestBuilderGetQueryParameters(
            start_date_time=start_iso,
            end_date_time=end_iso,
            top=100,  # Get more events to filter client-side
            orderby=["start/dateTime"],
            select=["id", "subject", "start", "end", "location", "attendees", "isAllDay", "organizer", "isCancelled", "bodyPreview"]
        )
        request_config = RequestConfiguration(query_parameters=query_params)

        # Get events
        events_response = await graph_client.users.by_user_id(mailbox_id).calendar.calendar_view.get(
            request_configuration=request_config
        )

        events = events_response.value if events_response and events_response.value else []

        # Filter out cancelled events
        events = [e for e in events if not e.is_cancelled]

        # Client-side search (more flexible than OData filter)
        query_lower = query.lower()
        matching_events = []

        for event in events:
            # Search in subject
            if event.subject and query_lower in event.subject.lower():
                matching_events.append(event)
                continue

            # Search in location
            if event.location and event.location.display_name:
                if query_lower in event.location.display_name.lower():
                    matching_events.append(event)
                    continue

            # Search in body preview
            if event.body_preview and query_lower in event.body_preview.lower():
                matching_events.append(event)
                continue

            # Search in attendees
            if event.attendees:
                for attendee in event.attendees:
                    if attendee.email_address:
                        name = attendee.email_address.name or ""
                        email = attendee.email_address.address or ""
                        if query_lower in name.lower() or query_lower in email.lower():
                            matching_events.append(event)
                            break

        # Limit results
        matching_events = matching_events[:count]

        if not matching_events:
            return (
                f"No events found matching '{query}'\n\n"
                f"Search range: {start_dt.strftime('%d %b %Y')} - {end_dt.strftime('%d %b %Y')}\n"
                f"Mailbox: {mailbox_id}"
            )

        # Format output
        lines = [
            f"Search Results for '{query}'",
            f"Mailbox: {mailbox_id}",
            f"Range: {start_dt.strftime('%d %b %Y')} - {end_dt.strftime('%d %b %Y')}",
            ""
        ]

        for i, event in enumerate(matching_events, 1):
            # Parse times
            if event.is_all_day:
                start_str = "All day"
                end_str = ""
            else:
                start_time = event.start
                end_time = event.end
                if start_time and start_time.date_time:
                    start_parsed = datetime.fromisoformat(start_time.date_time.replace("Z", ""))
                    start_str = start_parsed.strftime("%a %d %b %Y, %H:%M")
                else:
                    start_str = "Unknown"

                if end_time and end_time.date_time:
                    end_parsed = datetime.fromisoformat(end_time.date_time.replace("Z", ""))
                    end_str = f"-{end_parsed.strftime('%H:%M')}"
                else:
                    end_str = ""

            lines.append(f"{i}. {event.subject or '(No subject)'}")
            lines.append(f"   When: {start_str}{end_str}")

            if event.location and event.location.display_name:
                lines.append(f"   Location: {event.location.display_name}")

            if event.attendees:
                attendee_count = len(event.attendees)
                if attendee_count <= 3:
                    names = [a.email_address.name or a.email_address.address for a in event.attendees if a.email_address]
                    lines.append(f"   Attendees: {', '.join(names)}")
                else:
                    lines.append(f"   Attendees: {attendee_count} people")

            lines.append(f"   [ID: {event.id}]")
            lines.append("")

        lines.append("---")
        lines.append(f"Found {len(matching_events)} matching event(s)")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        return (
            f"❌ Error searching calendar events: {error_type}\n\n"
            f"Query: {query}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def check_availability(
    mailbox_id: str,
    start_datetime: str,
    end_datetime: str,
    attendees: str = ""
) -> str:
    """
    Check availability for a specific time slot.

    Can check single user or multiple attendees for meeting scheduling.
    Uses the Microsoft Graph getSchedule API to check free/busy information.

    Args:
        mailbox_id: Primary mailbox to check (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        start_datetime: Proposed start time (YYYY-MM-DDTHH:MM format)
        end_datetime: Proposed end time (YYYY-MM-DDTHH:MM format)
        attendees: Additional attendees to check (comma-separated email addresses)

    Returns:
        Availability status for each person, conflicts if any

    Examples:
        check_availability(mailbox_id="thomas@sixpillar.co.uk",
                          start_datetime="2025-12-01T14:00",
                          end_datetime="2025-12-01T15:00")
        check_availability(mailbox_id="thomas@sixpillar.co.uk",
                          start_datetime="2025-12-01T14:00",
                          end_datetime="2025-12-01T15:00",
                          attendees="john@example.com,jane@example.com")
    """
    try:
        graph_client = await get_graph_client()

        # Parse datetimes
        try:
            if "T" in start_datetime:
                start_dt = datetime.fromisoformat(start_datetime)
            else:
                return f"❌ Invalid start_datetime format: {start_datetime}. Use YYYY-MM-DDTHH:MM."
        except ValueError:
            return f"❌ Invalid start_datetime format: {start_datetime}. Use YYYY-MM-DDTHH:MM."

        try:
            if "T" in end_datetime:
                end_dt = datetime.fromisoformat(end_datetime)
            else:
                return f"❌ Invalid end_datetime format: {end_datetime}. Use YYYY-MM-DDTHH:MM."
        except ValueError:
            return f"❌ Invalid end_datetime format: {end_datetime}. Use YYYY-MM-DDTHH:MM."

        # Build list of schedules to check
        schedules_to_check = [mailbox_id]
        if attendees:
            for email in attendees.split(","):
                email = email.strip()
                if email and email not in schedules_to_check:
                    schedules_to_check.append(email)

        # Format times for API
        start_iso = start_dt.strftime("%Y-%m-%dT%H:%M:%S")
        end_iso = end_dt.strftime("%Y-%m-%dT%H:%M:%S")

        # Use getSchedule API via calendar view for each user's availability
        # Note: getSchedule API requires specific permissions and may not work for external users
        # Fall back to checking calendar view for the primary user

        results = []
        conflicts = []
        free_count = 0
        external_users = []

        for schedule_email in schedules_to_check:
            try:
                # Check calendar view for this time range
                query_params = CalendarViewRequestBuilder.CalendarViewRequestBuilderGetQueryParameters(
                    start_date_time=start_iso,
                    end_date_time=end_iso,
                    select=["id", "subject", "start", "end", "showAs", "isCancelled"]
                )
                request_config = RequestConfiguration(query_parameters=query_params)

                events_response = await graph_client.users.by_user_id(schedule_email).calendar.calendar_view.get(
                    request_configuration=request_config
                )

                events = events_response.value if events_response and events_response.value else []
                # Filter to non-cancelled events that show as busy
                busy_events = [
                    e for e in events
                    if not e.is_cancelled and e.show_as in [FreeBusyStatus.Busy, FreeBusyStatus.Tentative, FreeBusyStatus.Oof]
                ]

                if busy_events:
                    # Has conflicts
                    event_names = [e.subject or "(No subject)" for e in busy_events[:2]]
                    if len(busy_events) > 2:
                        event_names.append(f"and {len(busy_events) - 2} more")

                    status = "BUSY"
                    conflict_detail = f"({', '.join(event_names)})"
                    conflicts.append(schedule_email)
                else:
                    status = "FREE"
                    conflict_detail = ""
                    free_count += 1

                results.append((schedule_email, status, conflict_detail))

            except Exception as user_error:
                error_msg = str(user_error)
                if "ResourceNotFound" in error_msg or "MailboxNotFound" in error_msg:
                    # External user or no access
                    results.append((schedule_email, "UNKNOWN", "(external user or no access)"))
                    external_users.append(schedule_email)
                else:
                    results.append((schedule_email, "ERROR", f"({str(user_error)[:50]})"))

        # Format output
        lines = [
            f"Availability Check: {start_dt.strftime('%a %d %b %Y')}, {start_dt.strftime('%H:%M')}-{end_dt.strftime('%H:%M')}",
            ""
        ]

        for email, status, detail in results:
            status_icon = "✅" if status == "FREE" else "❌" if status == "BUSY" else "❓"
            lines.append(f"{email}: {status_icon} {status} {detail}")

        lines.append("")

        # Summary
        if conflicts:
            lines.append(f"Result: ❌ Time slot NOT available ({len(conflicts)} conflict(s))")
            if free_count > 0:
                lines.append(f"        {free_count} attendee(s) are free")
        elif external_users and len(external_users) == len(schedules_to_check):
            lines.append(f"Result: ❓ Could not determine availability (all external users)")
        else:
            lines.append(f"Result: ✅ Time slot AVAILABLE (all {free_count} attendee(s) free)")

        if external_users:
            lines.append("")
            lines.append(f"Note: Could not check {len(external_users)} external user(s)")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        return (
            f"❌ Error checking availability: {error_type}\n\n"
            f"Time: {start_datetime} - {end_datetime}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def find_meeting_times(
    mailbox_id: str,
    attendees: str,
    duration_minutes: int = 60,
    start_date: str = "",
    end_date: str = "",
    working_hours_only: bool = True
) -> str:
    """
    Find available meeting times for multiple attendees.

    Analyzes calendars and suggests optimal time slots when all (or most)
    attendees are available.

    Args:
        mailbox_id: Organizer's mailbox (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        attendees: Required attendees (comma-separated email addresses)
        duration_minutes: Meeting duration in minutes (default: 60)
        start_date: Start searching from in YYYY-MM-DD format (default: tomorrow)
        end_date: Search until in YYYY-MM-DD format (default: 7 days out)
        working_hours_only: Limit to business hours 09:00-17:00 (default: True)

    Returns:
        List of suggested time slots when attendees are available

    Examples:
        find_meeting_times(mailbox_id="thomas@sixpillar.co.uk",
                          attendees="john@example.com,jane@example.com")
        find_meeting_times(mailbox_id="thomas@sixpillar.co.uk",
                          attendees="john@example.com",
                          duration_minutes=30,
                          start_date="2025-12-01",
                          end_date="2025-12-07")
    """
    try:
        if not attendees:
            return "❌ Error: attendees is required (comma-separated email addresses)"

        graph_client = await get_graph_client()

        # Parse attendee list
        attendee_list = [mailbox_id]  # Include organizer
        for email in attendees.split(","):
            email = email.strip()
            if email and email not in attendee_list:
                attendee_list.append(email)

        # Parse dates - default to tomorrow to +7 days
        if start_date:
            try:
                start_dt = datetime.fromisoformat(start_date)
            except ValueError:
                return f"❌ Invalid start_date format: {start_date}. Use YYYY-MM-DD."
        else:
            start_dt = datetime.now() + timedelta(days=1)
        start_dt = start_dt.replace(hour=0, minute=0, second=0, microsecond=0)

        if end_date:
            try:
                end_dt = datetime.fromisoformat(end_date)
            except ValueError:
                return f"❌ Invalid end_date format: {end_date}. Use YYYY-MM-DD."
        else:
            end_dt = start_dt + timedelta(days=7)
        end_dt = end_dt.replace(hour=23, minute=59, second=59)

        # Working hours
        work_start_hour = 9 if working_hours_only else 0
        work_end_hour = 17 if working_hours_only else 23

        # Find free slots by analyzing each day
        suggested_slots = []

        current_date = start_dt
        while current_date <= end_dt and len(suggested_slots) < 10:
            # Skip weekends if working hours only
            if working_hours_only and current_date.weekday() >= 5:
                current_date += timedelta(days=1)
                continue

            # Check slots throughout the day
            slot_start = current_date.replace(hour=work_start_hour, minute=0)
            day_end = current_date.replace(hour=work_end_hour, minute=0)

            while slot_start + timedelta(minutes=duration_minutes) <= day_end and len(suggested_slots) < 10:
                slot_end = slot_start + timedelta(minutes=duration_minutes)

                # Check availability for all attendees
                all_free = True
                some_unknown = False

                for attendee_email in attendee_list:
                    try:
                        start_iso = slot_start.strftime("%Y-%m-%dT%H:%M:%S")
                        end_iso = slot_end.strftime("%Y-%m-%dT%H:%M:%S")

                        query_params = CalendarViewRequestBuilder.CalendarViewRequestBuilderGetQueryParameters(
                            start_date_time=start_iso,
                            end_date_time=end_iso,
                            select=["id", "showAs", "isCancelled"]
                        )
                        request_config = RequestConfiguration(query_parameters=query_params)

                        events_response = await graph_client.users.by_user_id(attendee_email).calendar.calendar_view.get(
                            request_configuration=request_config
                        )

                        events = events_response.value if events_response and events_response.value else []
                        busy_events = [
                            e for e in events
                            if not e.is_cancelled and e.show_as in [FreeBusyStatus.Busy, FreeBusyStatus.Tentative, FreeBusyStatus.Oof]
                        ]

                        if busy_events:
                            all_free = False
                            break

                    except Exception:
                        # External user or no access - mark as unknown
                        some_unknown = True

                if all_free:
                    confidence = "MEDIUM" if some_unknown else "HIGH"
                    suggested_slots.append((slot_start, slot_end, confidence))

                # Move to next slot (30 min increments)
                slot_start += timedelta(minutes=30)

            current_date += timedelta(days=1)

        if not suggested_slots:
            return (
                f"No available meeting times found\n\n"
                f"Duration: {duration_minutes} minutes\n"
                f"Search range: {start_dt.strftime('%d %b %Y')} - {end_dt.strftime('%d %b %Y')}\n"
                f"Attendees: {', '.join(attendee_list)}\n\n"
                f"Try extending the search range or reducing the duration."
            )

        # Format output
        lines = [
            f"Finding {duration_minutes}-minute meeting slots for:",
            f"  - {mailbox_id} (organizer)"
        ]
        for email in attendee_list[1:]:
            lines.append(f"  - {email}")
        lines.append("")
        lines.append(f"Suggested Times ({start_dt.strftime('%d %b')} - {end_dt.strftime('%d %b %Y')}):")
        lines.append("")

        for i, (slot_start, slot_end, confidence) in enumerate(suggested_slots, 1):
            conf_icon = "✅" if confidence == "HIGH" else "⚠️"
            conf_note = "(all attendees free)" if confidence == "HIGH" else "(some external attendees)"
            lines.append(f"{i}. {slot_start.strftime('%a %d %b %Y')}, {slot_start.strftime('%H:%M')}-{slot_end.strftime('%H:%M')}")
            lines.append(f"   Confidence: {conf_icon} {confidence} {conf_note}")
            lines.append("")

        lines.append("---")
        lines.append(f"Found {len(suggested_slots)} suitable slot(s)")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        return (
            f"❌ Error finding meeting times: {error_type}\n\n"
            f"Attendees: {attendees}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


# =============================================================================
# Phase 9 M3: Calendar Create & Manage Tools
# =============================================================================

@mcp.tool()
async def create_calendar_event(
    mailbox_id: str,
    subject: str,
    start_datetime: str,
    end_datetime: str,
    location: str = "",
    body: str = "",
    is_all_day: bool = False,
    reminder_minutes: int = 15,
    calendar_id: str = ""
) -> str:
    """
    Create a calendar event (no attendees - personal event).

    For events with attendees (meetings), use create_meeting instead.

    Args:
        mailbox_id: Email address of mailbox (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        subject: Event title (required)
        start_datetime: Start time in YYYY-MM-DDTHH:MM format (required)
        end_datetime: End time in YYYY-MM-DDTHH:MM format (required)
        location: Location text (optional)
        body: Event description/notes (optional)
        is_all_day: Create as all-day event (default: False)
        reminder_minutes: Reminder before event in minutes (default: 15, 0 to disable)
        calendar_id: Target calendar ID (default: primary calendar)

    Returns:
        Created event details with ID

    Examples:
        create_calendar_event(mailbox_id="thomas@sixpillar.co.uk",
                             subject="Project Review",
                             start_datetime="2025-12-01T14:00",
                             end_datetime="2025-12-01T15:00")
        create_calendar_event(mailbox_id="thomas@sixpillar.co.uk",
                             subject="Holiday",
                             start_datetime="2025-12-25",
                             end_datetime="2025-12-25",
                             is_all_day=True)
    """
    try:
        if not subject:
            return "❌ Error: subject is required"

        graph_client = await get_graph_client()

        # Parse datetimes
        if is_all_day:
            # All-day events use date only
            try:
                if "T" in start_datetime:
                    start_dt = datetime.fromisoformat(start_datetime.split("T")[0])
                else:
                    start_dt = datetime.fromisoformat(start_datetime)
                if "T" in end_datetime:
                    end_dt = datetime.fromisoformat(end_datetime.split("T")[0])
                else:
                    end_dt = datetime.fromisoformat(end_datetime)
                # All-day events end the next day in Graph API
                end_dt = end_dt + timedelta(days=1)
            except ValueError as ve:
                return f"❌ Invalid date format: {ve}. Use YYYY-MM-DD for all-day events."
        else:
            try:
                start_dt = datetime.fromisoformat(start_datetime)
                end_dt = datetime.fromisoformat(end_datetime)
            except ValueError as ve:
                return f"❌ Invalid datetime format: {ve}. Use YYYY-MM-DDTHH:MM."

        # Create event object
        new_event = Event()
        new_event.subject = subject

        # Set times
        start_tz = DateTimeTimeZone()
        end_tz = DateTimeTimeZone()

        if is_all_day:
            new_event.is_all_day = True
            start_tz.date_time = start_dt.strftime("%Y-%m-%dT00:00:00")
            start_tz.time_zone = "UTC"
            end_tz.date_time = end_dt.strftime("%Y-%m-%dT00:00:00")
            end_tz.time_zone = "UTC"
        else:
            new_event.is_all_day = False
            start_tz.date_time = start_dt.strftime("%Y-%m-%dT%H:%M:%S")
            start_tz.time_zone = MCP_MS_GRAPH_TIMEZONE  # Default timezone
            end_tz.date_time = end_dt.strftime("%Y-%m-%dT%H:%M:%S")
            end_tz.time_zone = MCP_MS_GRAPH_TIMEZONE

        new_event.start = start_tz
        new_event.end = end_tz

        # Optional fields
        if location:
            from msgraph.generated.models.location import Location
            loc = Location()
            loc.display_name = location
            new_event.location = loc

        if body:
            event_body = ItemBody()
            event_body.content = body
            event_body.content_type = BodyType.Text
            new_event.body = event_body

        if reminder_minutes > 0:
            new_event.is_reminder_on = True
            new_event.reminder_minutes_before_start = reminder_minutes
        else:
            new_event.is_reminder_on = False

        # Create the event
        if calendar_id:
            created = await graph_client.users.by_user_id(mailbox_id).calendars.by_calendar_id(calendar_id).events.post(body=new_event)
        else:
            created = await graph_client.users.by_user_id(mailbox_id).calendar.events.post(body=new_event)

        if not created:
            return "❌ Error: Event creation failed - no response"

        # Format output
        if is_all_day:
            time_str = f"{start_dt.strftime('%a %d %b %Y')} (all day)"
        else:
            time_str = f"{start_dt.strftime('%a %d %b %Y')}, {start_dt.strftime('%H:%M')}-{end_dt.strftime('%H:%M')}"

        lines = [
            "✅ Event Created Successfully",
            "=" * 30,
            "",
            f"Subject: {subject}",
            f"When: {time_str}",
        ]

        if location:
            lines.append(f"Location: {location}")

        if reminder_minutes > 0:
            lines.append(f"Reminder: {reminder_minutes} minutes before")

        lines.append("")
        lines.append(f"Event ID: {created.id}")
        lines.append(f"Mailbox: {mailbox_id}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        return (
            f"❌ Error creating event: {error_type}\n\n"
            f"Subject: {subject}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def create_meeting(
    mailbox_id: str,
    subject: str,
    start_datetime: str,
    end_datetime: str,
    required_attendees: str,
    optional_attendees: str = "",
    location: str = "",
    body: str = "",
    is_online_meeting: bool = True,
    reminder_minutes: int = 15
) -> str:
    """
    Create a meeting with attendees (sends calendar invites).

    Automatically creates a Teams meeting link if is_online_meeting=True.

    Args:
        mailbox_id: Organizer's mailbox (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        subject: Meeting title (required)
        start_datetime: Start time in YYYY-MM-DDTHH:MM format (required)
        end_datetime: End time in YYYY-MM-DDTHH:MM format (required)
        required_attendees: Required attendees (comma-separated email addresses)
        optional_attendees: Optional attendees (comma-separated email addresses)
        location: Physical location (optional, overridden if online meeting)
        body: Meeting agenda/description (optional)
        is_online_meeting: Create Teams meeting link (default: True)
        reminder_minutes: Reminder before meeting in minutes (default: 15)

    Returns:
        Created meeting details with ID and Teams link (if applicable)

    Examples:
        create_meeting(mailbox_id="thomas@sixpillar.co.uk",
                      subject="Project Kickoff",
                      start_datetime="2025-12-01T10:00",
                      end_datetime="2025-12-01T11:00",
                      required_attendees="john@example.com,jane@example.com")
    """
    try:
        if not subject:
            return "❌ Error: subject is required"
        if not required_attendees:
            return "❌ Error: required_attendees is required (use create_calendar_event for events without attendees)"

        graph_client = await get_graph_client()

        # Parse datetimes
        try:
            start_dt = datetime.fromisoformat(start_datetime)
            end_dt = datetime.fromisoformat(end_datetime)
        except ValueError as ve:
            return f"❌ Invalid datetime format: {ve}. Use YYYY-MM-DDTHH:MM."

        # Create event object
        new_event = Event()
        new_event.subject = subject
        new_event.is_all_day = False

        # Set times
        start_tz = DateTimeTimeZone()
        start_tz.date_time = start_dt.strftime("%Y-%m-%dT%H:%M:%S")
        start_tz.time_zone = MCP_MS_GRAPH_TIMEZONE
        new_event.start = start_tz

        end_tz = DateTimeTimeZone()
        end_tz.date_time = end_dt.strftime("%Y-%m-%dT%H:%M:%S")
        end_tz.time_zone = MCP_MS_GRAPH_TIMEZONE
        new_event.end = end_tz

        # Parse attendees
        attendees_list = []
        from msgraph.generated.models.attendee_type import AttendeeType

        for email in required_attendees.split(","):
            email = email.strip()
            if email:
                attendee = Attendee()
                attendee.email_address = EmailAddress()
                attendee.email_address.address = email
                attendee.type = AttendeeType.Required
                attendees_list.append(attendee)

        if optional_attendees:
            for email in optional_attendees.split(","):
                email = email.strip()
                if email:
                    attendee = Attendee()
                    attendee.email_address = EmailAddress()
                    attendee.email_address.address = email
                    attendee.type = AttendeeType.Optional
                    attendees_list.append(attendee)

        new_event.attendees = attendees_list

        # Online meeting
        if is_online_meeting:
            new_event.is_online_meeting = True
            from msgraph.generated.models.online_meeting_provider_type import OnlineMeetingProviderType
            new_event.online_meeting_provider = OnlineMeetingProviderType.TeamsForBusiness

        # Location
        if location and not is_online_meeting:
            from msgraph.generated.models.location import Location
            loc = Location()
            loc.display_name = location
            new_event.location = loc

        # Body
        if body:
            event_body = ItemBody()
            event_body.content = body
            event_body.content_type = BodyType.Text
            new_event.body = event_body

        # Reminder
        if reminder_minutes > 0:
            new_event.is_reminder_on = True
            new_event.reminder_minutes_before_start = reminder_minutes
        else:
            new_event.is_reminder_on = False

        # Create the meeting
        created = await graph_client.users.by_user_id(mailbox_id).calendar.events.post(body=new_event)

        if not created:
            return "❌ Error: Meeting creation failed - no response"

        # Format output
        time_str = f"{start_dt.strftime('%a %d %b %Y')}, {start_dt.strftime('%H:%M')}-{end_dt.strftime('%H:%M')}"

        lines = [
            "✅ Meeting Created Successfully",
            "=" * 30,
            "",
            f"Subject: {subject}",
            f"When: {time_str}",
        ]

        if is_online_meeting:
            lines.append("Location: Microsoft Teams Meeting")
            if created.online_meeting and created.online_meeting.join_url:
                lines.append(f"Teams Link: {created.online_meeting.join_url}")
        elif location:
            lines.append(f"Location: {location}")

        lines.append("")
        lines.append("Attendees:")
        req_count = 0
        opt_count = 0
        for email in required_attendees.split(","):
            email = email.strip()
            if email:
                lines.append(f"  - {email} (required) - Invite sent")
                req_count += 1
        if optional_attendees:
            for email in optional_attendees.split(","):
                email = email.strip()
                if email:
                    lines.append(f"  - {email} (optional) - Invite sent")
                    opt_count += 1

        lines.append("")
        lines.append(f"Event ID: {created.id}")
        lines.append(f"Mailbox: {mailbox_id}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        return (
            f"❌ Error creating meeting: {error_type}\n\n"
            f"Subject: {subject}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def update_calendar_event(
    event_id: str,
    mailbox_id: str,
    subject: str = "",
    start_datetime: str = "",
    end_datetime: str = "",
    location: str = "",
    body: str = ""
) -> str:
    """
    Update an existing calendar event.

    Only provided parameters are updated - others remain unchanged.
    If the event has attendees, they will receive update notifications.

    Args:
        event_id: Event to update (required)
        mailbox_id: Mailbox containing event (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        subject: New subject (optional - leave empty to keep current)
        start_datetime: New start time in YYYY-MM-DDTHH:MM format (optional)
        end_datetime: New end time in YYYY-MM-DDTHH:MM format (optional)
        location: New location (optional - leave empty to keep current)
        body: New description (optional - leave empty to keep current)

    Returns:
        Updated event details

    Examples:
        update_calendar_event(event_id="AAMk...",
                             mailbox_id="thomas@sixpillar.co.uk",
                             subject="Updated: Project Review")
        update_calendar_event(event_id="AAMk...",
                             mailbox_id="thomas@sixpillar.co.uk",
                             start_datetime="2025-12-01T15:00",
                             end_datetime="2025-12-01T16:00")
    """
    try:
        if not event_id:
            return "❌ Error: event_id is required"

        # Check if any update is requested
        if not any([subject, start_datetime, end_datetime, location, body]):
            return "❌ Error: At least one field to update is required"

        graph_client = await get_graph_client()

        # Build update object with only specified fields
        update_event = Event()
        updated_fields = []

        if subject:
            update_event.subject = subject
            updated_fields.append("subject")

        if start_datetime:
            try:
                start_dt = datetime.fromisoformat(start_datetime)
                start_tz = DateTimeTimeZone()
                start_tz.date_time = start_dt.strftime("%Y-%m-%dT%H:%M:%S")
                start_tz.time_zone = MCP_MS_GRAPH_TIMEZONE
                update_event.start = start_tz
                updated_fields.append("start time")
            except ValueError as ve:
                return f"❌ Invalid start_datetime format: {ve}. Use YYYY-MM-DDTHH:MM."

        if end_datetime:
            try:
                end_dt = datetime.fromisoformat(end_datetime)
                end_tz = DateTimeTimeZone()
                end_tz.date_time = end_dt.strftime("%Y-%m-%dT%H:%M:%S")
                end_tz.time_zone = MCP_MS_GRAPH_TIMEZONE
                update_event.end = end_tz
                updated_fields.append("end time")
            except ValueError as ve:
                return f"❌ Invalid end_datetime format: {ve}. Use YYYY-MM-DDTHH:MM."

        if location:
            from msgraph.generated.models.location import Location
            loc = Location()
            loc.display_name = location
            update_event.location = loc
            updated_fields.append("location")

        if body:
            event_body = ItemBody()
            event_body.content = body
            event_body.content_type = BodyType.Text
            update_event.body = event_body
            updated_fields.append("body")

        # Update the event
        updated = await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).patch(body=update_event)

        if not updated:
            return "❌ Error: Event update failed - no response"

        # Format output
        lines = [
            "✅ Event Updated Successfully",
            "=" * 30,
            "",
            f"Updated fields: {', '.join(updated_fields)}",
            "",
            f"Subject: {updated.subject or '(unchanged)'}",
        ]

        if updated.start and updated.start.date_time:
            start_parsed = datetime.fromisoformat(updated.start.date_time.replace("Z", ""))
            lines.append(f"Start: {start_parsed.strftime('%a %d %b %Y, %H:%M')}")

        if updated.end and updated.end.date_time:
            end_parsed = datetime.fromisoformat(updated.end.date_time.replace("Z", ""))
            lines.append(f"End: {end_parsed.strftime('%H:%M')}")

        if updated.location and updated.location.display_name:
            lines.append(f"Location: {updated.location.display_name}")

        lines.append("")
        lines.append(f"Event ID: {event_id}")
        lines.append(f"Mailbox: {mailbox_id}")

        if updated.attendees and len(updated.attendees) > 0:
            lines.append("")
            lines.append("Note: Attendees have been notified of this change.")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Event not found\n\n"
                f"Event ID: {event_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error updating event: {error_type}\n\n"
            f"Event ID: {event_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def delete_calendar_event(
    event_id: str,
    mailbox_id: str,
    cancellation_message: str = ""
) -> str:
    """
    Delete a calendar event or cancel a meeting.

    If the event has attendees, sends cancellation notice to all attendees.

    Args:
        event_id: Event to delete (required)
        mailbox_id: Mailbox containing event (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        cancellation_message: Message to send with cancellation (optional)

    Returns:
        Confirmation of deletion/cancellation

    Examples:
        delete_calendar_event(event_id="AAMk...",
                             mailbox_id="thomas@sixpillar.co.uk")
        delete_calendar_event(event_id="AAMk...",
                             mailbox_id="thomas@sixpillar.co.uk",
                             cancellation_message="Meeting cancelled due to scheduling conflict")
    """
    try:
        if not event_id:
            return "❌ Error: event_id is required"

        graph_client = await get_graph_client()

        # First, get the event to check if it has attendees
        try:
            existing = await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).get()
            event_subject = existing.subject if existing else "(Unknown)"
            has_attendees = existing.attendees and len(existing.attendees) > 0 if existing else False
            attendee_count = len(existing.attendees) if has_attendees else 0
        except Exception:
            event_subject = "(Unknown)"
            has_attendees = False
            attendee_count = 0

        # If cancellation message provided and has attendees, we should use cancel endpoint
        # For now, just delete (which auto-sends cancellation)
        if cancellation_message and has_attendees:
            # Update body with cancellation message before deleting
            try:
                cancel_update = Event()
                cancel_body = ItemBody()
                cancel_body.content = f"CANCELLED: {cancellation_message}"
                cancel_body.content_type = BodyType.Text
                cancel_update.body = cancel_body
                await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).patch(body=cancel_update)
            except Exception:
                pass  # Continue with deletion even if message update fails

        # Delete the event
        await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).delete()

        # Format output
        lines = [
            "✅ Event Deleted Successfully",
            "=" * 30,
            "",
            f"Subject: {event_subject}",
            f"Event ID: {event_id}",
            f"Mailbox: {mailbox_id}",
        ]

        if has_attendees:
            lines.append("")
            lines.append(f"Note: Cancellation notices sent to {attendee_count} attendee(s)")
            if cancellation_message:
                lines.append(f"Message: {cancellation_message}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Event not found\n\n"
                f"Event ID: {event_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error deleting event: {error_type}\n\n"
            f"Event ID: {event_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


# =============================================================================
# Phase 9 M4: Meeting Response Tools
# =============================================================================

@mcp.tool()
async def respond_to_meeting(
    event_id: str,
    mailbox_id: str,
    response: str,
    message: str = "",
    send_response: bool = True
) -> str:
    """
    Respond to a meeting invitation.

    Args:
        event_id: Meeting event ID (required)
        mailbox_id: Your mailbox (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        response: Response type - "accept", "tentative", or "decline" (required)
        message: Optional message to organizer
        send_response: Send response to organizer (default: True)

    Returns:
        Confirmation of response

    Examples:
        respond_to_meeting(event_id="AAMk...",
                          mailbox_id="thomas@sixpillar.co.uk",
                          response="accept")
        respond_to_meeting(event_id="AAMk...",
                          mailbox_id="thomas@sixpillar.co.uk",
                          response="decline",
                          message="I have a conflict at this time")
    """
    try:
        if not event_id:
            return "❌ Error: event_id is required"

        response_lower = response.lower().strip()
        if response_lower not in ["accept", "tentative", "decline"]:
            return f"❌ Error: response must be 'accept', 'tentative', or 'decline' (got: {response})"

        graph_client = await get_graph_client()

        # Get event details first
        try:
            event = await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).get()
            event_subject = event.subject if event else "(Unknown)"
            organizer = event.organizer.email_address.address if event and event.organizer and event.organizer.email_address else "(Unknown)"
        except Exception:
            event_subject = "(Unknown)"
            organizer = "(Unknown)"

        # Build response body
        from msgraph.generated.users.item.events.item.accept.accept_post_request_body import AcceptPostRequestBody
        from msgraph.generated.users.item.events.item.tentatively_accept.tentatively_accept_post_request_body import TentativelyAcceptPostRequestBody
        from msgraph.generated.users.item.events.item.decline.decline_post_request_body import DeclinePostRequestBody

        if response_lower == "accept":
            request_body = AcceptPostRequestBody()
            request_body.send_response = send_response
            if message:
                request_body.comment = message
            await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).accept.post(body=request_body)
            action = "ACCEPTED"
            icon = "✅"

        elif response_lower == "tentative":
            request_body = TentativelyAcceptPostRequestBody()
            request_body.send_response = send_response
            if message:
                request_body.comment = message
            await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).tentatively_accept.post(body=request_body)
            action = "TENTATIVELY ACCEPTED"
            icon = "⚠️"

        else:  # decline
            request_body = DeclinePostRequestBody()
            request_body.send_response = send_response
            if message:
                request_body.comment = message
            await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).decline.post(body=request_body)
            action = "DECLINED"
            icon = "❌"

        # Format output
        lines = [
            f"{icon} Meeting {action}",
            "=" * 30,
            "",
            f"Subject: {event_subject}",
            f"Organizer: {organizer}",
            f"Response: {action}",
        ]

        if send_response:
            lines.append("Response sent to organizer: Yes")
        else:
            lines.append("Response sent to organizer: No (silent)")

        if message:
            lines.append(f"Message: {message}")

        lines.append("")
        lines.append(f"Event ID: {event_id}")
        lines.append(f"Mailbox: {mailbox_id}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Event not found\n\n"
                f"Event ID: {event_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error responding to meeting: {error_type}\n\n"
            f"Event ID: {event_id}\n"
            f"Response: {response}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def propose_new_time(
    event_id: str,
    mailbox_id: str,
    proposed_start: str,
    proposed_end: str,
    message: str = ""
) -> str:
    """
    Propose a new time for a meeting invitation.

    Declines the current time and suggests an alternative to the organizer.

    Args:
        event_id: Meeting event ID (required)
        mailbox_id: Your mailbox (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes
        proposed_start: Suggested new start time in YYYY-MM-DDTHH:MM format (required)
        proposed_end: Suggested new end time in YYYY-MM-DDTHH:MM format (required)
        message: Explanation for the change (optional)

    Returns:
        Confirmation that proposal was sent

    Examples:
        propose_new_time(event_id="AAMk...",
                        mailbox_id="thomas@sixpillar.co.uk",
                        proposed_start="2025-12-01T15:00",
                        proposed_end="2025-12-01T16:00",
                        message="I have a conflict at the original time")
    """
    try:
        if not event_id:
            return "❌ Error: event_id is required"

        # Parse proposed times
        try:
            start_dt = datetime.fromisoformat(proposed_start)
            end_dt = datetime.fromisoformat(proposed_end)
        except ValueError as ve:
            return f"❌ Invalid datetime format: {ve}. Use YYYY-MM-DDTHH:MM."

        graph_client = await get_graph_client()

        # Get event details first
        try:
            event = await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).get()
            event_subject = event.subject if event else "(Unknown)"
            organizer = event.organizer.email_address.address if event and event.organizer and event.organizer.email_address else "(Unknown)"
        except Exception:
            event_subject = "(Unknown)"
            organizer = "(Unknown)"

        # Build propose new time request
        from msgraph.generated.users.item.events.item.tentatively_accept.tentatively_accept_post_request_body import TentativelyAcceptPostRequestBody
        from msgraph.generated.models.time_slot import TimeSlot

        # Create proposed time slot
        proposed_time = TimeSlot()
        start_tz = DateTimeTimeZone()
        start_tz.date_time = start_dt.strftime("%Y-%m-%dT%H:%M:%S")
        start_tz.time_zone = MCP_MS_GRAPH_TIMEZONE
        proposed_time.start = start_tz

        end_tz = DateTimeTimeZone()
        end_tz.date_time = end_dt.strftime("%Y-%m-%dT%H:%M:%S")
        end_tz.time_zone = MCP_MS_GRAPH_TIMEZONE
        proposed_time.end = end_tz

        # Use tentatively accept with proposed new time
        request_body = TentativelyAcceptPostRequestBody()
        request_body.send_response = True
        request_body.proposed_new_time = proposed_time
        if message:
            request_body.comment = message
        else:
            request_body.comment = f"Proposing new time: {start_dt.strftime('%a %d %b %Y, %H:%M')}-{end_dt.strftime('%H:%M')}"

        await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).tentatively_accept.post(body=request_body)

        # Format output
        lines = [
            "✅ New Time Proposed",
            "=" * 30,
            "",
            f"Subject: {event_subject}",
            f"Organizer: {organizer}",
            "",
            f"Proposed Time: {start_dt.strftime('%a %d %b %Y')}, {start_dt.strftime('%H:%M')}-{end_dt.strftime('%H:%M')}",
        ]

        if message:
            lines.append(f"Message: {message}")

        lines.append("")
        lines.append("The organizer will receive your proposal and can choose to reschedule.")
        lines.append("")
        lines.append(f"Event ID: {event_id}")
        lines.append(f"Mailbox: {mailbox_id}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Event not found\n\n"
                f"Event ID: {event_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error proposing new time: {error_type}\n\n"
            f"Event ID: {event_id}\n"
            f"Proposed: {proposed_start} - {proposed_end}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_meeting_responses(
    event_id: str,
    mailbox_id: str
) -> str:
    """
    Get attendee responses for a meeting (organizer view).

    Shows who has accepted, declined, or not yet responded to a meeting invitation.

    Args:
        event_id: Meeting event ID (required)
        mailbox_id: Organizer's mailbox (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        List of attendees with their response status

    Examples:
        get_meeting_responses(event_id="AAMk...",
                             mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        if not event_id:
            return "❌ Error: event_id is required"

        graph_client = await get_graph_client()

        # Get event with attendees
        event = await graph_client.users.by_user_id(mailbox_id).events.by_event_id(event_id).get()

        if not event:
            return f"❌ Event not found: {event_id}"

        event_subject = event.subject or "(No subject)"

        # Parse event time
        if event.start and event.start.date_time:
            start_parsed = datetime.fromisoformat(event.start.date_time.replace("Z", ""))
            time_str = start_parsed.strftime("%a %d %b %Y, %H:%M")
            if event.end and event.end.date_time:
                end_parsed = datetime.fromisoformat(event.end.date_time.replace("Z", ""))
                time_str += f"-{end_parsed.strftime('%H:%M')}"
        else:
            time_str = "(Unknown time)"

        if not event.attendees or len(event.attendees) == 0:
            return (
                f"Meeting: {event_subject}\n"
                f"When: {time_str}\n\n"
                f"No attendees for this event."
            )

        # Categorize responses
        from msgraph.generated.models.response_type import ResponseType

        accepted = []
        tentative = []
        declined = []
        no_response = []

        for attendee in event.attendees:
            if not attendee.email_address:
                continue

            name = attendee.email_address.name or attendee.email_address.address
            email = attendee.email_address.address or ""

            response_status = attendee.status
            response_type = response_status.response if response_status else None

            attendee_info = {
                "name": name,
                "email": email,
                "message": ""
            }

            if response_type == ResponseType.Accepted:
                accepted.append(attendee_info)
            elif response_type == ResponseType.TentativelyAccepted:
                tentative.append(attendee_info)
            elif response_type == ResponseType.Declined:
                declined.append(attendee_info)
            else:
                no_response.append(attendee_info)

        # Format output
        lines = [
            f"Meeting Responses: {event_subject}",
            "=" * 40,
            "",
            f"When: {time_str}",
            "",
            "Responses:",
        ]

        lines.append(f"  ACCEPTED ({len(accepted)}):")
        if accepted:
            for a in accepted:
                lines.append(f"    - {a['name']} <{a['email']}>")
        else:
            lines.append("    (none)")

        lines.append("")
        lines.append(f"  TENTATIVE ({len(tentative)}):")
        if tentative:
            for a in tentative:
                lines.append(f"    - {a['name']} <{a['email']}>")
        else:
            lines.append("    (none)")

        lines.append("")
        lines.append(f"  DECLINED ({len(declined)}):")
        if declined:
            for a in declined:
                lines.append(f"    - {a['name']} <{a['email']}>")
        else:
            lines.append("    (none)")

        lines.append("")
        lines.append(f"  NO RESPONSE ({len(no_response)}):")
        if no_response:
            for a in no_response:
                lines.append(f"    - {a['name']} <{a['email']}>")
        else:
            lines.append("    (none)")

        lines.append("")
        lines.append(f"Summary: {len(accepted)} accepted, {len(tentative)} tentative, {len(declined)} declined, {len(no_response)} pending")
        lines.append("")
        lines.append(f"Event ID: {event_id}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Event not found\n\n"
                f"Event ID: {event_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error getting meeting responses: {error_type}\n\n"
            f"Event ID: {event_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


# =============================================================================
# PHASE 11: TO DO INTEGRATION
# =============================================================================


@mcp.tool()
async def list_todo_lists(
    mailbox_id: str = "me"
) -> str:
    """
    List all Microsoft To Do task lists.

    Returns all task lists the user has access to, including the default
    "Tasks" list and any custom lists.

    Args:
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
                   Use explicit address like "thomas@sixpillar.co.uk" for delegated mailboxes

    Returns:
        Formatted list of task lists with name, ID, and task count

    Examples:
        list_todo_lists(mailbox_id="me")
        list_todo_lists(mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        client = await get_graph_client()

        # Get task lists
        if mailbox_id == "me":
            result = await asyncio.wait_for(
                client.me.todo.lists.get(),
                timeout=API_TIMEOUT
            )
        else:
            result = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.get(),
                timeout=API_TIMEOUT
            )

        if not result or not result.value:
            return f"No task lists found for {mailbox_id}"

        lines = [
            "Microsoft To Do - Task Lists",
            "=" * 40,
            ""
        ]

        for task_list in result.value:
            display_name = task_list.display_name or "(Unnamed)"
            list_id = task_list.id or ""
            is_owner = task_list.is_owner if hasattr(task_list, 'is_owner') else True
            wellknown = task_list.wellknown_list_name if hasattr(task_list, 'wellknown_list_name') else None

            # Add marker for special lists
            marker = ""
            if wellknown and wellknown.value == "defaultList":
                marker = " [Default]"
            elif wellknown and wellknown.value == "flaggedEmails":
                marker = " [Flagged Emails]"

            owner_marker = "" if is_owner else " (shared)"

            lines.append(f"📋 {display_name}{marker}{owner_marker}")
            lines.append(f"   ID: {list_id}")
            lines.append("")

        lines.append(f"Total: {len(result.value)} list(s)")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        return (
            f"❌ Error listing task lists: {error_type}\n\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_todo_list(
    list_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Get details of a specific To Do task list.

    Args:
        list_id: The ID of the task list (from list_todo_lists)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox

    Returns:
        Task list details including name and properties

    Examples:
        get_todo_list(list_id="AAMk...", mailbox_id="me")
    """
    try:
        client = await get_graph_client()

        if mailbox_id == "me":
            task_list = await asyncio.wait_for(
                client.me.todo.lists.by_todo_task_list_id(list_id).get(),
                timeout=API_TIMEOUT
            )
        else:
            task_list = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.by_todo_task_list_id(list_id).get(),
                timeout=API_TIMEOUT
            )

        if not task_list:
            return f"Task list not found: {list_id}"

        display_name = task_list.display_name or "(Unnamed)"
        is_owner = task_list.is_owner if hasattr(task_list, 'is_owner') else True
        wellknown = task_list.wellknown_list_name if hasattr(task_list, 'wellknown_list_name') else None

        lines = [
            f"Task List: {display_name}",
            "=" * 40,
            "",
            f"ID: {task_list.id}",
            f"Owner: {'Yes' if is_owner else 'No (shared)'}",
        ]

        if wellknown:
            lines.append(f"Type: {wellknown.value}")

        lines.append("")
        lines.append("Use list_todo_tasks to see tasks in this list.")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task list not found\n\n"
                f"List ID: {list_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error getting task list: {error_type}\n\n"
            f"List ID: {list_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def list_todo_tasks(
    list_id: str,
    mailbox_id: str = "me",
    include_completed: bool = False,
    count: int = 25
) -> str:
    """
    List tasks in a Microsoft To Do task list.

    Args:
        list_id: The ID of the task list (from list_todo_lists)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
        include_completed: Include completed tasks (default: False, only active tasks)
        count: Maximum tasks to return (default: 25, max: 100)

    Returns:
        Formatted list of tasks with title, status, due date, importance

    Examples:
        list_todo_tasks(list_id="AAMk...", mailbox_id="me")
        list_todo_tasks(list_id="AAMk...", mailbox_id="me", include_completed=True)
    """
    try:
        client = await get_graph_client()

        # Clamp count
        count = max(1, min(count, 100))

        # Build filter for completed status
        filter_query = None if include_completed else "status ne 'completed'"

        # Build request configuration
        from kiota_abstractions.base_request_configuration import RequestConfiguration

        config = RequestConfiguration()
        config.query_parameters = {
            "$top": count,
            "$orderby": "importance desc,dueDateTime/dateTime asc",
        }
        if filter_query:
            config.query_parameters["$filter"] = filter_query

        if mailbox_id == "me":
            result = await asyncio.wait_for(
                client.me.todo.lists.by_todo_task_list_id(list_id).tasks.get(
                    request_configuration=config
                ),
                timeout=API_TIMEOUT
            )
        else:
            result = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.by_todo_task_list_id(list_id).tasks.get(
                    request_configuration=config
                ),
                timeout=API_TIMEOUT
            )

        if not result or not result.value:
            status_msg = "all tasks" if include_completed else "active tasks"
            return f"No {status_msg} found in this list"

        lines = [
            "To Do Tasks",
            "=" * 40,
            ""
        ]

        for task in result.value:
            title = task.title or "(No title)"
            task_id = task.id or ""

            # Status
            status = task.status.value if task.status else "notStarted"
            status_icon = "✅" if status == "completed" else "⬜"

            # Importance
            importance = task.importance.value if task.importance else "normal"
            importance_marker = "❗" if importance == "high" else ""

            # Due date
            due_str = ""
            if task.due_date_time:
                due_date = task.due_date_time.date_time
                if due_date:
                    # Parse and format
                    try:
                        dt = datetime.fromisoformat(due_date.replace("Z", "+00:00"))
                        due_str = f" (due: {dt.strftime('%Y-%m-%d')})"
                    except:
                        due_str = f" (due: {due_date})"

            lines.append(f"{status_icon} {importance_marker}{title}{due_str}")
            lines.append(f"   ID: {task_id}")

            # Show body preview if exists
            if task.body and task.body.content:
                preview = task.body.content[:100].strip()
                if preview:
                    lines.append(f"   Note: {preview}...")

            lines.append("")

        lines.append(f"Showing: {len(result.value)} task(s)")
        if not include_completed:
            lines.append("(Completed tasks hidden - use include_completed=True to show)")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task list not found\n\n"
                f"List ID: {list_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error listing tasks: {error_type}\n\n"
            f"List ID: {list_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_todo_task(
    list_id: str,
    task_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Get full details of a specific To Do task, including checklist items.

    Args:
        list_id: The ID of the task list containing the task
        task_id: The ID of the task (from list_todo_tasks)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox

    Returns:
        Complete task details including title, body, due date, checklist

    Examples:
        get_todo_task(list_id="AAMk...", task_id="AAMk...", mailbox_id="me")
    """
    try:
        client = await get_graph_client()

        # Get task with checklist items expanded
        config = RequestConfiguration()
        config.query_parameters = {
            "$expand": "checklistItems"
        }

        if mailbox_id == "me":
            task = await asyncio.wait_for(
                client.me.todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).get(
                    request_configuration=config
                ),
                timeout=API_TIMEOUT
            )
        else:
            task = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).get(
                    request_configuration=config
                ),
                timeout=API_TIMEOUT
            )

        if not task:
            return f"Task not found: {task_id}"

        title = task.title or "(No title)"

        # Status
        status = task.status.value if task.status else "notStarted"
        status_display = {
            "notStarted": "Not Started",
            "inProgress": "In Progress",
            "completed": "Completed",
            "waitingOnOthers": "Waiting on Others",
            "deferred": "Deferred"
        }.get(status, status)

        # Importance
        importance = task.importance.value if task.importance else "normal"
        importance_display = importance.capitalize()

        lines = [
            f"Task: {title}",
            "=" * 40,
            "",
            f"Status: {status_display}",
            f"Importance: {importance_display}",
        ]

        # Due date
        if task.due_date_time:
            due_date = task.due_date_time.date_time
            if due_date:
                try:
                    dt = datetime.fromisoformat(due_date.replace("Z", "+00:00"))
                    lines.append(f"Due: {dt.strftime('%Y-%m-%d %H:%M')}")
                except:
                    lines.append(f"Due: {due_date}")

        # Reminder
        if task.reminder_date_time:
            reminder = task.reminder_date_time.date_time
            if reminder:
                try:
                    dt = datetime.fromisoformat(reminder.replace("Z", "+00:00"))
                    lines.append(f"Reminder: {dt.strftime('%Y-%m-%d %H:%M')}")
                except:
                    lines.append(f"Reminder: {reminder}")

        # Recurrence
        if task.recurrence:
            lines.append(f"Recurrence: Yes (pattern configured)")

        # Created/Modified
        if task.created_date_time:
            try:
                dt = datetime.fromisoformat(str(task.created_date_time).replace("Z", "+00:00"))
                lines.append(f"Created: {dt.strftime('%Y-%m-%d %H:%M')}")
            except:
                pass

        if task.last_modified_date_time:
            try:
                dt = datetime.fromisoformat(str(task.last_modified_date_time).replace("Z", "+00:00"))
                lines.append(f"Modified: {dt.strftime('%Y-%m-%d %H:%M')}")
            except:
                pass

        # Body/notes
        lines.append("")
        if task.body and task.body.content:
            content = task.body.content.strip()
            if content:
                lines.append("Notes:")
                lines.append(content[:500])
                if len(content) > 500:
                    lines.append("... (truncated)")
        else:
            lines.append("Notes: (none)")

        # Checklist items
        lines.append("")
        if task.checklist_items and len(task.checklist_items) > 0:
            lines.append(f"Checklist ({len(task.checklist_items)} items):")
            for item in task.checklist_items:
                item_text = item.display_name or "(no text)"
                is_checked = item.is_checked if hasattr(item, 'is_checked') else False
                check_mark = "✅" if is_checked else "⬜"
                lines.append(f"  {check_mark} {item_text}")
                if item.id:
                    lines.append(f"      ID: {item.id}")
        else:
            lines.append("Checklist: (none)")

        # IDs for reference
        lines.append("")
        lines.append(f"Task ID: {task.id}")
        lines.append(f"List ID: {list_id}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task not found\n\n"
                f"Task ID: {task_id}\n"
                f"List ID: {list_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error getting task: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"List ID: {list_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def create_todo_list(
    display_name: str,
    mailbox_id: str = "me"
) -> str:
    """
    Create a new Microsoft To Do task list.

    Args:
        display_name: Name for the new task list (required)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox

    Returns:
        Confirmation with new list ID

    Examples:
        create_todo_list(display_name="Shopping", mailbox_id="me")
        create_todo_list(display_name="Work Projects", mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        client = await get_graph_client()

        # Create the task list
        new_list = TodoTaskList(
            display_name=display_name
        )

        if mailbox_id == "me":
            result = await asyncio.wait_for(
                client.me.todo.lists.post(new_list),
                timeout=API_TIMEOUT
            )
        else:
            result = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.post(new_list),
                timeout=API_TIMEOUT
            )

        if not result:
            return "❌ Failed to create task list"

        return (
            f"✅ Task list created successfully!\n\n"
            f"Name: {result.display_name}\n"
            f"ID: {result.id}\n\n"
            f"Use this ID with list_todo_tasks and create_todo_task."
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        return (
            f"❌ Error creating task list: {error_type}\n\n"
            f"Name: {display_name}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def create_todo_task(
    list_id: str,
    title: str,
    mailbox_id: str = "me",
    body: str = "",
    due_date: str = "",
    importance: str = "normal",
    reminder_date: str = ""
) -> str:
    """
    Create a new task in a Microsoft To Do list.

    Args:
        list_id: The ID of the task list (from list_todo_lists)
        title: Task title (required)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
        body: Task notes/description (optional)
        due_date: Due date in YYYY-MM-DD format (optional)
        importance: Task importance - "low", "normal", or "high" (default: "normal")
        reminder_date: Reminder date/time in YYYY-MM-DDTHH:MM format (optional)

    Returns:
        Confirmation with new task ID

    Examples:
        create_todo_task(list_id="AAMk...", title="Buy groceries", mailbox_id="me")
        create_todo_task(list_id="AAMk...", title="Submit report", due_date="2025-12-25", importance="high")
    """
    try:
        client = await get_graph_client()

        # Build task object
        new_task = TodoTask(
            title=title
        )

        # Set body if provided
        if body:
            new_task.body = ItemBody(
                content=body,
                content_type=BodyType.Text
            )

        # Set due date if provided
        if due_date:
            new_task.due_date_time = DateTimeTimeZone(
                date_time=f"{due_date}T00:00:00",
                time_zone="UTC"
            )

        # Set importance
        importance_map = {
            "low": Importance.Low,
            "normal": Importance.Normal,
            "high": Importance.High
        }
        new_task.importance = importance_map.get(importance.lower(), Importance.Normal)

        # Set reminder if provided
        if reminder_date:
            # Handle both date and datetime formats
            if "T" not in reminder_date:
                reminder_date = f"{reminder_date}T09:00:00"
            new_task.reminder_date_time = DateTimeTimeZone(
                date_time=reminder_date,
                time_zone="UTC"
            )
            new_task.is_reminder_on = True

        if mailbox_id == "me":
            result = await asyncio.wait_for(
                client.me.todo.lists.by_todo_task_list_id(list_id).tasks.post(new_task),
                timeout=API_TIMEOUT
            )
        else:
            result = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.by_todo_task_list_id(list_id).tasks.post(new_task),
                timeout=API_TIMEOUT
            )

        if not result:
            return "❌ Failed to create task"

        # Build confirmation message
        lines = [
            "✅ Task created successfully!",
            "",
            f"Title: {result.title}",
            f"Task ID: {result.id}",
            f"List ID: {list_id}",
        ]

        if due_date:
            lines.append(f"Due: {due_date}")

        if importance != "normal":
            lines.append(f"Importance: {importance.capitalize()}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task list not found\n\n"
                f"List ID: {list_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error creating task: {error_type}\n\n"
            f"Title: {title}\n"
            f"List ID: {list_id}\n"
            f"Mailbox: {mailbox_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def update_todo_task(
    list_id: str,
    task_id: str,
    mailbox_id: str = "me",
    title: str = "",
    body: str = "",
    due_date: str = "",
    importance: str = "",
    status: str = ""
) -> str:
    """
    Update an existing To Do task. Only provided fields are updated.

    Args:
        list_id: The ID of the task list
        task_id: The ID of the task to update
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox
        title: New task title (optional - leave empty to keep current)
        body: New task notes (optional - leave empty to keep current)
        due_date: New due date in YYYY-MM-DD format (optional)
        importance: New importance - "low", "normal", "high" (optional)
        status: New status - "notStarted", "inProgress", "completed" (optional)

    Returns:
        Confirmation with updated task details

    Examples:
        update_todo_task(list_id="AAMk...", task_id="AAMk...", title="Updated title")
        update_todo_task(list_id="AAMk...", task_id="AAMk...", status="inProgress")
    """
    try:
        client = await get_graph_client()

        # Build update object with only provided fields
        update_task = TodoTask()
        fields_updated = []

        if title:
            update_task.title = title
            fields_updated.append("title")

        if body:
            update_task.body = ItemBody(
                content=body,
                content_type=BodyType.Text
            )
            fields_updated.append("body")

        if due_date:
            update_task.due_date_time = DateTimeTimeZone(
                date_time=f"{due_date}T00:00:00",
                time_zone="UTC"
            )
            fields_updated.append("due_date")

        if importance:
            importance_map = {
                "low": Importance.Low,
                "normal": Importance.Normal,
                "high": Importance.High
            }
            update_task.importance = importance_map.get(importance.lower(), Importance.Normal)
            fields_updated.append("importance")

        if status:
            status_map = {
                "notstarted": TaskStatus.NotStarted,
                "inprogress": TaskStatus.InProgress,
                "completed": TaskStatus.Completed,
                "waitingonothers": TaskStatus.WaitingOnOthers,
                "deferred": TaskStatus.Deferred
            }
            update_task.status = status_map.get(status.lower(), TaskStatus.NotStarted)
            fields_updated.append("status")

        if not fields_updated:
            return "❌ No fields provided to update"

        if mailbox_id == "me":
            result = await asyncio.wait_for(
                client.me.todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).patch(update_task),
                timeout=API_TIMEOUT
            )
        else:
            result = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).patch(update_task),
                timeout=API_TIMEOUT
            )

        if not result:
            return "❌ Failed to update task"

        return (
            f"✅ Task updated successfully!\n\n"
            f"Title: {result.title}\n"
            f"Updated fields: {', '.join(fields_updated)}\n"
            f"Task ID: {result.id}"
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task not found\n\n"
                f"Task ID: {task_id}\n"
                f"List ID: {list_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error updating task: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"List ID: {list_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def complete_todo_task(
    list_id: str,
    task_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Mark a To Do task as completed.

    Args:
        list_id: The ID of the task list
        task_id: The ID of the task to complete
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox

    Returns:
        Confirmation of completion

    Examples:
        complete_todo_task(list_id="AAMk...", task_id="AAMk...", mailbox_id="me")
    """
    try:
        client = await get_graph_client()

        # Update task status to completed
        update_task = TodoTask(
            status=TaskStatus.Completed
        )

        if mailbox_id == "me":
            result = await asyncio.wait_for(
                client.me.todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).patch(update_task),
                timeout=API_TIMEOUT
            )
        else:
            result = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).patch(update_task),
                timeout=API_TIMEOUT
            )

        if not result:
            return "❌ Failed to complete task"

        return (
            f"✅ Task completed!\n\n"
            f"Title: {result.title}\n"
            f"Status: Completed\n"
            f"Task ID: {result.id}"
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task not found\n\n"
                f"Task ID: {task_id}\n"
                f"List ID: {list_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error completing task: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"List ID: {list_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def delete_todo_task(
    list_id: str,
    task_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Delete a To Do task permanently.

    Args:
        list_id: The ID of the task list
        task_id: The ID of the task to delete
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox

    Returns:
        Confirmation of deletion

    Examples:
        delete_todo_task(list_id="AAMk...", task_id="AAMk...", mailbox_id="me")
    """
    try:
        client = await get_graph_client()

        if mailbox_id == "me":
            await asyncio.wait_for(
                client.me.todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).delete(),
                timeout=API_TIMEOUT
            )
        else:
            await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).delete(),
                timeout=API_TIMEOUT
            )

        return (
            f"✅ Task deleted successfully!\n\n"
            f"Task ID: {task_id}\n"
            f"List ID: {list_id}"
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task not found\n\n"
                f"Task ID: {task_id}\n"
                f"List ID: {list_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error deleting task: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"List ID: {list_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def delete_todo_list(
    list_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Delete a To Do task list and all its tasks permanently.

    WARNING: This deletes the entire list and all tasks within it.

    Args:
        list_id: The ID of the task list to delete
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox

    Returns:
        Confirmation of deletion

    Examples:
        delete_todo_list(list_id="AAMk...", mailbox_id="me")
    """
    try:
        client = await get_graph_client()

        if mailbox_id == "me":
            await asyncio.wait_for(
                client.me.todo.lists.by_todo_task_list_id(list_id).delete(),
                timeout=API_TIMEOUT
            )
        else:
            await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.by_todo_task_list_id(list_id).delete(),
                timeout=API_TIMEOUT
            )

        return (
            f"✅ Task list deleted successfully!\n\n"
            f"List ID: {list_id}\n\n"
            f"All tasks in this list have been deleted."
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task list not found\n\n"
                f"List ID: {list_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        if "cannot delete" in error_msg.lower() or "default" in error_msg.lower():
            return (
                f"❌ Cannot delete this list\n\n"
                f"The default Tasks list cannot be deleted.\n"
                f"List ID: {list_id}"
            )

        return (
            f"❌ Error deleting task list: {error_type}\n\n"
            f"List ID: {list_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def add_checklist_item(
    list_id: str,
    task_id: str,
    display_name: str,
    mailbox_id: str = "me"
) -> str:
    """
    Add a checklist item (subtask) to a To Do task.

    Args:
        list_id: The ID of the task list
        task_id: The ID of the task to add checklist item to
        display_name: Text for the checklist item (required)
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox

    Returns:
        Confirmation with new checklist item ID

    Examples:
        add_checklist_item(list_id="AAMk...", task_id="AAMk...", display_name="Buy milk")
    """
    try:
        client = await get_graph_client()

        # Create checklist item
        new_item = ChecklistItem(
            display_name=display_name
        )

        if mailbox_id == "me":
            result = await asyncio.wait_for(
                client.me.todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).checklist_items.post(new_item),
                timeout=API_TIMEOUT
            )
        else:
            result = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).checklist_items.post(new_item),
                timeout=API_TIMEOUT
            )

        if not result:
            return "❌ Failed to add checklist item"

        return (
            f"✅ Checklist item added!\n\n"
            f"Item: {result.display_name}\n"
            f"Item ID: {result.id}\n"
            f"Task ID: {task_id}"
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task not found\n\n"
                f"Task ID: {task_id}\n"
                f"List ID: {list_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error adding checklist item: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"List ID: {list_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def toggle_checklist_item(
    list_id: str,
    task_id: str,
    item_id: str,
    is_checked: bool,
    mailbox_id: str = "me"
) -> str:
    """
    Toggle the checked status of a checklist item.

    Args:
        list_id: The ID of the task list
        task_id: The ID of the task containing the checklist item
        item_id: The ID of the checklist item (from get_todo_task)
        is_checked: True to mark as checked, False to uncheck
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox

    Returns:
        Confirmation of status change

    Examples:
        toggle_checklist_item(list_id="AAMk...", task_id="AAMk...", item_id="AAMk...", is_checked=True)
    """
    try:
        client = await get_graph_client()

        # Update checklist item
        update_item = ChecklistItem(
            is_checked=is_checked
        )

        if mailbox_id == "me":
            result = await asyncio.wait_for(
                client.me.todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).checklist_items.by_checklist_item_id(item_id).patch(update_item),
                timeout=API_TIMEOUT
            )
        else:
            result = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).checklist_items.by_checklist_item_id(item_id).patch(update_item),
                timeout=API_TIMEOUT
            )

        if not result:
            return "❌ Failed to update checklist item"

        status = "checked" if is_checked else "unchecked"
        icon = "✅" if is_checked else "⬜"

        return (
            f"{icon} Checklist item {status}!\n\n"
            f"Item: {result.display_name}\n"
            f"Item ID: {result.id}"
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Checklist item not found\n\n"
                f"Item ID: {item_id}\n"
                f"Task ID: {task_id}\n"
                f"List ID: {list_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error updating checklist item: {error_type}\n\n"
            f"Item ID: {item_id}\n"
            f"Task ID: {task_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def delete_checklist_item(
    list_id: str,
    task_id: str,
    item_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Delete a checklist item from a To Do task.

    Args:
        list_id: The ID of the task list
        task_id: The ID of the task containing the checklist item
        item_id: The ID of the checklist item to delete
        mailbox_id: Email address of mailbox to access (required)
                   Use "me" for athena@'s own mailbox

    Returns:
        Confirmation of deletion

    Examples:
        delete_checklist_item(list_id="AAMk...", task_id="AAMk...", item_id="AAMk...")
    """
    try:
        client = await get_graph_client()

        if mailbox_id == "me":
            await asyncio.wait_for(
                client.me.todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).checklist_items.by_checklist_item_id(item_id).delete(),
                timeout=API_TIMEOUT
            )
        else:
            await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).todo.lists.by_todo_task_list_id(list_id).tasks.by_todo_task_id(task_id).checklist_items.by_checklist_item_id(item_id).delete(),
                timeout=API_TIMEOUT
            )

        return (
            f"✅ Checklist item deleted!\n\n"
            f"Item ID: {item_id}\n"
            f"Task ID: {task_id}"
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Checklist item not found\n\n"
                f"Item ID: {item_id}\n"
                f"Task ID: {task_id}\n"
                f"List ID: {list_id}\n"
                f"Mailbox: {mailbox_id}"
            )

        return (
            f"❌ Error deleting checklist item: {error_type}\n\n"
            f"Item ID: {item_id}\n"
            f"Task ID: {task_id}\n"
            f"Error: {error_msg}"
        )


# ============================================================================
# PLANNER TOOLS (Phase 11 - Milestone 4 & 5)
# ============================================================================
# Microsoft Planner integration for team/project task management
# Requires Group.Read.All scope for accessing plans
# ============================================================================


@mcp.tool()
async def list_planner_plans(
    mailbox_id: str = "me"
) -> str:
    """
    List all Planner plans the user has access to.

    Returns plans from all groups the user is a member of.

    Args:
        mailbox_id: User to get plans for ("me" or email address)

    Returns:
        Formatted list of plans with ID, title, owner group
    """
    try:
        client = await get_graph_client()

        # Get plans assigned to the user
        if mailbox_id == "me":
            result = await asyncio.wait_for(
                client.me.planner.plans.get(),
                timeout=API_TIMEOUT
            )
        else:
            result = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).planner.plans.get(),
                timeout=API_TIMEOUT
            )

        if not result or not result.value:
            return (
                "📋 No Planner plans found\n\n"
                "You don't have access to any Planner plans, or the Group.Read.All "
                "permission may be required."
            )

        lines = [f"📋 **Planner Plans** ({len(result.value)} found)\n"]

        for plan in result.value:
            lines.append(f"---")
            lines.append(f"**{plan.title or 'Untitled Plan'}**")
            lines.append(f"  Plan ID: `{plan.id}`")
            if plan.owner:
                lines.append(f"  Owner Group: {plan.owner}")
            if plan.created_date_time:
                lines.append(f"  Created: {plan.created_date_time.strftime('%Y-%m-%d')}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                f"❌ Permission denied\n\n"
                f"The Group.Read.All permission is required to access Planner plans.\n"
                f"Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error listing Planner plans: {error_type}\n\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def list_planner_buckets(
    plan_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    List all buckets (columns) in a Planner plan.

    Args:
        plan_id: The Planner plan ID (from list_planner_plans)
        mailbox_id: User context ("me" or email address)

    Returns:
        Formatted list of buckets with ID, name, order
    """
    try:
        client = await get_graph_client()

        result = await asyncio.wait_for(
            client.planner.plans.by_planner_plan_id(plan_id).buckets.get(),
            timeout=API_TIMEOUT
        )

        if not result or not result.value:
            return (
                f"📦 No buckets found in plan\n\n"
                f"Plan ID: {plan_id}\n"
                f"The plan may be empty or you may not have access."
            )

        # Sort by order_hint for proper display order
        buckets = sorted(result.value, key=lambda b: b.order_hint or "")

        lines = [f"📦 **Plan Buckets** ({len(buckets)} found)\n"]
        lines.append(f"Plan ID: `{plan_id}`\n")

        for i, bucket in enumerate(buckets, 1):
            lines.append(f"---")
            lines.append(f"**{i}. {bucket.name or 'Unnamed Bucket'}**")
            lines.append(f"  Bucket ID: `{bucket.id}`")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Plan not found\n\n"
                f"Plan ID: {plan_id}\n"
                f"The plan may have been deleted or you don't have access."
            )

        return (
            f"❌ Error listing buckets: {error_type}\n\n"
            f"Plan ID: {plan_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def list_planner_tasks(
    plan_id: str,
    bucket_id: str = "",
    include_completed: bool = True,
    mailbox_id: str = "me"
) -> str:
    """
    List tasks in a Planner plan, optionally filtered by bucket.

    Args:
        plan_id: The Planner plan ID (from list_planner_plans)
        bucket_id: Optional bucket ID to filter tasks (from list_planner_buckets)
        include_completed: Include completed tasks (default: True)
        mailbox_id: User context ("me" or email address)

    Returns:
        Formatted list of tasks with details
    """
    try:
        client = await get_graph_client()

        if bucket_id:
            # Get tasks for specific bucket
            result = await asyncio.wait_for(
                client.planner.buckets.by_planner_bucket_id(bucket_id).tasks.get(),
                timeout=API_TIMEOUT
            )
        else:
            # Get all tasks in plan
            result = await asyncio.wait_for(
                client.planner.plans.by_planner_plan_id(plan_id).tasks.get(),
                timeout=API_TIMEOUT
            )

        if not result or not result.value:
            filter_msg = f" in bucket `{bucket_id}`" if bucket_id else ""
            return (
                f"📋 No tasks found{filter_msg}\n\n"
                f"Plan ID: {plan_id}"
            )

        tasks = result.value

        # Filter completed if requested
        if not include_completed:
            tasks = [t for t in tasks if t.percent_complete != 100]

        if not tasks:
            return (
                f"📋 No incomplete tasks found\n\n"
                f"Plan ID: {plan_id}\n"
                f"All tasks may be completed."
            )

        # Sort by priority then due date
        def sort_key(t):
            priority = t.priority if t.priority is not None else 5
            due = t.due_date_time.isoformat() if t.due_date_time else "9999-12-31"
            return (priority, due)

        tasks = sorted(tasks, key=sort_key)

        lines = [f"📋 **Planner Tasks** ({len(tasks)} found)\n"]
        lines.append(f"Plan ID: `{plan_id}`")
        if bucket_id:
            lines.append(f"Bucket: `{bucket_id}`")
        lines.append("")

        priority_map = {1: "🔴 Urgent", 3: "🟠 Important", 5: "🟡 Medium", 9: "🟢 Low"}

        for task in tasks:
            lines.append("---")
            # Status indicator
            if task.percent_complete == 100:
                status = "✅"
            elif task.percent_complete and task.percent_complete > 0:
                status = "🔄"
            else:
                status = "⬜"

            lines.append(f"{status} **{task.title or 'Untitled Task'}**")
            lines.append(f"  Task ID: `{task.id}`")

            if task.bucket_id:
                lines.append(f"  Bucket: `{task.bucket_id}`")

            if task.priority is not None:
                priority_label = priority_map.get(task.priority, f"Priority {task.priority}")
                lines.append(f"  Priority: {priority_label}")

            if task.due_date_time:
                lines.append(f"  Due: {task.due_date_time.strftime('%Y-%m-%d')}")

            if task.percent_complete is not None and task.percent_complete > 0:
                lines.append(f"  Progress: {task.percent_complete}%")

            # Check for assignments
            if task.assignments and task.assignments.additional_data:
                assignee_count = len(task.assignments.additional_data)
                lines.append(f"  Assignees: {assignee_count}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Plan or bucket not found\n\n"
                f"Plan ID: {plan_id}\n"
                f"Bucket ID: {bucket_id if bucket_id else 'None'}"
            )

        return (
            f"❌ Error listing tasks: {error_type}\n\n"
            f"Plan ID: {plan_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_planner_task(
    task_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Get full details of a specific Planner task.

    Includes task details, checklist items, and attachments.

    Args:
        task_id: The task ID (from list_planner_tasks)
        mailbox_id: User context ("me" or email address)

    Returns:
        Complete task details with description, checklist, and references
    """
    try:
        client = await get_graph_client()

        # Get task basic info
        task = await asyncio.wait_for(
            client.planner.tasks.by_planner_task_id(task_id).get(),
            timeout=API_TIMEOUT
        )

        if not task:
            return (
                f"❌ Task not found\n\n"
                f"Task ID: {task_id}"
            )

        # Get task details (description, checklist, references)
        try:
            details = await asyncio.wait_for(
                client.planner.tasks.by_planner_task_id(task_id).details.get(),
                timeout=API_TIMEOUT
            )
        except Exception:
            details = None

        lines = ["📋 **Planner Task Details**\n"]

        # Status indicator
        if task.percent_complete == 100:
            status = "✅ Completed"
        elif task.percent_complete and task.percent_complete > 0:
            status = f"🔄 In Progress ({task.percent_complete}%)"
        else:
            status = "⬜ Not Started"

        lines.append(f"**{task.title or 'Untitled Task'}**")
        lines.append(f"Status: {status}")
        lines.append(f"Task ID: `{task.id}`")
        lines.append(f"Plan ID: `{task.plan_id}`")

        if task.bucket_id:
            lines.append(f"Bucket ID: `{task.bucket_id}`")

        # Priority
        priority_map = {1: "🔴 Urgent", 3: "🟠 Important", 5: "🟡 Medium", 9: "🟢 Low"}
        if task.priority is not None:
            priority_label = priority_map.get(task.priority, f"Priority {task.priority}")
            lines.append(f"Priority: {priority_label}")

        # Dates
        if task.start_date_time:
            lines.append(f"Start: {task.start_date_time.strftime('%Y-%m-%d')}")
        if task.due_date_time:
            lines.append(f"Due: {task.due_date_time.strftime('%Y-%m-%d')}")
        if task.completed_date_time:
            lines.append(f"Completed: {task.completed_date_time.strftime('%Y-%m-%d %H:%M')}")

        # Assignments
        if task.assignments and task.assignments.additional_data:
            lines.append("\n**Assignees:**")
            for user_id in task.assignments.additional_data.keys():
                lines.append(f"  - User: `{user_id}`")

        # Description from details
        if details and details.description:
            lines.append("\n**Description:**")
            lines.append(details.description[:500])
            if len(details.description) > 500:
                lines.append("... (truncated)")

        # Checklist from details
        if details and details.checklist and details.checklist.additional_data:
            checklist_items = details.checklist.additional_data
            lines.append(f"\n**Checklist ({len(checklist_items)} items):**")
            for item_id, item_data in checklist_items.items():
                if isinstance(item_data, dict):
                    is_checked = item_data.get("isChecked", False)
                    title = item_data.get("title", "Untitled")
                    check_mark = "☑️" if is_checked else "⬜"
                    lines.append(f"  {check_mark} {title}")

        # References (attachments) from details
        if details and details.references and details.references.additional_data:
            refs = details.references.additional_data
            lines.append(f"\n**References ({len(refs)} items):**")
            for ref_url, ref_data in refs.items():
                if isinstance(ref_data, dict):
                    alias = ref_data.get("alias", ref_url)
                    lines.append(f"  - {alias}")

        # Metadata
        lines.append("\n**Metadata:**")
        if task.created_date_time:
            lines.append(f"  Created: {task.created_date_time.strftime('%Y-%m-%d %H:%M')}")
        if task.created_by:
            lines.append(f"  Created by: {task.created_by.user.id if task.created_by.user else 'Unknown'}")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task not found\n\n"
                f"Task ID: {task_id}\n"
                f"The task may have been deleted or you don't have access."
            )

        return (
            f"❌ Error getting task: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def create_planner_task(
    plan_id: str,
    title: str,
    bucket_id: str = "",
    due_date: str = "",
    priority: int = 5,
    notes: str = "",
    mailbox_id: str = "me"
) -> str:
    """
    Create a new task in a Planner plan.

    Args:
        plan_id: The Planner plan ID (from list_planner_plans) - REQUIRED
        title: Task title - REQUIRED
        bucket_id: Optional bucket ID to place task in (from list_planner_buckets)
        due_date: Optional due date in YYYY-MM-DD format
        priority: Task priority (1=Urgent, 3=Important, 5=Medium, 9=Low, default: 5)
        notes: Optional task description/notes
        mailbox_id: User context ("me" or email address)

    Returns:
        Created task details with ID
    """
    try:
        client = await get_graph_client()

        # Build task object
        new_task = PlannerTask()
        new_task.plan_id = plan_id
        new_task.title = title

        if bucket_id:
            new_task.bucket_id = bucket_id

        if due_date:
            try:
                from datetime import datetime, timezone
                due_dt = datetime.strptime(due_date, "%Y-%m-%d")
                # Set to end of day in UTC
                due_dt = due_dt.replace(hour=23, minute=59, second=59, tzinfo=timezone.utc)
                new_task.due_date_time = due_dt
            except ValueError:
                return (
                    f"❌ Invalid due date format\n\n"
                    f"Expected: YYYY-MM-DD\n"
                    f"Got: {due_date}"
                )

        # Priority: 1=Urgent, 3=Important, 5=Medium, 9=Low
        if priority in [1, 3, 5, 9]:
            new_task.priority = priority
        else:
            new_task.priority = 5

        # Create the task
        result = await asyncio.wait_for(
            client.planner.tasks.post(new_task),
            timeout=API_TIMEOUT
        )

        if not result:
            return "❌ Failed to create task - no result returned"

        # If notes provided, update task details
        if notes and result.id:
            try:
                # Get current details for etag
                details = await asyncio.wait_for(
                    client.planner.tasks.by_planner_task_id(result.id).details.get(),
                    timeout=API_TIMEOUT
                )
                if details:
                    # Update with notes
                    details.description = notes
                    # Get etag for If-Match header (required by Planner API)
                    details_etag = details.additional_data.get("@odata.etag", "")
                    details_config = RequestConfiguration()
                    details_config.headers.try_add("If-Match", details_etag)

                    await asyncio.wait_for(
                        client.planner.tasks.by_planner_task_id(result.id).details.patch(details, request_configuration=details_config),
                        timeout=API_TIMEOUT
                    )
            except Exception:
                # Notes update failed but task was created
                pass

        priority_map = {1: "Urgent", 3: "Important", 5: "Medium", 9: "Low"}
        priority_label = priority_map.get(result.priority or 5, "Medium")

        lines = [
            f"✅ Planner task created!\n",
            f"**{result.title}**",
            f"Task ID: `{result.id}`",
            f"Plan ID: `{result.plan_id}`",
        ]

        if result.bucket_id:
            lines.append(f"Bucket ID: `{result.bucket_id}`")

        lines.append(f"Priority: {priority_label}")

        if result.due_date_time:
            lines.append(f"Due: {result.due_date_time.strftime('%Y-%m-%d')}")

        if notes:
            lines.append(f"Notes: Added")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                f"❌ Permission denied\n\n"
                f"The Group.Read.All permission is required to create Planner tasks.\n"
                f"Please add this permission in Azure Portal and re-authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Plan or bucket not found\n\n"
                f"Plan ID: {plan_id}\n"
                f"Bucket ID: {bucket_id if bucket_id else 'None'}"
            )

        return (
            f"❌ Error creating task: {error_type}\n\n"
            f"Title: {title}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def update_planner_task(
    task_id: str,
    title: str = "",
    bucket_id: str = "",
    due_date: str = "",
    priority: int = 0,
    percent_complete: int = -1,
    notes: str = "",
    mailbox_id: str = "me"
) -> str:
    """
    Update an existing Planner task.

    Only provided parameters are updated - others remain unchanged.

    Args:
        task_id: The task ID (from list_planner_tasks) - REQUIRED
        title: New task title (optional)
        bucket_id: Move to different bucket (optional)
        due_date: New due date in YYYY-MM-DD format (optional)
        priority: New priority 1=Urgent, 3=Important, 5=Medium, 9=Low (optional)
        percent_complete: Progress 0-100 (optional, use complete_planner_task for 100%)
        notes: Update description/notes (optional)
        mailbox_id: User context ("me" or email address)

    Returns:
        Updated task details
    """
    try:
        # Load token for direct API call (SDK has etag header issues)
        token_data = load_token_cache()
        if not token_data or not is_token_valid(token_data):
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        access_token = token_data['access_token']

        # Build update payload
        update_payload = {}
        updates = []

        if title:
            update_payload["title"] = title
            updates.append(f"Title: {title}")

        if bucket_id:
            update_payload["bucketId"] = bucket_id
            updates.append(f"Bucket: {bucket_id}")

        if due_date:
            try:
                from datetime import datetime, timezone
                due_dt = datetime.strptime(due_date, "%Y-%m-%d")
                due_dt = due_dt.replace(hour=23, minute=59, second=59, tzinfo=timezone.utc)
                update_payload["dueDateTime"] = due_dt.isoformat()
                updates.append(f"Due: {due_date}")
            except ValueError:
                return (
                    f"❌ Invalid due date format\n\n"
                    f"Expected: YYYY-MM-DD\n"
                    f"Got: {due_date}"
                )

        if priority in [1, 3, 5, 9]:
            update_payload["priority"] = priority
            priority_map = {1: "Urgent", 3: "Important", 5: "Medium", 9: "Low"}
            updates.append(f"Priority: {priority_map.get(priority)}")

        if percent_complete >= 0 and percent_complete <= 100:
            update_payload["percentComplete"] = percent_complete
            updates.append(f"Progress: {percent_complete}%")

        if not update_payload and not notes:
            return (
                f"❌ No updates specified\n\n"
                f"Provide at least one field to update."
            )

        async with httpx.AsyncClient() as http_client:
            # Update task fields if we have any
            if update_payload:
                # Get current task to get etag
                get_resp = await http_client.get(
                    f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}",
                    headers={"Authorization": f"Bearer {access_token}"},
                    timeout=API_TIMEOUT
                )

                if get_resp.status_code == 404:
                    return f"❌ Task not found\n\nTask ID: {task_id}"

                if get_resp.status_code != 200:
                    return f"❌ Error fetching task\n\nStatus: {get_resp.status_code}\n{get_resp.text}"

                task_data = get_resp.json()
                etag = task_data.get("@odata.etag", "")

                # Patch the task with If-Match header
                patch_resp = await http_client.patch(
                    f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}",
                    headers={
                        "Authorization": f"Bearer {access_token}",
                        "Content-Type": "application/json",
                        "If-Match": etag
                    },
                    json=update_payload,
                    timeout=API_TIMEOUT
                )

                if patch_resp.status_code == 412:
                    return (
                        f"❌ Task was modified by someone else\n\n"
                        f"Task ID: {task_id}\n"
                        f"Please try again."
                    )

                if patch_resp.status_code not in [200, 204]:
                    return f"❌ Error updating task\n\nStatus: {patch_resp.status_code}\n{patch_resp.text}"

            # Update notes if provided (requires separate endpoint)
            if notes:
                try:
                    # Get task details to get its etag
                    details_resp = await http_client.get(
                        f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}/details",
                        headers={"Authorization": f"Bearer {access_token}"},
                        timeout=API_TIMEOUT
                    )

                    if details_resp.status_code == 200:
                        details_data = details_resp.json()
                        details_etag = details_data.get("@odata.etag", "")

                        # Patch the details
                        details_patch_resp = await http_client.patch(
                            f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}/details",
                            headers={
                                "Authorization": f"Bearer {access_token}",
                                "Content-Type": "application/json",
                                "If-Match": details_etag
                            },
                            json={"description": notes},
                            timeout=API_TIMEOUT
                        )

                        if details_patch_resp.status_code in [200, 204]:
                            updates.append("Notes: Updated")
                        else:
                            updates.append(f"Notes: Failed ({details_patch_resp.status_code})")
                    else:
                        updates.append("Notes: Failed to fetch details")
                except Exception:
                    updates.append("Notes: Failed to update")

        return (
            f"✅ Planner task updated!\n\n"
            f"Task ID: `{task_id}`\n\n"
            f"**Updates:**\n" + "\n".join(f"  - {u}" for u in updates)
        )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error updating task: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def assign_planner_task(
    task_id: str,
    assignee_id: str,
    remove: bool = False,
    mailbox_id: str = "me"
) -> str:
    """
    Assign or unassign a user to/from a Planner task.

    Args:
        task_id: The task ID (from list_planner_tasks) - REQUIRED
        assignee_id: The user's Azure AD object ID to assign - REQUIRED
        remove: If True, removes the user from the task instead of adding
        mailbox_id: User context ("me" or email address)

    Returns:
        Confirmation of assignment change
    """
    try:
        # Load token for direct API call (SDK has etag header issues)
        token_data = load_token_cache()
        if not token_data or not is_token_valid(token_data):
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        access_token = token_data['access_token']

        async with httpx.AsyncClient() as http_client:
            # Get current task to get etag and current assignments
            get_resp = await http_client.get(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}",
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=API_TIMEOUT
            )

            if get_resp.status_code == 404:
                return f"❌ Task not found\n\nTask ID: {task_id}"

            if get_resp.status_code != 200:
                return f"❌ Error fetching task\n\nStatus: {get_resp.status_code}\n{get_resp.text}"

            task_data = get_resp.json()
            etag = task_data.get("@odata.etag", "")

            # Get current assignments
            current_assignments = task_data.get("assignments", {})

            if remove:
                if assignee_id in current_assignments:
                    # Set to null to remove assignment
                    current_assignments[assignee_id] = None
                    action = "removed from"
                else:
                    return (
                        f"❌ User not assigned to this task\n\n"
                        f"Task ID: {task_id}\n"
                        f"User ID: {assignee_id}"
                    )
            else:
                if assignee_id in current_assignments:
                    return (
                        f"❌ User already assigned to this task\n\n"
                        f"Task ID: {task_id}\n"
                        f"User ID: {assignee_id}"
                    )
                # Add assignment with order hint
                current_assignments[assignee_id] = {
                    "@odata.type": "#microsoft.graph.plannerAssignment",
                    "orderHint": " !"
                }
                action = "assigned to"

            # Patch the task with If-Match header
            patch_resp = await http_client.patch(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}",
                headers={
                    "Authorization": f"Bearer {access_token}",
                    "Content-Type": "application/json",
                    "If-Match": etag
                },
                json={"assignments": current_assignments},
                timeout=API_TIMEOUT
            )

            if patch_resp.status_code == 412:
                return (
                    f"❌ Task was modified by someone else\n\n"
                    f"Task ID: {task_id}\n"
                    f"Please try again."
                )

            if patch_resp.status_code not in [200, 204]:
                return f"❌ Error updating assignment\n\nStatus: {patch_resp.status_code}\n{patch_resp.text}"

            return (
                f"✅ User {action} task!\n\n"
                f"Task ID: `{task_id}`\n"
                f"User ID: `{assignee_id}`\n"
                f"Action: {'Removed' if remove else 'Assigned'}"
            )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task or user not found\n\n"
                f"Task ID: {task_id}\n"
                f"User ID: {assignee_id}"
            )

        return (
            f"❌ Error updating assignment: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def complete_planner_task(
    task_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Mark a Planner task as complete (100% progress).

    Args:
        task_id: The task ID (from list_planner_tasks) - REQUIRED
        mailbox_id: User context ("me" or email address)

    Returns:
        Confirmation of completion
    """
    try:
        # Load token for direct API call (SDK has etag header issues)
        token_data = load_token_cache()
        if not token_data or not is_token_valid(token_data):
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        access_token = token_data['access_token']

        async with httpx.AsyncClient() as http_client:
            # Get current task to check status and get etag
            get_resp = await http_client.get(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}",
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=API_TIMEOUT
            )

            if get_resp.status_code == 404:
                return f"❌ Task not found\n\nTask ID: {task_id}"

            if get_resp.status_code != 200:
                return f"❌ Error fetching task\n\nStatus: {get_resp.status_code}\n{get_resp.text}"

            task_data = get_resp.json()
            task_title = task_data.get("title", "Untitled")
            etag = task_data.get("@odata.etag", "")
            current_percent = task_data.get("percentComplete", 0)

            if current_percent == 100:
                return (
                    f"ℹ️ Task is already complete\n\n"
                    f"**{task_title}**\n"
                    f"Task ID: `{task_id}`"
                )

            # Update to 100% with If-Match header
            patch_resp = await http_client.patch(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}",
                headers={
                    "Authorization": f"Bearer {access_token}",
                    "Content-Type": "application/json",
                    "If-Match": etag
                },
                json={"percentComplete": 100},
                timeout=API_TIMEOUT
            )

            if patch_resp.status_code in [200, 204]:
                return (
                    f"✅ Task completed!\n\n"
                    f"**{task_title}**\n"
                    f"Task ID: `{task_id}`\n"
                    f"Progress: 100%"
                )
            elif patch_resp.status_code == 412:
                return (
                    f"❌ Task was modified by someone else\n\n"
                    f"Task ID: {task_id}\n"
                    f"Please try again."
                )
            else:
                return f"❌ Error completing task\n\nStatus: {patch_resp.status_code}\n{patch_resp.text}"

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error completing task: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def delete_planner_task(
    task_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Delete a Planner task permanently.

    WARNING: This action cannot be undone.

    Args:
        task_id: The task ID (from list_planner_tasks) - REQUIRED
        mailbox_id: User context ("me" or email address)

    Returns:
        Confirmation of deletion
    """
    try:
        # Load token for direct API call (SDK has etag header issues)
        token_data = load_token_cache()
        if not token_data or not is_token_valid(token_data):
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        access_token = token_data['access_token']

        # Get task first to confirm it exists and get title/etag
        async with httpx.AsyncClient() as http_client:
            get_resp = await http_client.get(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}",
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=API_TIMEOUT
            )

            if get_resp.status_code == 404:
                return f"❌ Task not found\n\nTask ID: {task_id}"

            if get_resp.status_code != 200:
                return f"❌ Error fetching task\n\nStatus: {get_resp.status_code}\n{get_resp.text}"

            task_data = get_resp.json()
            task_title = task_data.get("title", "Untitled")
            etag = task_data.get("@odata.etag", "")

            # Delete the task with If-Match header
            delete_resp = await http_client.delete(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}",
                headers={
                    "Authorization": f"Bearer {access_token}",
                    "If-Match": etag
                },
                timeout=API_TIMEOUT
            )

            if delete_resp.status_code == 204:
                return (
                    f"✅ Planner task deleted!\n\n"
                    f"**{task_title}**\n"
                    f"Task ID: `{task_id}`\n\n"
                    f"⚠️ This action cannot be undone."
                )
            elif delete_resp.status_code == 409:
                return (
                    f"❌ Task was modified by someone else\n\n"
                    f"Task ID: {task_id}\n"
                    f"Please try again."
                )
            else:
                return f"❌ Error deleting task\n\nStatus: {delete_resp.status_code}\n{delete_resp.text}"

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Task not found\n\n"
                f"Task ID: {task_id}\n"
                f"The task may have already been deleted."
            )

        return (
            f"❌ Error deleting task: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def add_planner_checklist_item(
    task_id: str,
    title: str,
    is_checked: bool = False,
    mailbox_id: str = "me"
) -> str:
    """
    Add a checklist item to a Planner task.

    Planner tasks can have checklists (subtasks) that can be checked off.
    This function adds a new item to an existing task's checklist.

    Args:
        task_id: The task ID (from list_planner_tasks) - REQUIRED
        title: The checklist item text - REQUIRED
        is_checked: Whether the item starts as checked (default: False)
        mailbox_id: User context ("me" or email address)

    Returns:
        Confirmation with the new checklist item details
    """
    import uuid

    try:
        token_data = load_token_cache()
        if not token_data or not is_token_valid(token_data):
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        access_token = token_data['access_token']

        async with httpx.AsyncClient() as http_client:
            # Get task details (contains checklist and etag)
            details_resp = await http_client.get(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}/details",
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=API_TIMEOUT
            )

            if details_resp.status_code == 404:
                return f"❌ Task not found\n\nTask ID: {task_id}"

            if details_resp.status_code != 200:
                return f"❌ Error fetching task details\n\nStatus: {details_resp.status_code}\n{details_resp.text}"

            details_data = details_resp.json()
            details_etag = details_data.get("@odata.etag", "")

            # Get existing checklist (may be None or empty dict)
            # Note: Planner API requires @odata.type annotation for all checklist items
            existing_checklist = {}
            raw_checklist = details_data.get("checklist", {}) or {}
            if hasattr(raw_checklist, "additional_data"):
                raw_checklist = raw_checklist.additional_data or {}

            for k, v in raw_checklist.items():
                if k.startswith("@"):
                    continue
                # Ensure each item has the required @odata.type annotation
                existing_checklist[k] = {
                    "@odata.type": "#microsoft.graph.plannerChecklistItem",
                    "isChecked": v.get("isChecked", False),
                    "title": v.get("title", "")
                }

            # Generate new item ID and add to checklist
            new_item_id = str(uuid.uuid4())
            existing_checklist[new_item_id] = {
                "@odata.type": "#microsoft.graph.plannerChecklistItem",
                "isChecked": is_checked,
                "title": title
            }

            # Patch task details with updated checklist
            patch_resp = await http_client.patch(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}/details",
                headers={
                    "Authorization": f"Bearer {access_token}",
                    "Content-Type": "application/json",
                    "If-Match": details_etag
                },
                json={"checklist": existing_checklist},
                timeout=API_TIMEOUT
            )

            if patch_resp.status_code == 412:
                return (
                    f"❌ Task was modified by someone else\n\n"
                    f"Task ID: {task_id}\n"
                    f"Please try again."
                )

            if patch_resp.status_code not in [200, 204]:
                return f"❌ Error updating checklist\n\nStatus: {patch_resp.status_code}\n{patch_resp.text}"

            status = "☑️" if is_checked else "☐"
            return (
                f"✅ Checklist item added!\n\n"
                f"Task ID: `{task_id}`\n"
                f"Item: {status} {title}\n"
                f"Item ID: `{new_item_id}`\n\n"
                f"Total checklist items: {len(existing_checklist)}"
            )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error adding checklist item: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def update_planner_checklist_item(
    task_id: str,
    item_id: str,
    is_checked: bool = None,
    title: str = "",
    mailbox_id: str = "me"
) -> str:
    """
    Update or toggle a checklist item in a Planner task.

    Use this to check/uncheck items or update their text.

    Args:
        task_id: The task ID (from list_planner_tasks) - REQUIRED
        item_id: The checklist item ID (from get_planner_task) - REQUIRED
        is_checked: Set checked state (True/False), omit to toggle
        title: New title text (optional, keeps existing if not provided)
        mailbox_id: User context ("me" or email address)

    Returns:
        Updated checklist item details
    """
    try:
        token_data = load_token_cache()
        if not token_data or not is_token_valid(token_data):
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        access_token = token_data['access_token']

        async with httpx.AsyncClient() as http_client:
            # Get task details
            details_resp = await http_client.get(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}/details",
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=API_TIMEOUT
            )

            if details_resp.status_code == 404:
                return f"❌ Task not found\n\nTask ID: {task_id}"

            if details_resp.status_code != 200:
                return f"❌ Error fetching task details\n\nStatus: {details_resp.status_code}\n{details_resp.text}"

            details_data = details_resp.json()
            details_etag = details_data.get("@odata.etag", "")

            # Get existing checklist with @odata.type annotations
            existing_checklist = {}
            raw_checklist = details_data.get("checklist", {}) or {}
            if hasattr(raw_checklist, "additional_data"):
                raw_checklist = raw_checklist.additional_data or {}

            for k, v in raw_checklist.items():
                if k.startswith("@"):
                    continue
                existing_checklist[k] = {
                    "@odata.type": "#microsoft.graph.plannerChecklistItem",
                    "isChecked": v.get("isChecked", False),
                    "title": v.get("title", "")
                }

            # Find the item
            if item_id not in existing_checklist:
                available_ids = list(existing_checklist.keys())[:5]
                return (
                    f"❌ Checklist item not found\n\n"
                    f"Item ID: {item_id}\n"
                    f"Available IDs: {available_ids}"
                )

            current_item = existing_checklist[item_id]
            current_checked = current_item.get("isChecked", False)
            current_title = current_item.get("title", "")

            # Determine new values
            if is_checked is None:
                # Toggle mode
                new_checked = not current_checked
            else:
                new_checked = is_checked

            new_title = title if title else current_title

            # Update the item with @odata.type annotation
            existing_checklist[item_id] = {
                "@odata.type": "#microsoft.graph.plannerChecklistItem",
                "isChecked": new_checked,
                "title": new_title
            }

            # Patch task details
            patch_resp = await http_client.patch(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}/details",
                headers={
                    "Authorization": f"Bearer {access_token}",
                    "Content-Type": "application/json",
                    "If-Match": details_etag
                },
                json={"checklist": existing_checklist},
                timeout=API_TIMEOUT
            )

            if patch_resp.status_code == 412:
                return (
                    f"❌ Task was modified by someone else\n\n"
                    f"Task ID: {task_id}\n"
                    f"Please try again."
                )

            if patch_resp.status_code not in [200, 204]:
                return f"❌ Error updating checklist\n\nStatus: {patch_resp.status_code}\n{patch_resp.text}"

            old_status = "☑️" if current_checked else "☐"
            new_status = "☑️" if new_checked else "☐"

            return (
                f"✅ Checklist item updated!\n\n"
                f"Task ID: `{task_id}`\n"
                f"Item ID: `{item_id}`\n\n"
                f"Before: {old_status} {current_title}\n"
                f"After: {new_status} {new_title}"
            )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error updating checklist item: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def delete_planner_checklist_item(
    task_id: str,
    item_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Delete a checklist item from a Planner task.

    Args:
        task_id: The task ID (from list_planner_tasks) - REQUIRED
        item_id: The checklist item ID (from get_planner_task) - REQUIRED
        mailbox_id: User context ("me" or email address)

    Returns:
        Confirmation of deletion
    """
    try:
        token_data = load_token_cache()
        if not token_data or not is_token_valid(token_data):
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        access_token = token_data['access_token']

        async with httpx.AsyncClient() as http_client:
            # Get task details
            details_resp = await http_client.get(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}/details",
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=API_TIMEOUT
            )

            if details_resp.status_code == 404:
                return f"❌ Task not found\n\nTask ID: {task_id}"

            if details_resp.status_code != 200:
                return f"❌ Error fetching task details\n\nStatus: {details_resp.status_code}\n{details_resp.text}"

            details_data = details_resp.json()
            details_etag = details_data.get("@odata.etag", "")

            # Get existing checklist with @odata.type annotations
            existing_checklist = {}
            raw_checklist = details_data.get("checklist", {}) or {}
            if hasattr(raw_checklist, "additional_data"):
                raw_checklist = raw_checklist.additional_data or {}

            for k, v in raw_checklist.items():
                if k.startswith("@"):
                    continue
                existing_checklist[k] = {
                    "@odata.type": "#microsoft.graph.plannerChecklistItem",
                    "isChecked": v.get("isChecked", False),
                    "title": v.get("title", "")
                }

            # Find and remove the item
            if item_id not in existing_checklist:
                return (
                    f"❌ Checklist item not found\n\n"
                    f"Item ID: {item_id}"
                )

            deleted_item = existing_checklist.pop(item_id)
            deleted_title = deleted_item.get("title", "Untitled")

            # Patch task details with item removed
            patch_resp = await http_client.patch(
                f"https://graph.microsoft.com/v1.0/planner/tasks/{task_id}/details",
                headers={
                    "Authorization": f"Bearer {access_token}",
                    "Content-Type": "application/json",
                    "If-Match": details_etag
                },
                json={"checklist": existing_checklist},
                timeout=API_TIMEOUT
            )

            if patch_resp.status_code == 412:
                return (
                    f"❌ Task was modified by someone else\n\n"
                    f"Task ID: {task_id}\n"
                    f"Please try again."
                )

            if patch_resp.status_code not in [200, 204]:
                return f"❌ Error updating checklist\n\nStatus: {patch_resp.status_code}\n{patch_resp.text}"

            return (
                f"✅ Checklist item deleted!\n\n"
                f"Task ID: `{task_id}`\n"
                f"Deleted: {deleted_title}\n\n"
                f"Remaining checklist items: {len(existing_checklist)}"
            )

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        return (
            f"❌ Error deleting checklist item: {error_type}\n\n"
            f"Task ID: {task_id}\n"
            f"Error: {error_msg}"
        )


# ============================================================================
# M365 GROUPS (Phase 11 - Groups)
# ============================================================================
# Group listing for Planner plan creation
# Requires Group.Read.All scope
# ============================================================================


@mcp.tool()
async def list_groups(
    mailbox_id: str = "me"
) -> str:
    """
    List M365 Groups the user is a member of.

    Returns groups that can be used for creating Planner plans.
    Only M365 groups (unified groups) support Planner plans.

    Args:
        mailbox_id: User to get groups for ("me" or email address)

    Returns:
        Formatted list of groups with ID, name, description

    Examples:
        list_groups(mailbox_id="me")
        list_groups(mailbox_id="thomas@sixpillar.co.uk")
    """
    try:
        client = await get_graph_client()

        # Get groups user is member of
        # /me/memberOf returns DirectoryObjects, filter to groups only
        if mailbox_id == "me":
            result = await asyncio.wait_for(
                client.me.member_of.get(),
                timeout=API_TIMEOUT
            )
        else:
            result = await asyncio.wait_for(
                client.users.by_user_id(mailbox_id).member_of.get(),
                timeout=API_TIMEOUT
            )

        if not result or not result.value:
            return (
                "👥 No groups found\n\n"
                "You are not a member of any M365 groups, or the Group.Read.All "
                "permission may be required."
            )

        # Filter to only M365 groups (unified groups) - these support Planner
        # DirectoryObjects have odata_type to identify the actual type
        groups = []
        for obj in result.value:
            # Check if it's a group (not a role or other directory object)
            odata_type = getattr(obj, 'odata_type', None)
            if odata_type and 'group' in odata_type.lower():
                groups.append(obj)

        if not groups:
            return (
                "👥 No M365 groups found\n\n"
                "You are a member of directory objects but no M365 groups.\n"
                "Only M365 groups (unified groups) can have Planner plans."
            )

        lines = [f"👥 **M365 Groups** ({len(groups)} found)\n"]

        for group in groups:
            lines.append(f"---")
            lines.append(f"**{group.display_name or 'Unnamed Group'}**")
            lines.append(f"  Group ID: `{group.id}`")
            if hasattr(group, 'description') and group.description:
                # Truncate long descriptions
                desc = group.description[:100] + "..." if len(group.description) > 100 else group.description
                lines.append(f"  Description: {desc}")
            if hasattr(group, 'mail') and group.mail:
                lines.append(f"  Email: {group.mail}")
            if hasattr(group, 'visibility') and group.visibility:
                lines.append(f"  Visibility: {group.visibility}")

        lines.append("")
        lines.append("---")
        lines.append("Use the Group ID with `create_planner_plan()` to create a new plan.")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                f"❌ Permission denied\n\n"
                f"The Group.Read.All permission is required to list groups.\n"
                f"Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error listing groups: {error_type}\n\n"
            f"Error: {error_msg}"
        )


# ============================================================================
# M365 GROUPS - Extended Group Management (Phase 11b)
# ============================================================================
# Tenant-wide group visibility and management
# Requires Group.Read.All (read) or Group.ReadWrite.All (write)
# ============================================================================

# Admin role group patterns to exclude from management
ADMIN_GROUP_PATTERNS = [
    "global admin",
    "exchange admin",
    "sharepoint admin",
    "teams admin",
    "security admin",
    "compliance admin",
    "helpdesk admin",
    "user admin",
    "billing admin",
    "dynamics 365 admin",
    "power platform admin",
    "azure ad admin",
    "intune admin",
    "license admin",
    "password admin",
    "privileged",
    "role-",
]


def is_admin_group(group) -> bool:
    """Check if a group appears to be an admin/role group that should be excluded."""
    display_name = (getattr(group, 'display_name', None) or '').lower()

    # Check display name against admin patterns
    for pattern in ADMIN_GROUP_PATTERNS:
        if pattern in display_name:
            return True

    # Check if it's an Azure AD role group (has roleTemplateId or specific odata type)
    odata_type = getattr(group, 'odata_type', None) or ''
    if 'directoryRole' in odata_type:
        return True

    return False


def get_group_type(group) -> str:
    """Determine the type of group (M365, Security, Distribution, Dynamic)."""
    group_types = getattr(group, 'group_types', None) or []
    mail_enabled = getattr(group, 'mail_enabled', False)
    security_enabled = getattr(group, 'security_enabled', False)

    # Check for dynamic group
    if 'DynamicMembership' in group_types:
        return 'Dynamic'

    # Check for M365 group (Unified)
    if 'Unified' in group_types:
        return 'M365'

    # Mail-enabled security group
    if mail_enabled and security_enabled:
        return 'Mail-enabled Security'

    # Distribution group (mail only)
    if mail_enabled and not security_enabled:
        return 'Distribution'

    # Security group
    if security_enabled:
        return 'Security'

    return 'Unknown'


@mcp.tool()
async def list_all_groups(
    group_type: str = "",
    include_dynamic: bool = False,
    count: int = 50,
    mailbox_id: str = "me"
) -> str:
    """
    List all groups in the tenant (not just user memberships).

    Returns M365 groups, security groups, and distribution groups.
    Admin role groups are excluded for safety.

    Args:
        group_type: Filter by type - "M365", "Security", "Distribution", or "" for all
        include_dynamic: Include dynamic membership groups (default: False)
        count: Maximum groups to return (default: 50, max: 100)
        mailbox_id: User context for API call ("me" or email address)

    Returns:
        Formatted list of groups with ID, name, type, and details

    Examples:
        list_all_groups()  # All groups (except dynamic)
        list_all_groups(group_type="M365")  # Only M365 groups
        list_all_groups(group_type="Distribution")  # Only distribution lists
        list_all_groups(include_dynamic=True)  # Include dynamic groups
    """
    try:
        client = await get_graph_client()

        # Cap count at 100
        count = min(count, 100)

        # Build request configuration
        from msgraph.generated.groups.groups_request_builder import GroupsRequestBuilder
        from kiota_abstractions.base_request_configuration import RequestConfiguration

        query_params = GroupsRequestBuilder.GroupsRequestBuilderGetQueryParameters(
            top=count,
            select=["id", "displayName", "description", "mail", "mailEnabled",
                    "securityEnabled", "groupTypes", "visibility", "createdDateTime"],
            orderby=["displayName"]
        )
        config = RequestConfiguration(query_parameters=query_params)

        result = await asyncio.wait_for(
            client.groups.get(request_configuration=config),
            timeout=API_TIMEOUT
        )

        if not result or not result.value:
            return (
                "👥 No groups found\n\n"
                "The tenant may have no groups, or Group.Read.All permission "
                "may be required."
            )

        # Filter groups
        filtered_groups = []
        excluded_admin = 0
        excluded_dynamic = 0

        for group in result.value:
            # Skip admin groups
            if is_admin_group(group):
                excluded_admin += 1
                continue

            grp_type = get_group_type(group)

            # Skip dynamic groups unless requested
            if grp_type == 'Dynamic' and not include_dynamic:
                excluded_dynamic += 1
                continue

            # Filter by type if specified
            if group_type:
                if group_type.lower() == "m365" and grp_type != "M365":
                    continue
                elif group_type.lower() == "security" and grp_type != "Security":
                    continue
                elif group_type.lower() == "distribution" and grp_type != "Distribution":
                    continue

            filtered_groups.append((group, grp_type))

        if not filtered_groups:
            filter_msg = f" matching type '{group_type}'" if group_type else ""
            return (
                f"👥 No groups found{filter_msg}\n\n"
                f"Excluded: {excluded_admin} admin groups, {excluded_dynamic} dynamic groups"
            )

        lines = [f"👥 **All Tenant Groups** ({len(filtered_groups)} shown)\n"]

        if excluded_admin > 0 or excluded_dynamic > 0:
            lines.append(f"*Excluded: {excluded_admin} admin groups, {excluded_dynamic} dynamic groups*\n")

        for group, grp_type in filtered_groups:
            lines.append("---")

            # Type badge
            type_badge = {
                'M365': '🟦 M365',
                'Security': '🔒 Security',
                'Distribution': '📧 Distribution',
                'Mail-enabled Security': '📧🔒 Mail+Security',
                'Dynamic': '⚡ Dynamic',
            }.get(grp_type, '❓ Unknown')

            lines.append(f"**{group.display_name or 'Unnamed Group'}** [{type_badge}]")
            lines.append(f"  ID: `{group.id}`")

            if group.description:
                desc = group.description[:80] + "..." if len(group.description) > 80 else group.description
                lines.append(f"  Description: {desc}")

            if group.mail:
                lines.append(f"  Email: {group.mail}")

            if hasattr(group, 'visibility') and group.visibility:
                lines.append(f"  Visibility: {group.visibility}")

        lines.append("")
        lines.append("---")
        lines.append("Use `get_group_by_id()` for full details, `list_group_members()` for membership.")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The Group.Read.All permission is required to list all groups.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error listing groups: {error_type}\n\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def get_group_by_id(
    group_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Get complete details of a specific group.

    Returns full group information including settings, type, and metadata.

    Args:
        group_id: The group ID (from list_all_groups or list_groups) - REQUIRED
        mailbox_id: User context for API call ("me" or email address)

    Returns:
        Full group details including type, settings, creation date

    Examples:
        get_group_by_id(group_id="591f854c-f04a-472d-9587-31c53e2ef219")
    """
    try:
        if not group_id:
            return "❌ Error: group_id is required"

        client = await get_graph_client()

        result = await asyncio.wait_for(
            client.groups.by_group_id(group_id).get(),
            timeout=API_TIMEOUT
        )

        if not result:
            return (
                f"❌ Group not found\n\n"
                f"Group ID: {group_id}"
            )

        # Check if admin group
        if is_admin_group(result):
            return (
                f"⚠️ Admin group (read-only)\n\n"
                f"**{result.display_name}**\n"
                f"This is an admin/role group and cannot be modified.\n"
                f"Group ID: `{group_id}`"
            )

        grp_type = get_group_type(result)

        # Type badge
        type_badge = {
            'M365': '🟦 M365 Group',
            'Security': '🔒 Security Group',
            'Distribution': '📧 Distribution List',
            'Mail-enabled Security': '📧🔒 Mail-enabled Security Group',
            'Dynamic': '⚡ Dynamic Group',
        }.get(grp_type, '❓ Unknown Type')

        lines = [f"👥 **Group Details**\n"]
        lines.append(f"**{result.display_name or 'Unnamed Group'}**")
        lines.append(f"Type: {type_badge}")
        lines.append(f"Group ID: `{result.id}`")

        if result.description:
            lines.append(f"\n**Description:**")
            lines.append(result.description)

        lines.append(f"\n**Settings:**")

        if result.mail:
            lines.append(f"  Email: {result.mail}")

        if hasattr(result, 'visibility') and result.visibility:
            lines.append(f"  Visibility: {result.visibility}")

        lines.append(f"  Mail Enabled: {'Yes' if result.mail_enabled else 'No'}")
        lines.append(f"  Security Enabled: {'Yes' if result.security_enabled else 'No'}")

        group_types = getattr(result, 'group_types', None) or []
        if group_types:
            lines.append(f"  Group Types: {', '.join(group_types)}")

        if hasattr(result, 'created_date_time') and result.created_date_time:
            created = result.created_date_time.strftime("%Y-%m-%d %H:%M")
            lines.append(f"  Created: {created}")

        # Capabilities based on type
        lines.append(f"\n**Capabilities:**")
        if grp_type == 'M365':
            lines.append("  ✅ Planner plans")
            lines.append("  ✅ SharePoint site")
            lines.append("  ✅ Teams (if enabled)")
            lines.append("  ✅ Shared mailbox")
        elif grp_type == 'Distribution':
            lines.append("  ✅ Email distribution")
            lines.append("  ❌ Planner/SharePoint/Teams")
        elif grp_type == 'Security':
            lines.append("  ✅ Access control")
            lines.append("  ❌ Email/Planner/SharePoint")
        elif grp_type == 'Dynamic':
            lines.append("  ⚡ Automatic membership based on rules")
            lines.append("  ⚠️ Members cannot be manually managed")

        lines.append("")
        lines.append("---")
        lines.append("Use `list_group_members()` and `list_group_owners()` for membership details.")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Group not found\n\n"
                f"Group ID: {group_id}\n"
                f"The group may have been deleted or the ID is incorrect."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The Group.Read.All permission is required.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error getting group: {error_type}\n\n"
            f"Group ID: {group_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def list_group_members(
    group_id: str,
    count: int = 50,
    mailbox_id: str = "me"
) -> str:
    """
    List all members of a specific group.

    Returns user details for each member including name and email.

    Args:
        group_id: The group ID (from list_all_groups or list_groups) - REQUIRED
        count: Maximum members to return (default: 50, max: 100)
        mailbox_id: User context for API call ("me" or email address)

    Returns:
        Formatted list of group members with name, email, and ID

    Examples:
        list_group_members(group_id="591f854c-f04a-472d-9587-31c53e2ef219")
        list_group_members(group_id="591f854c-...", count=20)
    """
    try:
        if not group_id:
            return "❌ Error: group_id is required"

        client = await get_graph_client()
        count = min(count, 100)

        # First get group name for context
        group = await asyncio.wait_for(
            client.groups.by_group_id(group_id).get(),
            timeout=API_TIMEOUT
        )
        group_name = group.display_name if group else "Unknown Group"

        # Check if dynamic group
        grp_type = get_group_type(group) if group else "Unknown"
        if grp_type == "Dynamic":
            # Still list members but add warning
            pass

        # Get members
        from msgraph.generated.groups.item.members.members_request_builder import MembersRequestBuilder
        from kiota_abstractions.base_request_configuration import RequestConfiguration

        query_params = MembersRequestBuilder.MembersRequestBuilderGetQueryParameters(
            top=count
        )
        config = RequestConfiguration(query_parameters=query_params)

        result = await asyncio.wait_for(
            client.groups.by_group_id(group_id).members.get(request_configuration=config),
            timeout=API_TIMEOUT
        )

        if not result or not result.value:
            return (
                f"👥 **{group_name}** - Members\n\n"
                f"No members found in this group.\n"
                f"Group ID: `{group_id}`"
            )

        lines = [f"👥 **{group_name}** - Members ({len(result.value)} shown)\n"]

        if grp_type == "Dynamic":
            lines.append("*⚡ Dynamic group - membership is automatic*\n")

        lines.append(f"Group ID: `{group_id}`\n")

        for member in result.value:
            lines.append("---")

            # Members can be users or other objects
            odata_type = getattr(member, 'odata_type', '') or ''

            display_name = getattr(member, 'display_name', None) or 'Unknown'
            member_id = getattr(member, 'id', 'Unknown')

            if 'user' in odata_type.lower():
                mail = getattr(member, 'mail', None) or getattr(member, 'user_principal_name', None)
                job_title = getattr(member, 'job_title', None)

                lines.append(f"👤 **{display_name}**")
                if mail:
                    lines.append(f"   Email: {mail}")
                if job_title:
                    lines.append(f"   Title: {job_title}")
                lines.append(f"   User ID: `{member_id}`")
            else:
                # Could be a group, service principal, etc.
                member_type = odata_type.split('.')[-1] if odata_type else 'Unknown'
                lines.append(f"📦 **{display_name}** ({member_type})")
                lines.append(f"   ID: `{member_id}`")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Group not found\n\n"
                f"Group ID: {group_id}\n"
                f"The group may have been deleted or the ID is incorrect."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The Group.Read.All or GroupMember.Read.All permission is required.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error listing members: {error_type}\n\n"
            f"Group ID: {group_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def list_group_owners(
    group_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    List owners of a specific group.

    Owners have administrative rights over the group including adding/removing
    members and changing group settings.

    Args:
        group_id: The group ID (from list_all_groups or list_groups) - REQUIRED
        mailbox_id: User context for API call ("me" or email address)

    Returns:
        Formatted list of group owners with name, email, and ID

    Examples:
        list_group_owners(group_id="591f854c-f04a-472d-9587-31c53e2ef219")
    """
    try:
        if not group_id:
            return "❌ Error: group_id is required"

        client = await get_graph_client()

        # First get group name for context
        group = await asyncio.wait_for(
            client.groups.by_group_id(group_id).get(),
            timeout=API_TIMEOUT
        )
        group_name = group.display_name if group else "Unknown Group"

        # Get owners
        result = await asyncio.wait_for(
            client.groups.by_group_id(group_id).owners.get(),
            timeout=API_TIMEOUT
        )

        if not result or not result.value:
            return (
                f"👑 **{group_name}** - Owners\n\n"
                f"No owners found for this group.\n"
                f"Group ID: `{group_id}`\n\n"
                f"⚠️ Groups should have at least one owner."
            )

        lines = [f"👑 **{group_name}** - Owners ({len(result.value)} found)\n"]
        lines.append(f"Group ID: `{group_id}`\n")

        for owner in result.value:
            lines.append("---")

            display_name = getattr(owner, 'display_name', None) or 'Unknown'
            owner_id = getattr(owner, 'id', 'Unknown')
            mail = getattr(owner, 'mail', None) or getattr(owner, 'user_principal_name', None)

            lines.append(f"👑 **{display_name}**")
            if mail:
                lines.append(f"   Email: {mail}")
            lines.append(f"   User ID: `{owner_id}`")

        lines.append("")
        lines.append("---")
        lines.append("Owners can manage group settings and membership.")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Group not found\n\n"
                f"Group ID: {group_id}\n"
                f"The group may have been deleted or the ID is incorrect."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The Group.Read.All permission is required.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error listing owners: {error_type}\n\n"
            f"Group ID: {group_id}\n"
            f"Error: {error_msg}"
        )


# ============================================================================
# M365 GROUPS - Phase 2: Group CRUD Operations
# ============================================================================
# Create, update, and delete M365 groups
# Requires Group.ReadWrite.All permission
# ============================================================================

@mcp.tool()
async def create_m365_group(
    display_name: str,
    mail_nickname: str,
    description: str = "",
    visibility: str = "Private",
    mailbox_id: str = "me"
) -> str:
    """
    Create a new M365 group.

    M365 groups provide: Shared mailbox, SharePoint site, Planner plans, Teams support.
    The creating user is automatically added as an owner.

    Args:
        display_name: Group display name (required)
        mail_nickname: Unique email prefix for the group (required)
                      Will become: mail_nickname@yourdomain.com
        description: Optional group description
        visibility: "Private" (members only) or "Public" (anyone can join) - default: Private
        mailbox_id: User context ("me" or email address)

    Returns:
        Created group details with ID

    Examples:
        create_m365_group(display_name="Project Alpha", mail_nickname="projectalpha")
        create_m365_group(display_name="Marketing Team", mail_nickname="marketing", visibility="Public")
    """
    try:
        if not display_name:
            return "❌ Error: display_name is required"
        if not mail_nickname:
            return "❌ Error: mail_nickname is required"

        # Validate mail_nickname (no spaces or special chars)
        import re
        if not re.match(r'^[a-zA-Z0-9_-]+$', mail_nickname):
            return (
                "❌ Error: mail_nickname invalid\n\n"
                "mail_nickname must contain only letters, numbers, underscores, and hyphens.\n"
                f"Provided: '{mail_nickname}'"
            )

        # Validate visibility
        visibility = visibility.capitalize()
        if visibility not in ["Private", "Public"]:
            return (
                "❌ Error: visibility must be 'Private' or 'Public'\n\n"
                f"Provided: '{visibility}'"
            )

        client = await get_graph_client()

        # Build group object for M365 group
        from msgraph.generated.models.group import Group
        new_group = Group()
        new_group.display_name = display_name
        new_group.mail_nickname = mail_nickname
        new_group.description = description or None
        new_group.visibility = visibility
        new_group.group_types = ["Unified"]  # M365 group
        new_group.mail_enabled = True
        new_group.security_enabled = False

        result = await asyncio.wait_for(
            client.groups.post(new_group),
            timeout=API_TIMEOUT
        )

        if not result:
            return "❌ Failed to create group - no result returned"

        grp_type = get_group_type(result)

        lines = [
            "✅ M365 Group created!\n",
            f"**{result.display_name}**",
            f"Group ID: `{result.id}`",
            f"Type: 🟦 {grp_type}",
        ]

        if result.mail:
            lines.append(f"Email: {result.mail}")

        lines.append(f"Visibility: {result.visibility or 'Private'}")

        if result.description:
            lines.append(f"Description: {result.description}")

        if hasattr(result, 'created_date_time') and result.created_date_time:
            lines.append(f"Created: {result.created_date_time.strftime('%Y-%m-%d %H:%M')}")

        lines.append("")
        lines.append("---")
        lines.append("**Next steps:**")
        lines.append(f"- Use `add_group_member(group_id=\"{result.id}\", ...)` to add members")
        lines.append(f"- Use `create_planner_plan(group_id=\"{result.id}\", ...)` to create a Planner plan")
        lines.append(f"- Use `list_group_members(group_id=\"{result.id}\")` to see membership")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The Group.ReadWrite.All permission is required to create groups.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        if "mailNickname" in error_msg.lower() and ("already" in error_msg.lower() or "exist" in error_msg.lower()):
            return (
                f"❌ Mail nickname already exists\n\n"
                f"The mail_nickname '{mail_nickname}' is already in use.\n"
                f"Please choose a different unique nickname."
            )

        if "displayName" in error_msg and "already" in error_msg.lower():
            return (
                f"❌ Group name already exists\n\n"
                f"A group with name '{display_name}' may already exist.\n"
                f"Please choose a different name."
            )

        return (
            f"❌ Error creating group: {error_type}\n\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def update_group(
    group_id: str,
    display_name: str = "",
    description: str = "",
    visibility: str = "",
    mailbox_id: str = "me"
) -> str:
    """
    Update an existing group's properties.

    Only provided parameters are updated - others remain unchanged.
    Cannot update admin groups or system groups.

    Args:
        group_id: The group ID to update (required)
        display_name: New display name (optional - leave empty to keep current)
        description: New description (optional - leave empty to keep current)
        visibility: New visibility "Private" or "Public" (optional - leave empty to keep current)
        mailbox_id: User context ("me" or email address)

    Returns:
        Updated group details

    Examples:
        update_group(group_id="abc123...", display_name="New Team Name")
        update_group(group_id="abc123...", description="Updated project description")
        update_group(group_id="abc123...", visibility="Public")
    """
    try:
        if not group_id:
            return "❌ Error: group_id is required"

        if not display_name and not description and not visibility:
            return (
                "❌ Error: No updates provided\n\n"
                "Please provide at least one of: display_name, description, visibility"
            )

        client = await get_graph_client()

        # First check if it's an admin group
        current = await asyncio.wait_for(
            client.groups.by_group_id(group_id).get(),
            timeout=API_TIMEOUT
        )

        if not current:
            return f"❌ Group not found: {group_id}"

        if is_admin_group(current):
            return (
                f"❌ Cannot update admin group\n\n"
                f"**{current.display_name}** is an admin/role group and cannot be modified."
            )

        grp_type = get_group_type(current)
        if grp_type == "Dynamic":
            return (
                f"❌ Cannot update dynamic group\n\n"
                f"**{current.display_name}** is a dynamic group with automatic membership.\n"
                f"Dynamic group properties are managed by Azure AD."
            )

        # Build update object
        from msgraph.generated.models.group import Group
        update_group_obj = Group()
        updated_fields = []

        if display_name:
            update_group_obj.display_name = display_name
            updated_fields.append(f"display_name: '{display_name}'")

        if description:
            update_group_obj.description = description
            updated_fields.append(f"description: '{description[:50]}...' " if len(description) > 50 else f"description: '{description}'")

        if visibility:
            visibility = visibility.capitalize()
            if visibility not in ["Private", "Public"]:
                return "❌ Error: visibility must be 'Private' or 'Public'"
            update_group_obj.visibility = visibility
            updated_fields.append(f"visibility: '{visibility}'")

        # Perform update
        await asyncio.wait_for(
            client.groups.by_group_id(group_id).patch(update_group_obj),
            timeout=API_TIMEOUT
        )

        # Get updated group
        updated = await asyncio.wait_for(
            client.groups.by_group_id(group_id).get(),
            timeout=API_TIMEOUT
        )

        lines = [
            "✅ Group updated!\n",
            f"**{updated.display_name}**",
            f"Group ID: `{updated.id}`",
            "",
            "**Updated fields:**",
        ]

        for field in updated_fields:
            lines.append(f"  - {field}")

        lines.append("")
        lines.append("---")
        lines.append(f"Use `get_group_by_id(group_id=\"{group_id}\")` to see full details.")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Group not found\n\n"
                f"Group ID: {group_id}\n"
                f"The group may have been deleted or the ID is incorrect."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The Group.ReadWrite.All permission is required to update groups.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error updating group: {error_type}\n\n"
            f"Group ID: {group_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def delete_group(
    group_id: str,
    confirm: bool = False,
    mailbox_id: str = "me"
) -> str:
    """
    Delete a group permanently.

    WARNING: This action cannot be undone. The group and all its resources
    (Planner plans, SharePoint content, Teams, shared mailbox) will be deleted.

    Args:
        group_id: The group ID to delete (required)
        confirm: Must be True to confirm deletion (required safety check)
        mailbox_id: User context ("me" or email address)

    Returns:
        Confirmation of deletion or error message

    Examples:
        delete_group(group_id="abc123...", confirm=True)
    """
    try:
        if not group_id:
            return "❌ Error: group_id is required"

        if not confirm:
            return (
                "⚠️ Deletion not confirmed\n\n"
                "To delete a group, you must set `confirm=True`.\n\n"
                "**Warning:** This will permanently delete:\n"
                "- The group and all members\n"
                "- All Planner plans in the group\n"
                "- The SharePoint site (if M365 group)\n"
                "- The Teams team (if connected)\n"
                "- The shared mailbox and all emails\n\n"
                "This action cannot be undone."
            )

        client = await get_graph_client()

        # First get the group to confirm what we're deleting
        current = await asyncio.wait_for(
            client.groups.by_group_id(group_id).get(),
            timeout=API_TIMEOUT
        )

        if not current:
            return f"❌ Group not found: {group_id}"

        group_name = current.display_name
        grp_type = get_group_type(current)

        # Safety checks
        if is_admin_group(current):
            return (
                f"❌ Cannot delete admin group\n\n"
                f"**{group_name}** is an admin/role group and cannot be deleted."
            )

        if grp_type == "Security":
            return (
                f"⚠️ Security group deletion restricted\n\n"
                f"**{group_name}** is a Security group.\n"
                f"Security groups may control access to resources and should be managed carefully.\n"
                f"If you're certain, please delete via Azure Portal directly."
            )

        # Delete the group
        await asyncio.wait_for(
            client.groups.by_group_id(group_id).delete(),
            timeout=API_TIMEOUT
        )

        lines = [
            "✅ Group deleted\n",
            f"**{group_name}** has been permanently deleted.",
            f"Group ID: `{group_id}`",
            f"Type: {grp_type}",
            "",
            "Associated resources (Planner, SharePoint, Teams, mailbox) will be deleted.",
            "",
            "Note: Some resources may take time to fully remove from the system.",
        ]

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Group not found\n\n"
                f"Group ID: {group_id}\n"
                f"The group may have already been deleted."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The Group.ReadWrite.All permission is required to delete groups.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error deleting group: {error_type}\n\n"
            f"Group ID: {group_id}\n"
            f"Error: {error_msg}"
        )


# ============================================================================
# M365 GROUPS - Phase 3: Membership Management
# ============================================================================
# Add/remove members and owners, lookup users by email
# Requires GroupMember.ReadWrite.All and User.Read.All permissions
# ============================================================================

@mcp.tool()
async def lookup_user_by_email(
    email: str,
    mailbox_id: str = "me"
) -> str:
    """
    Look up a user's Azure AD object ID by email address.

    Needed for membership operations which require user IDs, not email addresses.

    Args:
        email: Email address to look up (required)
        mailbox_id: User context ("me" or email address)

    Returns:
        User details including ID, or error if not found

    Examples:
        lookup_user_by_email(email="john@example.com")
    """
    try:
        if not email:
            return "❌ Error: email is required"

        client = await get_graph_client()

        # Use filter to find user by mail or userPrincipalName
        from msgraph.generated.users.users_request_builder import UsersRequestBuilder
        from kiota_abstractions.base_request_configuration import RequestConfiguration

        # Try both mail and userPrincipalName
        filter_query = f"mail eq '{email}' or userPrincipalName eq '{email}'"

        query_params = UsersRequestBuilder.UsersRequestBuilderGetQueryParameters(
            filter=filter_query,
            select=["id", "displayName", "mail", "userPrincipalName", "jobTitle", "department"],
            top=1
        )
        config = RequestConfiguration(query_parameters=query_params)

        result = await asyncio.wait_for(
            client.users.get(request_configuration=config),
            timeout=API_TIMEOUT
        )

        if not result or not result.value or len(result.value) == 0:
            return (
                f"❌ User not found\n\n"
                f"No user found with email: {email}\n\n"
                f"Possible reasons:\n"
                f"- The email address is incorrect\n"
                f"- The user is not in this Azure AD tenant\n"
                f"- The user is a guest user (try the full email)"
            )

        user = result.value[0]

        lines = [
            "✅ User found\n",
            f"**{user.display_name or 'Unknown'}**",
            f"User ID: `{user.id}`",
        ]

        if user.mail:
            lines.append(f"Email: {user.mail}")
        if user.user_principal_name and user.user_principal_name != user.mail:
            lines.append(f"UPN: {user.user_principal_name}")
        if user.job_title:
            lines.append(f"Title: {user.job_title}")
        if user.department:
            lines.append(f"Department: {user.department}")

        lines.append("")
        lines.append("---")
        lines.append(f"Use this User ID for membership operations:")
        lines.append(f"- `add_group_member(group_id=\"...\", user_id=\"{user.id}\")`")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The User.Read.All permission is required to lookup users.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error looking up user: {error_type}\n\n"
            f"Email: {email}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def add_group_member(
    group_id: str,
    user_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Add a user as a member of a group.

    Use lookup_user_by_email to get the user_id if you only have an email address.

    Args:
        group_id: The group ID to add member to (required)
        user_id: The Azure AD user ID to add (required)
        mailbox_id: User context ("me" or email address)

    Returns:
        Confirmation of membership addition

    Examples:
        add_group_member(group_id="abc123...", user_id="user456...")
    """
    try:
        if not group_id:
            return "❌ Error: group_id is required"
        if not user_id:
            return "❌ Error: user_id is required"

        client = await get_graph_client()

        # First check the group
        group = await asyncio.wait_for(
            client.groups.by_group_id(group_id).get(),
            timeout=API_TIMEOUT
        )

        if not group:
            return f"❌ Group not found: {group_id}"

        group_name = group.display_name
        grp_type = get_group_type(group)

        if is_admin_group(group):
            return (
                f"❌ Cannot modify admin group\n\n"
                f"**{group_name}** is an admin/role group and cannot have members added."
            )

        if grp_type == "Dynamic":
            return (
                f"❌ Cannot modify dynamic group\n\n"
                f"**{group_name}** is a dynamic group with automatic membership.\n"
                f"Members are managed by Azure AD based on rules."
            )

        # Get user info for confirmation
        user = await asyncio.wait_for(
            client.users.by_user_id(user_id).get(),
            timeout=API_TIMEOUT
        )
        user_name = user.display_name if user else "Unknown User"

        # Add member using POST to members/$ref
        from msgraph.generated.models.reference_create import ReferenceCreate
        reference = ReferenceCreate()
        reference.odata_id = f"https://graph.microsoft.com/v1.0/directoryObjects/{user_id}"

        await asyncio.wait_for(
            client.groups.by_group_id(group_id).members.ref.post(reference),
            timeout=API_TIMEOUT
        )

        lines = [
            "✅ Member added\n",
            f"**{user_name}** is now a member of **{group_name}**",
            "",
            f"User ID: `{user_id}`",
            f"Group ID: `{group_id}`",
            "",
            "---",
            f"Use `list_group_members(group_id=\"{group_id}\")` to see all members.",
        ]

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg:
            if "User" in error_msg or user_id in error_msg:
                return (
                    f"❌ User not found\n\n"
                    f"User ID: {user_id}\n"
                    f"Use `lookup_user_by_email()` to find the correct user ID."
                )
            return (
                f"❌ Group not found\n\n"
                f"Group ID: {group_id}"
            )

        if "already exist" in error_msg.lower() or "added" in error_msg.lower():
            return (
                f"ℹ️ User is already a member\n\n"
                f"User ID: `{user_id}`\n"
                f"Group ID: `{group_id}`\n\n"
                f"No action needed - user was already a member of this group."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The GroupMember.ReadWrite.All permission is required to manage membership.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error adding member: {error_type}\n\n"
            f"Group ID: {group_id}\n"
            f"User ID: {user_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def remove_group_member(
    group_id: str,
    user_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Remove a user from group membership.

    Note: This removes member access but not owner access. To remove an owner,
    use remove_group_owner.

    Args:
        group_id: The group ID to remove member from (required)
        user_id: The Azure AD user ID to remove (required)
        mailbox_id: User context ("me" or email address)

    Returns:
        Confirmation of membership removal

    Examples:
        remove_group_member(group_id="abc123...", user_id="user456...")
    """
    try:
        if not group_id:
            return "❌ Error: group_id is required"
        if not user_id:
            return "❌ Error: user_id is required"

        client = await get_graph_client()

        # First check the group
        group = await asyncio.wait_for(
            client.groups.by_group_id(group_id).get(),
            timeout=API_TIMEOUT
        )

        if not group:
            return f"❌ Group not found: {group_id}"

        group_name = group.display_name
        grp_type = get_group_type(group)

        if is_admin_group(group):
            return (
                f"❌ Cannot modify admin group\n\n"
                f"**{group_name}** is an admin/role group and cannot have members removed."
            )

        if grp_type == "Dynamic":
            return (
                f"❌ Cannot modify dynamic group\n\n"
                f"**{group_name}** is a dynamic group with automatic membership.\n"
                f"Members are managed by Azure AD based on rules."
            )

        # Get user info for confirmation
        try:
            user = await asyncio.wait_for(
                client.users.by_user_id(user_id).get(),
                timeout=API_TIMEOUT
            )
            user_name = user.display_name if user else "Unknown User"
        except:
            user_name = "Unknown User"

        # Remove member using DELETE to members/{id}/$ref
        await asyncio.wait_for(
            client.groups.by_group_id(group_id).members.by_directory_object_id(user_id).ref.delete(),
            timeout=API_TIMEOUT
        )

        lines = [
            "✅ Member removed\n",
            f"**{user_name}** is no longer a member of **{group_name}**",
            "",
            f"User ID: `{user_id}`",
            f"Group ID: `{group_id}`",
            "",
            "---",
            f"Use `list_group_members(group_id=\"{group_id}\")` to verify.",
        ]

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"ℹ️ User is not a member\n\n"
                f"User ID: `{user_id}`\n"
                f"Group ID: `{group_id}`\n\n"
                f"The user may have already been removed or was never a member."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The GroupMember.ReadWrite.All permission is required to manage membership.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error removing member: {error_type}\n\n"
            f"Group ID: {group_id}\n"
            f"User ID: {user_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def add_group_owner(
    group_id: str,
    user_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Add a user as an owner of a group.

    Owners have administrative rights: add/remove members, change settings, delete group.
    Adding as owner automatically adds as member if not already.

    Args:
        group_id: The group ID to add owner to (required)
        user_id: The Azure AD user ID to add as owner (required)
        mailbox_id: User context ("me" or email address)

    Returns:
        Confirmation of owner addition

    Examples:
        add_group_owner(group_id="abc123...", user_id="user456...")
    """
    try:
        if not group_id:
            return "❌ Error: group_id is required"
        if not user_id:
            return "❌ Error: user_id is required"

        client = await get_graph_client()

        # First check the group
        group = await asyncio.wait_for(
            client.groups.by_group_id(group_id).get(),
            timeout=API_TIMEOUT
        )

        if not group:
            return f"❌ Group not found: {group_id}"

        group_name = group.display_name

        if is_admin_group(group):
            return (
                f"❌ Cannot modify admin group\n\n"
                f"**{group_name}** is an admin/role group and cannot have owners added."
            )

        # Get user info for confirmation
        user = await asyncio.wait_for(
            client.users.by_user_id(user_id).get(),
            timeout=API_TIMEOUT
        )
        user_name = user.display_name if user else "Unknown User"

        # Add owner using POST to owners/$ref
        from msgraph.generated.models.reference_create import ReferenceCreate
        reference = ReferenceCreate()
        reference.odata_id = f"https://graph.microsoft.com/v1.0/directoryObjects/{user_id}"

        await asyncio.wait_for(
            client.groups.by_group_id(group_id).owners.ref.post(reference),
            timeout=API_TIMEOUT
        )

        lines = [
            "✅ Owner added\n",
            f"**{user_name}** is now an owner of **{group_name}**",
            "",
            f"User ID: `{user_id}`",
            f"Group ID: `{group_id}`",
            "",
            "**Owner privileges:**",
            "- Add and remove members",
            "- Modify group settings",
            "- Delete the group",
            "- Manage Planner plans and SharePoint",
            "",
            "---",
            f"Use `list_group_owners(group_id=\"{group_id}\")` to see all owners.",
        ]

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg:
            if "User" in error_msg or user_id in error_msg:
                return (
                    f"❌ User not found\n\n"
                    f"User ID: {user_id}\n"
                    f"Use `lookup_user_by_email()` to find the correct user ID."
                )
            return (
                f"❌ Group not found\n\n"
                f"Group ID: {group_id}"
            )

        if "already exist" in error_msg.lower() or "added" in error_msg.lower():
            return (
                f"ℹ️ User is already an owner\n\n"
                f"User ID: `{user_id}`\n"
                f"Group ID: `{group_id}`\n\n"
                f"No action needed - user was already an owner of this group."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The GroupMember.ReadWrite.All permission is required to manage ownership.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error adding owner: {error_type}\n\n"
            f"Group ID: {group_id}\n"
            f"User ID: {user_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def remove_group_owner(
    group_id: str,
    user_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Remove a user from group ownership.

    Note: Cannot remove the last owner of a group - at least one owner must remain.
    This removes owner access but NOT member access.

    Args:
        group_id: The group ID to remove owner from (required)
        user_id: The Azure AD user ID to remove as owner (required)
        mailbox_id: User context ("me" or email address)

    Returns:
        Confirmation of owner removal

    Examples:
        remove_group_owner(group_id="abc123...", user_id="user456...")
    """
    try:
        if not group_id:
            return "❌ Error: group_id is required"
        if not user_id:
            return "❌ Error: user_id is required"

        client = await get_graph_client()

        # First check the group
        group = await asyncio.wait_for(
            client.groups.by_group_id(group_id).get(),
            timeout=API_TIMEOUT
        )

        if not group:
            return f"❌ Group not found: {group_id}"

        group_name = group.display_name

        if is_admin_group(group):
            return (
                f"❌ Cannot modify admin group\n\n"
                f"**{group_name}** is an admin/role group and cannot have owners removed."
            )

        # Check current owner count
        owners_result = await asyncio.wait_for(
            client.groups.by_group_id(group_id).owners.get(),
            timeout=API_TIMEOUT
        )

        if owners_result and owners_result.value:
            owner_count = len(owners_result.value)
            if owner_count <= 1:
                return (
                    f"❌ Cannot remove last owner\n\n"
                    f"**{group_name}** has only {owner_count} owner(s).\n"
                    f"Every group must have at least one owner.\n\n"
                    f"To remove this owner, first add another owner using `add_group_owner()`."
                )

        # Get user info for confirmation
        try:
            user = await asyncio.wait_for(
                client.users.by_user_id(user_id).get(),
                timeout=API_TIMEOUT
            )
            user_name = user.display_name if user else "Unknown User"
        except:
            user_name = "Unknown User"

        # Remove owner using DELETE to owners/{id}/$ref
        await asyncio.wait_for(
            client.groups.by_group_id(group_id).owners.by_directory_object_id(user_id).ref.delete(),
            timeout=API_TIMEOUT
        )

        lines = [
            "✅ Owner removed\n",
            f"**{user_name}** is no longer an owner of **{group_name}**",
            "",
            f"User ID: `{user_id}`",
            f"Group ID: `{group_id}`",
            "",
            "Note: User may still be a member (use `remove_group_member()` to remove completely).",
            "",
            "---",
            f"Use `list_group_owners(group_id=\"{group_id}\")` to verify.",
        ]

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                "❌ Authentication required\n\n"
                "Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"ℹ️ User is not an owner\n\n"
                f"User ID: `{user_id}`\n"
                f"Group ID: `{group_id}`\n\n"
                f"The user may have already been removed or was never an owner."
            )

        if "last owner" in error_msg.lower() or "at least one" in error_msg.lower():
            return (
                f"❌ Cannot remove last owner\n\n"
                f"This is the last owner of the group.\n"
                f"Every group must have at least one owner.\n\n"
                f"To remove this owner, first add another owner using `add_group_owner()`."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                "❌ Permission denied\n\n"
                "The GroupMember.ReadWrite.All permission is required to manage ownership.\n"
                "Please add this permission in Azure Portal and re-authenticate."
            )

        return (
            f"❌ Error removing owner: {error_type}\n\n"
            f"Group ID: {group_id}\n"
            f"User ID: {user_id}\n"
            f"Error: {error_msg}"
        )


# ============================================================================
# PLANNER - Plan CRUD Operations
# ============================================================================

@mcp.tool()
async def create_planner_plan(
    group_id: str,
    title: str,
    mailbox_id: str = "me"
) -> str:
    """
    Create a new Planner plan in an existing M365 Group.

    Plans are created within a group context - use list_groups() to find
    available groups and their IDs.

    Args:
        group_id: The M365 Group ID to create the plan in (from list_groups) - REQUIRED
        title: Plan title - REQUIRED
        mailbox_id: User context ("me" or email address)

    Returns:
        Created plan details with ID

    Examples:
        create_planner_plan(group_id="abc123...", title="Q1 Projects")
        create_planner_plan(group_id="abc123...", title="Marketing Campaign 2025")
    """
    try:
        client = await get_graph_client()

        # Build plan object
        new_plan = PlannerPlan()
        new_plan.owner = group_id  # owner is the group ID
        new_plan.title = title

        # Create the plan - POST does NOT need ETag
        result = await asyncio.wait_for(
            client.planner.plans.post(new_plan),
            timeout=API_TIMEOUT
        )

        if not result:
            return "❌ Failed to create plan - no result returned"

        lines = [
            f"✅ Planner plan created!\n",
            f"**{result.title}**",
            f"Plan ID: `{result.id}`",
            f"Owner Group: `{result.owner}`",
        ]

        if result.created_date_time:
            lines.append(f"Created: {result.created_date_time.strftime('%Y-%m-%d %H:%M')}")

        lines.append("")
        lines.append("---")
        lines.append("Next steps:")
        lines.append(f"- Use `list_planner_buckets(plan_id=\"{result.id}\")` to see columns")
        lines.append(f"- Use `create_planner_task(plan_id=\"{result.id}\", ...)` to add tasks")

        return "\n".join(lines)

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "Authorization_RequestDenied" in error_msg or "Forbidden" in error_msg:
            return (
                f"❌ Permission denied\n\n"
                f"The Group.ReadWrite.All permission may be required to create Planner plans.\n"
                f"Please check permissions in Azure Portal and re-authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Group not found\n\n"
                f"Group ID: {group_id}\n"
                f"The group may not exist or you may not have access.\n"
                f"Use list_groups() to find available groups."
            )

        if "InvalidRequest" in error_msg or "BadRequest" in error_msg:
            return (
                f"❌ Could not create plan\n\n"
                f"Title: {title}\n"
                f"Group ID: {group_id}\n"
                f"The request was invalid. Check the group ID is correct."
            )

        return (
            f"❌ Error creating plan: {error_type}\n\n"
            f"Title: {title}\n"
            f"Group ID: {group_id}\n"
            f"Error: {error_msg}"
        )


@mcp.tool()
async def delete_planner_plan(
    plan_id: str,
    mailbox_id: str = "me"
) -> str:
    """
    Delete a Planner plan permanently.

    WARNING: This action cannot be undone. All tasks and buckets in the plan
    will also be deleted.

    Args:
        plan_id: The plan ID (from list_planner_plans) - REQUIRED
        mailbox_id: User context ("me" or email address)

    Returns:
        Confirmation of deletion

    Examples:
        delete_planner_plan(plan_id="abc123...")
    """
    try:
        # Load token for direct API call (SDK has etag header issues)
        token_data = load_token_cache()
        if not token_data or not is_token_valid(token_data):
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        access_token = token_data['access_token']

        # Get plan first to confirm it exists and get title/etag
        async with httpx.AsyncClient() as http_client:
            get_resp = await http_client.get(
                f"https://graph.microsoft.com/v1.0/planner/plans/{plan_id}",
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=API_TIMEOUT
            )

            if get_resp.status_code == 404:
                return f"❌ Plan not found\n\nPlan ID: {plan_id}"

            if get_resp.status_code != 200:
                return f"❌ Error fetching plan\n\nStatus: {get_resp.status_code}\n{get_resp.text}"

            plan_data = get_resp.json()
            plan_title = plan_data.get("title", "Untitled")
            etag = plan_data.get("@odata.etag", "")

            # Delete the plan with If-Match header
            delete_resp = await http_client.delete(
                f"https://graph.microsoft.com/v1.0/planner/plans/{plan_id}",
                headers={
                    "Authorization": f"Bearer {access_token}",
                    "If-Match": etag
                },
                timeout=API_TIMEOUT
            )

            if delete_resp.status_code == 204:
                return (
                    f"✅ Planner plan deleted!\n\n"
                    f"**{plan_title}**\n"
                    f"Plan ID: `{plan_id}`\n\n"
                    f"⚠️ This action cannot be undone. All tasks and buckets were also deleted."
                )
            elif delete_resp.status_code == 409:
                return (
                    f"❌ Plan was modified by someone else\n\n"
                    f"Plan ID: {plan_id}\n"
                    f"Please try again."
                )
            elif delete_resp.status_code == 403:
                return (
                    f"❌ Permission denied\n\n"
                    f"Plan ID: {plan_id}\n"
                    f"You may not have permission to delete this plan.\n"
                    f"Only plan owners or group admins can delete plans."
                )
            else:
                return f"❌ Error deleting plan\n\nStatus: {delete_resp.status_code}\n{delete_resp.text}"

    except Exception as e:
        error_type = type(e).__name__
        error_msg = str(e)

        if "AuthenticationRequired" in error_msg or "InvalidAuthenticationToken" in error_msg:
            return (
                f"❌ Authentication required\n\n"
                f"Please call test_connection first to authenticate."
            )

        if "ResourceNotFound" in error_msg or "does not exist" in error_msg:
            return (
                f"❌ Plan not found\n\n"
                f"Plan ID: {plan_id}\n"
                f"The plan may have already been deleted."
            )

        return (
            f"❌ Error deleting plan: {error_type}\n\n"
            f"Plan ID: {plan_id}\n"
            f"Error: {error_msg}"
        )


# Run server with STDIO transport (wrapped by mcp-proxy for HTTP)
if __name__ == "__main__":
    import sys

    # Log startup to stderr (stdout is reserved for MCP protocol)
    print("=" * 60, file=sys.stderr, flush=True)
    print("Microsoft Graph Email MCP Server", file=sys.stderr, flush=True)
    print("=" * 60, file=sys.stderr, flush=True)
    print(f"Server: microsoft-graph-email", file=sys.stderr, flush=True)
    print(f"Token cache: {TOKEN_CACHE_DIR}/{TOKEN_CACHE_NAME}", file=sys.stderr, flush=True)
    print(f"Transport: STDIO (wrapped by mcp-proxy for HTTP)", file=sys.stderr, flush=True)
    print("=" * 60, file=sys.stderr, flush=True)
    print(file=sys.stderr, flush=True)

    # Run with STDIO transport (mcp-proxy will wrap this)
    mcp.run()
